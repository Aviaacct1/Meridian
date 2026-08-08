#!/usr/bin/env python3
"""Offline test of the visual layer. No engine, no stores, no network.

    py -3.12 test_visual_layer.py

Every number in FC below is a TEST FIXTURE, invented to exercise the code path.
Nothing here is a forecast of anything and nothing here may be quoted.

Avia Solutions Limited. All rights reserved.
"""
import os
import sys

import tempfile

V4 = os.path.dirname(os.path.abspath(__file__))
# Found by landmark, not by folder depth. See deck_paths.py: the previous line counted two
# levels up and appended "app", which broke the moment the renderer changed folder depth,
# and this suite failing on ModuleNotFoundError is how the live entry point's identical
# bug was found.
import deck_paths as _DP                                   # noqa: E402
APP = _DP.on_path(V4, who="test_visual_layer")
OUT = os.path.join(tempfile.gettempdir(), "avia_visual_test")
os.makedirs(OUT, exist_ok=True)

FAIL = []


CHECKS = 0


def check(name, cond, detail=""):
    global CHECKS
    CHECKS += 1
    print("%-46s %s %s" % (name, "PASS" if cond else "FAIL", detail))
    if not cond:
        FAIL.append(name)


# ---------------------------------------------------------------------------
# Dependency preflight. A traceback out of an import twenty lines into a test
# tells the reader that something broke, not what is missing or how to fix it,
# and it is the same failure the tool is supposed to report rather than throw.
# ---------------------------------------------------------------------------
NEEDED = [
    ("matplotlib", "matplotlib", "every chart and the route map"),
    ("mpl_toolkits.basemap", "basemap", "the route map only; the deck still "
                                        "builds without it, one page shorter"),
    ("pptx", "python-pptx", "the PowerPoint renderer"),
    ("PIL", "pillow", "image sizing and the picture fit"),
]


def preflight():
    missing = []
    for mod, pkg, what in NEEDED:
        try:
            __import__(mod)
        except ImportError:
            missing.append((mod, pkg, what))
    if not missing:
        return True
    exe = os.path.basename(sys.executable)
    print("MISSING DEPENDENCIES on %s (Python %s)"
          % (exe, ".".join(str(v) for v in sys.version_info[:3])))
    for mod, pkg, what in missing:
        print("   %-22s needed for %s" % (mod, what))
    print("\n   %s -m pip install %s"
          % (sys.executable, " ".join(p for _m, p, _w in missing)))
    hard = [m for m, _p, _w in missing if m != "mpl_toolkits.basemap"]
    if hard:
        print("\nSTOPPING: %s is required and absent." % ", ".join(hard))
        sys.exit(2)
    print("\nCONTINUING WITHOUT BASEMAP. The map checks will report the "
          "absence, which is the behaviour under test.\n")
    return False


HAVE_BASEMAP = preflight()


# --- a synthetic engine forecast, shaped exactly as calibrated_forecast returns
FC = {
    "ok": True,
    "origin": {"iata": "EDI", "city": "Edinburgh", "country": "United Kingdom"},
    "dest": {"iata": "AUS", "city": "Austin", "country": "United States"},
    "catchment": {"origin_ll": [55.950, -3.372], "dest_ll": [30.194, -97.670]},
    "demand": {"natural": 118400, "current": 71200, "stimulation": 1.15,
               "qsi_share": 0.412, "p2p_carried": 41900, "feed_behind": 9600,
               "feed_beyond": 14300, "total": 65800, "pdew_total": 180.3,
               "avg_fare": 412, "beyond_pdew": []},
    "capacity": {"carried": 65800, "spill": 1200, "load": 0.812,
                 "annual_capacity": 81000, "aircraft": "A21X", "freq": 5,
                 "recommendation": "Fits 5x/week A21X at 81% load."},
    "schedule": {"outbound": {"sector": "EDI-AUS", "dep": "10:45", "arr": "15:10"},
                 "inbound": {"sector": "AUS-EDI", "dep": "17:20", "arr": "08:35"}},
    "distance_nm": 4193, "block_min": 585, "year": 2026,
}

# --- 1. the chart -----------------------------------------------------------
import avia_charts as AC
p = AC.demand_build(
    os.path.join(OUT, "chart.png"), market=FC["demand"]["natural"],
    p2p_carried=FC["demand"]["p2p_carried"],
    feed_behind=FC["demand"]["feed_behind"],
    feed_beyond=FC["demand"]["feed_beyond"], carried=FC["demand"]["total"],
    load=FC["capacity"]["load"], origin_city="Edinburgh", dest_city="Austin",
    year=2026)
check("chart drawn", bool(p) and os.path.getsize(p) > 20000,
      "%d bytes" % (os.path.getsize(p) if p else 0))
check("chart returns None on empty legs",
      AC.demand_build(os.path.join(OUT, "x.png"), market=100, p2p_carried=0,
                      feed_behind=0, feed_beyond=0, carried=0, load=None,
                      origin_city="A", dest_city="B") is None)

# --- 2. deck_figures --------------------------------------------------------
import deck_figures as DF
figs, notes = DF.build(FC, OUT)
WANT = ["demand_build", "route_map"] if HAVE_BASEMAP else ["demand_build"]
check("figures drawn", sorted(figs) == WANT, "%s notes=%s" % (sorted(figs), notes))
if HAVE_BASEMAP:
    check("route map file non-trivial",
          bool(figs.get("route_map")) and os.path.getsize(figs["route_map"]) > 50000,
          "%d bytes" % (os.path.getsize(figs["route_map"]) if figs.get("route_map") else 0))
else:
    check("missing basemap is reported, not swallowed",
          "route_map" in notes and "basemap" in notes["route_map"],
          notes.get("route_map"))
facts = DF.route_facts(FC)
check("route facts read from the engine", len(facts) == 3, facts)
check("block time formatted", "9h 45m" in facts[1], facts[1])

stub = {"origin": {"iata": "EDI", "city": "Edinburgh"},
        "dest": {"iata": "AUS", "city": "Austin"}, "demand": {}, "capacity": {}}
f2, n2 = DF.build(stub, OUT)
check("stub draws nothing", f2 == {}, f2)
check("stub reports both reasons", sorted(n2) == ["demand_build", "route_map"],
      list(n2.values()))

# --- 3. forecast_spec -------------------------------------------------------
import forecast_spec as FS
full = FS.from_forecast(FC, currency="GBP", charted=False)
trim = FS.from_forecast(FC, currency="GBP", charted=True)
nf = len(full["segments"]["rows"])
nt = len(trim["segments"]["rows"])
check("full segments table keeps every row", nf == 10, "%d rows" % nf)
check("charted segments table drops the five chart rows", nt == 5,
      "%d rows" % nt)
check("charted table renames itself",
      trim["segments"]["title"] == "How the build works",
      trim["segments"]["title"])
check("no volume row survives the trim",
      not any("market in the catchment" in r[0] or "Total carried" in r[0]
              for r in trim["segments"]["rows"]),
      [r[0] for r in trim["segments"]["rows"]])

# --- 4. schedule viability, one implementation ------------------------------
import schedule_viability as SV
check("viability silent at a healthy load",
      SV.schedule_viability(FC) is None)
thin = dict(FC, capacity=dict(FC["capacity"], load=0.38, freq=5))
v = SV.schedule_viability(thin)
check("viability fires at 38%", v and v["band"] == "NOT A PROPOSITION", v and v["band"])
check("viability sizes the schedule", v and v["sized_frequency"] == 2,
      v and v["sized_frequency"])
check("John's wording is the question",
      v and v["question"].startswith("At 38% load factor this is unlikely"),
      v and v["question"])
check("deck side delegates to the same function",
      FS.schedule_viability(thin) == v)

# --- 5. prose_size ----------------------------------------------------------
import render_pptx as RPX
one = [(None, "x" * 430)]
four = [(None, "y" * 430) for _ in range(4)]
s1 = RPX.prose_size(one, RPX.COL * 0.62, 750)
s4 = RPX.prose_size(four, RPX.COL, 750)
check("a single short paragraph sets larger", s1 > 12.5, "%.1fpt" % s1)
check("a full page stays at the base size", s4 == 12.5, "%.1fpt" % s4)
check("the step up is capped", s1 <= RPX.PROSE_PT_MAX, "%.1fpt" % s1)

# --- 6. spec assembly and budget --------------------------------------------
import spec_from_research as SFR
import deck_spec as S
research = {"origin_city": "Edinburgh", "destination_city": "Austin",
            "executive_summary": "A test fixture paragraph. " * 6,
            "blocks": [{"block_id": "economic_context",
                        "block_name": "The market", "relevance": "include",
                        "summary": "", "presentation_text": "Test prose. " * 10,
                        "data_gaps": [],
                        "findings": [{"claim": "Austin metro output reached $268bn",
                                      "caption": "Austin metro output, a test fixture.",
                                      "value": "$268bn", "unit": "", "year": "2024",
                                      "relevance_to_case": "", "is_single_source": True,
                                      "citations": [{"source_name": "Test source",
                                                     "title": "", "date": "2024",
                                                     "url": "https://example.invalid"}]}]}]}
spec = SFR.build_spec(research, forecast=trim, codename="Forth",
                      title="A direct link between\nEdinburgh and Austin",
                      strap="Test", prepared_for="Edinburgh Airport",
                      date="7 August 2026", figures=figs, route_facts=facts)
kinds = [sl["type"] for sl in spec["slides"]]
check("figure slides in the spec", kinds.count("figure") == (2 if HAVE_BASEMAP else 1), kinds)
titles = [sl.get("title") for sl in spec["slides"]]
check("route map opens the forecast section",
      ("The route" in titles) == HAVE_BASEMAP, titles)
check("demand chart is on the page", "Where the traffic comes from" in titles)
check("segments table renamed in the spec", "How the build works" in titles)
check("every figure slide is sourced",
      all(sl.get("source") for sl in spec["slides"] if sl["type"] == "figure"))
warn = S.check(spec, verbose=False)
check("content budget clean", not warn, warn[:3])

# no figures at all: the slides must be absent, not placeholders
spec2 = SFR.build_spec(research, forecast=full, codename="Forth", title="T",
                       strap="", prepared_for="X", date="")
check("no figures means no figure slides",
      [sl["type"] for sl in spec2["slides"]].count("figure") == 0)
check("uncharted table keeps its old title",
      "Point to point demand" in [sl.get("title") for sl in spec2["slides"]])

# --- 7. render --------------------------------------------------------------
out = os.path.join(OUT, "test_deck.pptx")
RPX.render(spec, out, safe_fonts=True, assets_dir=os.path.join(V4, "assets_obs"))
check("deck rendered", os.path.exists(out) and os.path.getsize(out) > 100000,
      "%d bytes" % os.path.getsize(out))
from pptx import Presentation
pr = Presentation(out)
pics = sum(1 for sl in pr.slides for sh in sl.shapes if sh.shape_type == 13)
check("pictures placed in the pptx", pics >= (2 if HAVE_BASEMAP else 1),
      "%d pictures" % pics)
check("author is the Observatory",
      pr.core_properties.author == "The Aviation Observatory",
      pr.core_properties.author)
import zipfile, re as _re
with zipfile.ZipFile(out) as z:
    bad = set()
    for n in z.namelist():
        if n.endswith(".xml"):
            for m in _re.findall(r'lang="([^"]+)"', z.read(n).decode("utf-8", "ignore")):
                if m != "en-GB":
                    bad.add(m)
check("no language other than en-GB", not bad, sorted(bad))

print("\n%d checks, %d failed%s" % (CHECKS, len(FAIL),
      "" if HAVE_BASEMAP else "   (no basemap: the map checks ran as the absent case)"))
if FAIL:
    print("FAILED: %s" % ", ".join(FAIL))
sys.exit(1 if FAIL else 0)
