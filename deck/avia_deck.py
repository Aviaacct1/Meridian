"""
Avia Solutions route-pitch deck generator: house-style library.

House style reverse-engineered from the live China Airlines TPE-SJC deck
(August 2026), which is the current bespoke standard:
  4:3 slide, navy 021D49 header band, body text 002060, orange FFA800
  callout panels, client logo top right, dense evidence tables, full-bleed
  photographic section dividers, source attribution on every figure.

Author metadata on every generated file: Avia Solutions.
Proofing language: en-GB.

Avia Solutions Limited. All rights reserved.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ----------------------------------------------------------------------------
# Palette and metrics
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x02, 0x1D, 0x49)        # header band, circles, table headers
BODY = RGBColor(0x00, 0x20, 0x60)        # body copy
ORANGE = RGBColor(0xFF, 0xA8, 0x00)      # callout panels
CYAN = RGBColor(0x00, 0xB0, 0xF0)        # headline figures on map slides
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x59, 0x59, 0x59)        # source notes
LIGHT = RGBColor(0xEA, 0xF0, 0xF8)       # table banding
MIDBLUE = RGBColor(0x1F, 0x6F, 0xB2)     # chart series
TEAL = RGBColor(0x14, 0x5A, 0x6E)        # right-hand evidence panel
GREEN = RGBColor(0x2E, 0x8B, 0x57)
RED = RGBColor(0xC0, 0x39, 0x2B)

FONT = "Arial"

SW = 10.0        # slide width, inches
SH = 7.5         # slide height, inches
BAND_H = 1.18    # header band height
M = 0.30         # left / right margin
TOP = 1.34       # first usable y below the band
BOTTOM = 7.06    # last usable y above the source line


def _pt(v):
    return Pt(v)


class AviaDeck:
    """Builds an Avia route-pitch deck in the house style."""

    def __init__(self, deck_title, event_line, assets_dir,
                 client_logo=None, airline_logo=None, watermark=None):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SW)
        self.prs.slide_height = Inches(SH)
        self.blank = self.prs.slide_layouts[6]
        self.deck_title = deck_title
        self.event_line = event_line
        self.assets = assets_dir
        self.client_logo = client_logo
        self.airline_logo = airline_logo
        self.watermark = watermark
        self.n = 0
        self.contents = []

    # -- asset helper --------------------------------------------------------
    def a(self, name):
        p = os.path.join(self.assets, name)
        return p if os.path.exists(p) else None

    # -- primitives ----------------------------------------------------------
    def _slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _rect(self, s, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE,
              line_w=1.0):
        sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid()
            sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(line_w)
        sp.shadow.inherit = False
        return sp

    def _text(self, s, x, y, w, h, runs, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, space_after=4, line_spacing=None):
        """runs: list of (text, size, bold, colour) or list of lists for paragraphs."""
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        if runs and not isinstance(runs[0], list):
            runs = [runs]
        first = True
        for para in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(space_after)
            if line_spacing:
                p.line_spacing = line_spacing
            if not isinstance(para, list):
                para = [para]
            for item in para:
                txt, sz, bold, col = item
                r = p.add_run()
                r.text = txt
                r.font.size = Pt(sz)
                r.font.bold = bold
                r.font.color.rgb = col
                r.font.name = FONT
        return tb

    def _pic(self, s, name, x, y, w=None, h=None, crop_to=None):
        p = self.a(name)
        if p is None:
            from PIL import Image  # noqa: F401
            return self._placeholder(s, x, y, w or 3, h or 2, name)
        if crop_to:
            return self._pic_cover(s, p, x, y, crop_to[0], crop_to[1])
        kw = {}
        if w:
            kw["width"] = Inches(w)
        if h:
            kw["height"] = Inches(h)
        return s.shapes.add_picture(p, Inches(x), Inches(y), **kw)

    def _pic_cover(self, s, path, x, y, w, h):
        """Place an image filling the box, cropped centrally (cover fit)."""
        from PIL import Image
        iw, ih = Image.open(path).size
        box_ar = w / h
        img_ar = iw / ih
        if img_ar > box_ar:      # image wider: crop left/right
            new_w = h * img_ar
            pic = s.shapes.add_picture(path, Inches(x), Inches(y),
                                       width=Inches(new_w), height=Inches(h))
            over = (new_w - w) / new_w / 2
            pic.crop_left = over
            pic.crop_right = over
            pic.left, pic.width = Inches(x), Inches(w)
        else:                    # image taller: crop top/bottom
            new_h = w / img_ar
            pic = s.shapes.add_picture(path, Inches(x), Inches(y),
                                       width=Inches(w), height=Inches(new_h))
            over = (new_h - h) / new_h / 2
            pic.crop_top = over
            pic.crop_bottom = over
            pic.top, pic.height = Inches(y), Inches(h)
        return pic

    def _placeholder(self, s, x, y, w, h, label):
        sp = self._rect(s, x, y, w, h, fill=RGBColor(0xE8, 0xEC, 0xF2),
                        line=RGBColor(0x99, 0xA8, 0xBE))
        sp.line.dash_style = 4
        self._text(s, x + 0.1, y + h / 2 - 0.28, w - 0.2, 0.56,
                   [[("IMAGE SLOT", 11, True, RGBColor(0x66, 0x77, 0x91))],
                    [(label.replace("_", " ").rsplit(".", 1)[0], 9, False,
                      RGBColor(0x66, 0x77, 0x91))]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return sp

    # -- chrome --------------------------------------------------------------
    def _chrome(self, s, page=True):
        self._rect(s, -0.02, -0.02, SW + 0.04, BAND_H, fill=NAVY)
        if self.client_logo and self.a(self.client_logo):
            self._pic(s, self.client_logo, 8.36, 0.20, w=1.42)
        if page:
            self.n += 1
            self._rect(s, 9.55, 7.19, 0.45, 0.31, fill=NAVY)
            self._text(s, 9.55, 7.24, 0.45, 0.22, [(str(self.n), 9, False, WHITE)],
                       align=PP_ALIGN.CENTER)

    def source(self, s, text, y=7.16, x=M, w=8.9, size=8.5):
        self._text(s, x, y, w, 0.26, [(text, size, False, GREY)])

    # ------------------------------------------------------------------
    # Slide types
    # ------------------------------------------------------------------
    def _hero(self, s, image, x, y, w, h, label=None):
        """Full-bleed image, or a designed navy fallback when the pool is thin."""
        p = self.a(image)
        if p:
            return self._pic_cover(s, p, x, y, w, h)
        self._rect(s, x, y, w, h, fill=NAVY)
        for i in range(9):
            self._rect(s, x + w * (0.08 + i * 0.105), y + h * 0.16,
                       w * 0.045, h * (0.10 + 0.075 * ((i * 5) % 9)),
                       fill=RGBColor(0x0A, 0x33, 0x70))
        self._rect(s, x, y + h - 0.03, w, 0.03, fill=CYAN)
        if label:
            self._text(s, x + 0.24, y + h - 0.36, w - 0.48, 0.26,
                       [("Image slot: %s" % label, 9, False,
                         RGBColor(0x5A, 0x7C, 0xA8))])
        return None

    def cover(self, image, title_lines, strapline, date_line, status=None):
        s = self._slide()
        self._hero(s, image, 0, 0, SW, 5.80, "cover photograph")
        self._rect(s, 0, 5.80, SW, SH - 5.80, fill=WHITE)
        if self.client_logo and self.a(self.client_logo):
            self._rect(s, 0.30, 5.98, 1.62, 0.92, fill=NAVY)
            self._pic(s, self.client_logo, 0.34, 6.10, w=1.54)
        runs = [[(t, 20, True, NAVY)] for t in title_lines]
        self._text(s, 2.10, 5.96, 5.70, 0.64, runs, align=PP_ALIGN.CENTER,
                   space_after=1)
        self._text(s, 1.35, 6.66, 7.20, 0.24, [(strapline, 11, True, BODY)],
                   align=PP_ALIGN.CENTER)
        self._text(s, 1.35, 6.94, 7.20, 0.24,
                   [[(date_line, 9.5, True, RGBColor(0x8A, 0x8A, 0x8A))]],
                   align=PP_ALIGN.CENTER)
        self._rect(s, 0, 7.22, SW, 0.28, fill=NAVY)
        self._text(s, 0, 7.27, SW, 0.20,
                   [("Avia Solutions Limited   |   aviasolutions.com", 8.5, False,
                     RGBColor(0x9F, 0xC3, 0xE4))], align=PP_ALIGN.CENTER)
        if self.airline_logo and self.a(self.airline_logo):
            self._pic(s, self.airline_logo, 8.00, 6.06, w=1.70)
        else:
            self._placeholder(s, 8.00, 5.98, 1.70, 0.92, "airline logo")
        if status:
            self._rect(s, 0, 0, 2.3, 0.42, fill=RED)
            self._text(s, 0, 0.08, 2.3, 0.28, [(status, 11, True, WHITE)],
                       align=PP_ALIGN.CENTER)
        return s

    def divider(self, image, number, title, strapline=None):
        s = self._slide()
        self._hero(s, image, 0, 0, SW, SH, "section divider")
        self._rect(s, 0, 2.55, SW, 2.05, fill=NAVY).fill.fore_color.rgb = NAVY
        # translucency via a solid band keeps it print-safe
        self._text(s, M + 0.25, 2.78, 9.0, 0.5,
                   [("Section %s" % number, 17, True, RGBColor(0x7F, 0xC6, 0xF0))])
        self._text(s, M + 0.25, 3.24, 9.0, 0.7, [(title, 32, True, WHITE)])
        if strapline:
            self._text(s, M + 0.25, 3.96, 9.0, 0.42, [(strapline, 14, False,
                                                       RGBColor(0xD5, 0xE2, 0xF2))])
        self.n += 1
        self.contents.append((number, title, self.n))
        return s

    def content(self, title, subtitle=None):
        s = self._slide()
        self._chrome(s)
        n = len(title)
        size = 24 if n <= 44 else (21.5 if n <= 50 else (19 if n <= 57 else 17))
        self._text(s, M, 0.15, 7.90, 0.50, [(title, size, True, WHITE)])
        if subtitle:
            ssz = 14.5 if len(subtitle) <= 92 else 12.5
            self._text(s, M, 0.70, 7.90, 0.42,
                       [(subtitle, ssz, False, RGBColor(0xB9, 0xCD, 0xE5))])
        return s

    # -- content components --------------------------------------------------
    def table(self, s, x, y, w, headers, rows, col_w=None, total_row=False,
              size=9.0, hdr_size=9.0, row_h=0.235, hdr_h=0.44, aligns=None,
              hdr_wrap=True):
        ncol = len(headers)
        nrow = len(rows) + 1
        h = hdr_h + row_h * len(rows)
        gf = s.shapes.add_table(nrow, ncol, Inches(x), Inches(y),
                                Inches(w), Inches(h))
        tbl = gf.table
        tbl.first_row = False
        # kill the default banded style
        tblPr = tbl._tbl.find(qn('a:tblPr'))
        for tag in ('bandRow', 'firstRow'):
            if tblPr is not None and tblPr.get(tag) is not None:
                tblPr.set(tag, '0')
        if col_w:
            tot = sum(col_w)
            for i, cw in enumerate(col_w):
                tbl.columns[i].width = Emu(int(Inches(w) * cw / tot))
        tbl.rows[0].height = Inches(hdr_h)
        for i in range(len(rows)):
            tbl.rows[i + 1].height = Inches(row_h)

        def fmt(cell, text, sz, bold, colour, fill, align):
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = str(text)
            r.font.size = Pt(sz)
            r.font.bold = bold
            r.font.color.rgb = colour
            r.font.name = FONT

        if aligns is None:
            aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * (ncol - 1)
        for j, htxt in enumerate(headers):
            fmt(tbl.cell(0, j), htxt, hdr_size, True, WHITE, NAVY,
                PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)
        for i, row in enumerate(rows):
            last = total_row and i == len(rows) - 1
            fill = NAVY if last else (LIGHT if i % 2 else WHITE)
            col = WHITE if last else BODY
            for j, val in enumerate(row):
                fmt(tbl.cell(i + 1, j), val, size, last, col, fill, aligns[j])
        return gf

    def callout(self, s, x, y, w, h, text, size=13, fill=ORANGE, colour=NAVY,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE):
        self._rect(s, x, y, w, h, fill=fill, shape=shape)
        lines = text if isinstance(text, list) else [text]
        self._text(s, x + 0.14, y + 0.08, w - 0.28, h - 0.16,
                   [[(t, size, True, colour)] for t in lines],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=3)

    def stat(self, s, x, y, w, value, label, vsize=30, lsize=11,
             vcol=NAVY, lcol=BODY, align=PP_ALIGN.LEFT):
        self._text(s, x, y, w, 0.52, [(value, vsize, True, vcol)], align=align)
        self._text(s, x, y + 0.50, w, 0.50, [(label, lsize, False, lcol)], align=align)

    def keynumbers(self, s, items, x=M, y=1.44, w=9.4, circle=0.78, gap=0.05):
        """items: list of (circle_text_lines, statement)."""
        rowh = (BOTTOM - y) / len(items)
        for i, (circ, stmt) in enumerate(items):
            cy = y + i * rowh + (rowh - circle) / 2
            self._rect(s, x + 0.12, cy, circle, circle, fill=NAVY,
                       shape=MSO_SHAPE.OVAL)
            lines = circ if isinstance(circ, list) else [circ]
            sz = 11 if max(len(l) for l in lines) <= 7 else 9
            self._text(s, x + 0.02, cy + circle / 2 - 0.11 * len(lines),
                       circle + 0.20, circle,
                       [[(l, sz, True, WHITE)] for l in lines],
                       align=PP_ALIGN.CENTER, space_after=0)
            self._text(s, x + 1.18, cy + 0.06, w - 1.30, circle,
                       [(stmt, 15, True, BODY)], anchor=MSO_ANCHOR.MIDDLE)

    def bullets(self, s, x, y, w, items, size=12, colour=BODY, dot=CYAN,
                gap=0.05, bold_lead=True):
        cy = y
        for it in items:
            if isinstance(it, tuple):
                lead, rest = it
            else:
                lead, rest = None, it
            self._rect(s, x, cy + 0.055, 0.115, 0.115, fill=dot,
                       shape=MSO_SHAPE.OVAL)
            runs = []
            if lead:
                runs.append((lead, size, True, colour))
                runs.append((" " + rest, size, False, colour))
            else:
                runs.append((rest, size, False, colour))
            est = 0.02 + 0.185 * (1 + (len(lead or "") + len(rest)) // int(w * 13.5))
            self._text(s, x + 0.24, cy, w - 0.24, est, [runs])
            cy += est + gap
        return cy

    def panel(self, s, x, y, w, h, title, items, fill=TEAL, size=12):
        self._rect(s, x, y, w, h, fill=fill)
        cy = y + 0.16
        if title:
            tsz = 14 if len(title) <= int((w - 0.4) * 8.2) else 12
            lines = 1 + len(title) // max(1, int((w - 0.4) * (7.6 if tsz == 14 else 8.8)))
            self._text(s, x + 0.18, cy, w - 0.36, 0.30 * lines,
                       [(title, tsz, True, WHITE)])
            cy += 0.30 * lines + 0.12
        for it in items:
            est = 0.02 + 0.18 * (1 + len(it) // int((w - 0.5) * 12.5))
            self._rect(s, x + 0.18, cy + 0.05, 0.10, 0.10,
                       fill=RGBColor(0x7F, 0xC6, 0xF0), shape=MSO_SHAPE.OVAL)
            self._text(s, x + 0.38, cy, w - 0.56, est, [(it, size, False, WHITE)])
            cy += est + 0.06

    def quotecard(self, s, x, y, w, h, quote, attrib):
        self._rect(s, x, y, w, h, fill=WHITE, line=NAVY, line_w=1.25,
                   shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        self._text(s, x + 0.16, y + 0.12, w - 0.32, h - 0.44,
                   [[('"%s"' % quote, 11.5, True, BODY)]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._text(s, x + 0.16, y + h - 0.32, w - 0.32, 0.24,
                   [(attrib, 9, False, GREY)], align=PP_ALIGN.CENTER)

    def methodology(self, s, paragraphs, x=M, y=1.44, w=9.4):
        cy = y
        for head, body in paragraphs:
            if head:
                self._text(s, x, cy, w, 0.30, [(head, 14, True, NAVY)])
                cy += 0.32
            est = 0.04 + 0.175 * (1 + len(body) // int(w * 13.0))
            self._text(s, x, cy, w, est, [(body, 11.5, False, BODY)],
                       line_spacing=1.05)
            cy += est + 0.14
        return cy

    # ------------------------------------------------------------------
    def contents_slide(self, after_index=1):
        pass  # contents built in the content script for ordering control

    def save(self, path, author="Avia Solutions", title=None, subject=None):
        cp = self.prs.core_properties
        cp.author = author
        cp.last_modified_by = author
        cp.title = title or self.deck_title
        cp.subject = subject or ""
        cp.category = "Route forecast pitch"
        cp.comments = "Avia Solutions Limited. All rights reserved."
        self._set_language()
        self.prs.save(path)
        return path

    def _set_language(self):
        """Force en-GB at document default and strip any other run language."""
        for part in list(self.prs.part.package.iter_parts()):
            try:
                el = part._element
            except AttributeError:
                continue
            if el is None:
                continue
            for rpr in el.iter():
                tag = rpr.tag
                if tag in (qn('a:rPr'), qn('a:defRPr'), qn('a:endParaRPr')):
                    rpr.set('lang', 'en-GB')
                    rpr.set('altLang', 'en-GB')
                    if rpr.get('noProof'):
                        del rpr.attrib['noProof']
