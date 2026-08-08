"""Project Liguria: Genoa Cristoforo Colombo - New York route business case.

Second worked deck from the Avia house-style library. Deliberately a harder case
than Project Redwood: a secondary European airport with no long-haul history, a
thin licensable image pool, no published leakage study, and a set of failed
comparables that the deck presents as its own material.

Run: python3 build_goa_nyc.py
"""

import os
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from avia_deck import (AviaDeck, NAVY, BODY, ORANGE, CYAN, WHITE, GREY, LIGHT,
                       MIDBLUE, TEAL, RED, M, SW, SH)
import goa_figures as gf

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "Project_Liguria_GOA-NYC_Route_Business_Case.pptx")

CODENAME = "Project Liguria"
BASIS = ("Source: AviaSolutions analysis, Genoa - New York business case, central "
         "planning case. Route economics in US dollars; Italian economic and traffic "
         "figures in euro or as published.")

# --- the business case, read from Genoa_NYC_business_case.xlsx --------------
CASE = {
    "capture": 0.65, "freq": 7, "lf_cap": 0.85,
    "fare_y": 345, "fare_j": 1400, "fuel": 0.90, "y_share": 0.80,
    "catchment_pop": 19536487, "ny_market": 828760.74,
    "natural_catchment": 138608.178, "carried_today": 9043.41,
    "seats_y": 162, "seats_j": 20, "burn": 2500, "bh_turn": 18,
    "demand_each_way": 93260.509, "pax_turn": 309.4,
    "rev": {"Economy": 95013.0, "Business": 47600.0, "Cargo": 7259.84},
    "net_rev": 149872.84, "charges_recovery": 14340.69, "gross_rev": 164213.53,
    "var": [("Fuel", 40500.0), ("Maintenance", 14740.385), ("Catering", 819.291),
            ("Landing", 4429.0), ("Passenger airport charges", 15934.1),
            ("En-route navigation", 1134.999), ("Ground handling", 5871.0)],
    "fixed": [("Ownership", 16880.516), ("Insurance", 2156.682), ("Crew", 21600.0)],
    "indirect": [("Admin, 5% of net revenue", 7493.642),
                 ("Sales, 5% of net revenue", 7493.642)],
    "var_total": 83428.775, "fixed_total": 40637.198, "indirect_total": 14987.284,
    "total_cost": 139053.257, "profit_turn": 25160.273, "margin": 0.153,
    "breakeven_lf": 0.689, "annual_profit": 9158339.372, "annual_pax": 112621.6,
    "standalone_profit": 5127639.243,
    "block_hours": 6552, "aircraft": 2, "util_per_ac": 3276, "util_vs_full": 0.604,
    "spare_bh": 4298, "spare_turns": 238.778,
    "own_charged": 6144507.871, "own_dedicated": 10175208.0,
    "season_idx": [0.527, 0.448, 0.701, 1.173, 1.128, 1.264, 1.381, 1.637,
                   1.260, 1.047, 0.689, 0.747],
    "season_a": 9183500, "season_b": 7759505, "season_c": 8868071,
    "fuel_grid": [0.68, 0.80, 0.90, 1.00, 1.10],
    "profit_grid": [12.8, 10.8, 9.2, 7.5, 5.9],
    "network": [("GOA-JFK", "120 turns/yr", 33446.4, 18018925.44, 15147433.963, 2871491.477, 0.159, 0.398),
                ("GOA-EWR", "90 turns/yr", 25084.8, 13517927.712, 11306383.533, 2211544.179, 0.164, 0.301),
                ("GOA-BOS", "60 turns/yr", 16723.2, 9003240.0, 7329368.74, 1673871.26, 0.186, 0.198)],
    "network_total": (75254.4, 40540093.152, 33783186.236, 6756906.916, 0.167, 0.897),
}


def fmt(n, dp=0):
    return "{:,.{}f}".format(n, dp)


def pct(x, dp=1):
    return "{:.{}f}%".format(x * 100, dp)


def build_figures():
    p = lambda n: os.path.join(ASSETS, n)
    gf.route_map(p("goa_map_route.png"))
    gf.catchment_map(p("goa_map_catchment.png"))
    gf.traffic(p("goa_ch_traffic.png"))
    gf.seasonality(p("goa_ch_season.png"), CASE["season_idx"])
    gf.cost_stack(p("goa_ch_cost.png"),
                  CASE["var"] + CASE["fixed"] + CASE["indirect"],
                  CASE["gross_rev"])
    gf.fuel_sensitivity(p("goa_ch_fuel.png"), CASE["fuel_grid"], CASE["profit_grid"])


def main(pages=None):
    build_figures()
    d = AviaDeck(deck_title="%s - Genoa to New York route business case" % CODENAME,
                 event_line="World Routes 2026", assets_dir=ASSETS,
                 client_logo=None, airline_logo=None)

    # ------------------------------------------------------------ 1 cover
    d.cover("goa_cover.jpg",
            ["A Direct Link Between Genoa",
             "and New York"],
            "Prepared for Aeroporto di Genova and Regione Liguria",
            "World Routes 2026   |   5 August 2026   |   %s   |   Commercial in Confidence" % CODENAME,
            status="DRAFT")

    # --------------------------------------------------------- 2 contents
    s = d.content("Contents", "A Direct Link Between Genoa and New York")
    titles = ["The opportunity in one slide", "Genoa and its catchment",
              "The Genoa - New York corridor", "The Italian secondary-airport wave",
              "Route economics", "The counter-case",
              "Methodology and track record", "The airport, the aircraft and the ask",
              "Appendix"]
    pg = pages or [0] * len(titles)
    y = 1.56
    for i, t in enumerate(titles):
        page = pg[i] if i < len(pg) else 0
        d._rect(s, M + 0.10, y, 0.50, 0.50, fill=NAVY, shape=MSO_SHAPE.OVAL)
        d._text(s, M + 0.10, y + 0.135, 0.50, 0.30, [(str(i + 1), 13.5, True, WHITE)],
                align=PP_ALIGN.CENTER)
        d._text(s, M + 0.84, y + 0.09, 7.4, 0.34, [(t, 15.5, True, BODY)])
        d._text(s, 8.60, y + 0.09, 1.0, 0.34,
                [(str(page) if page else "-", 15.5, True, MIDBLUE)],
                align=PP_ALIGN.RIGHT)
        d._rect(s, M + 0.84, y + 0.53, 7.75, 0.012, fill=LIGHT)
        y += 0.615

    # ------------------------------------------------- 3 the case in one slide
    s = d.content("The case, and the case against",
                  "Both halves on one slide, because the airline will build the second half anyway")
    fors = [("An existing corridor, in freight", "Circa 37% of Genoa and Savona port traffic is United States linked, circa 336,000 container units a year, and MSC runs a direct nine-day Genoa to New York container service. The commercial corridor exists; the air link does not."),
            ("The fastest-growing of the three", "Genoa grew 18.1% to 1,587,761 passengers in 2025, against Malpensa's 8.6% and Fiumicino's 4.5%, on the regulator's own comparable table."),
            ("A gateway that already moves visitors", "1,630,593 cruise passengers used Genoa in 2025, more than the airport handled, and the cruise sector is worth EUR 346m to the Genoa and Savona economy."),
            ("The aircraft now reaches", "Genoa to New York is 3,509.6 nautical miles, circa 75% of the A321XLR's published 4,700 nm range, and shorter than Madrid to Boston, which has operated since 14 November 2024.")]
    againsts = [("The market has not grown", "The United States dropped out of Italy's top five international country markets in 2025, replaced by Poland. The secondary-airport route boom is redistribution, not growth."),
                ("The closest comparables failed", "SAS's Gothenburg to Newark and Aalborg to Newark A321LR services were cut. Aer Lingus is cutting Dublin to Minneapolis on the XLR after load factors bottoming at 30.0%."),
                ("There is no leakage study", "None exists, from ENAC, Assaeroporti, Regione Liguria, the chamber of commerce or academia. The catchment case rests on classification and drive times, not measurement.")]
    d._rect(s, M, 1.40, 4.62, 0.36, fill=NAVY)
    d._text(s, M + 0.14, 1.46, 4.4, 0.26, [("The case for", 12.5, True, WHITE)])
    yy = 1.86
    for h, b in fors:
        d._rect(s, M, yy, 4.62, 1.16, fill=LIGHT)
        d._rect(s, M, yy, 0.05, 1.16, fill=ORANGE)
        d._text(s, M + 0.16, yy + 0.08, 4.32, 0.26, [(h, 11.5, True, NAVY)])
        d._text(s, M + 0.16, yy + 0.34, 4.32, 0.76, [(b, 9.6, False, BODY)])
        yy += 1.24
    d._rect(s, 5.10, 1.40, 4.62, 0.36, fill=RED)
    d._text(s, 5.24, 1.46, 4.4, 0.26, [("The case against", 12.5, True, WHITE)])
    yy = 1.86
    for h, b in againsts:
        d._rect(s, 5.10, yy, 4.62, 1.16, fill=RGBColor(0xF7, 0xE9, 0xE7))
        d._rect(s, 5.10, yy, 0.05, 1.16, fill=RED)
        d._text(s, 5.26, yy + 0.08, 4.32, 0.26, [(h, 11.5, True, NAVY)])
        d._text(s, 5.26, yy + 0.34, 4.32, 0.76, [(b, 9.6, False, BODY)])
        yy += 1.24
    d.callout(s, 5.10, 5.58, 4.62, 1.10,
              ["Our position: the route works on the",
               "economics we can model, and the open",
               "items are named, not smoothed over"], size=11.5, fill=NAVY, colour=WHITE)
    d.source(s, "Sources: Autorita di Sistema Portuale del Mar Ligure Occidentale via Genova24, October 2024; ENAC Air Traffic Data 2025, January 2026; "
                "Stazioni Marittime SpA via ANSA, 13 January 2026; Airbus published range; US DOT and carrier statements for the comparables. "
                "Full citations in the evidence pack.", size=7.0)

    # ------------------------------------------- 4 summary of route forecast
    s = d.content("Summary of the proposition",
                  "Daily Airbus A321XLR, Genoa to New York, central planning case")
    d._pic_cover(s, d.a("goa_map_route.png"), 0.0, 1.18, SW, 3.62)
    d._text(s, 0.34, 1.34, 3.3, 0.28,
            [("Genoa natural New York catchment", 11.5, True, NAVY)])
    d._text(s, 0.34, 1.58, 3.3, 0.44,
            [(fmt(CASE["natural_catchment"]), 25, True, MIDBLUE)])
    d._text(s, 6.44, 1.34, 3.2, 0.28,
            [("Carried through Genoa today", 11.5, True, NAVY)], align=PP_ALIGN.RIGHT)
    d._text(s, 6.44, 1.58, 3.2, 0.44,
            [(fmt(CASE["carried_today"]), 25, True, MIDBLUE)], align=PP_ALIGN.RIGHT)
    d._text(s, 6.14, 3.66, 3.5, 0.28,
            [("Forecast passengers, year 1", 11.5, True, NAVY)], align=PP_ALIGN.RIGHT)
    d._text(s, 6.14, 3.90, 3.5, 0.44,
            [(fmt(CASE["annual_pax"]), 25, True, MIDBLUE)], align=PP_ALIGN.RIGHT)
    rows = [["GOA-JFK", "A321XLR", "Daily", "182", str(CASE["freq"]),
             fmt(CASE["seats_j"]), fmt(CASE["seats_y"]),
             fmt(CASE["annual_pax"]), pct(CASE["lf_cap"], 0)]]
    d._text(s, M, 4.92, 4.0, 0.26,
            [("Schedule and configuration", 12, True, NAVY)])
    d.table(s, M, 5.18, 8.90,
            ["Sector", "Aircraft", "Op. days", "Total seats", "Weekly freq.",
             "Business", "Economy", "Annual pax yr 1", "Planning LF"],
            rows, col_w=[1.0, 1.0, 0.85, 0.9, 0.95, 0.85, 0.85, 1.1, 0.95],
            size=10, hdr_size=8.6, row_h=0.34, hdr_h=0.50)
    d.callout(s, M, 5.98, 4.34, 0.66,
              "Route margin %s at a breakeven load factor of %s" % (
                  pct(CASE["margin"], 1), pct(CASE["breakeven_lf"], 1)), size=12)
    d.callout(s, 5.38, 5.98, 4.34, 0.66,
              "Annual profit $%.1fm on the network basis" % (CASE["annual_profit"] / 1e6),
              size=12, fill=NAVY, colour=WHITE)
    d.source(s, BASIS + " Catchment demand from AviaSolutions catchment allocation on GeoNames population and OSRM road times; "
                "New York market split and fares from Sabre ODPOO 2024. Distance 3,509.6 nm.", size=7.0)

    # ======================================================= SECTION 2
    d.divider("goa_div_city.jpg", "2", "Genoa and its catchment",
              "The airport, the region and the airports that take its traffic today")

    s = d.content("Genoa in six numbers",
                  "The regulator's own figures, and the two that are not yet measured")
    d.keynumbers(s, [
        (["1.59m"], "Passengers in 2025, up 18.1%, the fastest growth of Italy's three north-western gateways"),
        (["19th"], "Rank among the 44 Italian airports open to commercial traffic"),
        (["1.63m"], "Cruise passengers through Genoa in 2025, more than the airport itself handled"),
        (["37%"], "Share of Genoa and Savona port traffic that is United States linked"),
        (["Zero"], "Long-haul or transatlantic services at Genoa today, from any airline"),
        (["None"], "Published studies measuring how much Ligurian demand leaks to Milan, Turin, Nice or Pisa"),
    ])
    d.source(s, "Sources: ENAC, Air Traffic Data 2025 - Executive Summary, Airport Traffic Table, published January 2026; "
                "Stazioni Marittime SpA reported by ANSA, 13 January 2026; Autorita di Sistema Portuale del Mar Ligure Occidentale "
                "reported by Genova24, October 2024; AviaSolutions route survey, 5 August 2026.", size=7.5)

    s = d.content("Genoa is growing faster than Milan or Rome",
                  "Verified data points only; the unverified back series is deliberately omitted")
    d._pic(s, "goa_ch_traffic.png", M, 1.40, w=5.86)
    rows = [["Genoa", "1,587,761", "+18.1%", "19th"],
            ["Milan Malpensa", "not stated here", "+8.6%", "2nd"],
            ["Rome Fiumicino", "not stated here", "+4.5%", "1st"]]
    d.table(s, 6.32, 1.44, 3.40, ["Airport", "2025 passengers", "vs 2024", "Rank"],
            rows, col_w=[1.15, 1.10, 0.65, 0.50], size=8.6, hdr_size=8.0,
            row_h=0.46, hdr_h=0.44,
            aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])
    d.panel(s, 6.32, 3.36, 3.40, 1.94, "One basis, labelled", [
        "ENAC reports 1,587,761 for 2025; the airport operator's own figure is 1,577,159.",
        "The gap is scope. We use the ENAC basis because it is comparable across all 44 Italian airports.",
    ], size=10)
    d.panel(s, 6.32, 5.42, 3.40, 1.60, "What we will not print", [
        "The 2015 to 2023 annual series could not be tied to a named publisher. It is omitted rather than estimated.",
    ], size=10, fill=RED)
    d.source(s, "Sources: ENAC, Air Traffic Data 2025 - Executive Summary, Airport Traffic Table, published January 2026, for all three airports' "
                "2025 growth and rank. Absolute Malpensa and Fiumicino totals are not quoted because they were not taken from the same table in "
                "this pass.", size=7.5)

    s = d.content("The catchment, and what it is not",
                  "Six competing airports, all of them between one and two and a half hours away")
    d._pic(s, "goa_map_catchment.png", M, 1.40, h=5.16)
    rows = [["Milan Malpensa", "180-188 km", "1h58-2h06", "circa 2h20-2h35, two legs"],
            ["Milan Linate", "circa 151 km", "1h43", "circa 2h00-2h30, two legs"],
            ["Milan Bergamo", "190-191 km", "2h07-2h22", "not found"],
            ["Turin city", "171 km", "1h53", "1h41-1h49 direct"],
            ["Nice", "195-200 km", "2h22", "circa 3h, 3 direct trains"],
            ["Pisa", "158-176 km", "1h50-2h09", "circa 2h09"]]
    d.table(s, 5.06, 1.40, 4.66,
            ["Airport", "Road", "Drive", "Rail"], rows,
            col_w=[1.25, 0.95, 0.85, 1.61], size=8.6, hdr_size=8.2, row_h=0.40,
            hdr_h=0.42, aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT,
                                PP_ALIGN.LEFT])
    d.panel(s, 5.06, 4.30, 4.66, 1.42, "The regulator's own classification", [
        "The Piano Nazionale degli Aeroporti 2026-2035 places Genoa in the same Nord-Ovest traffic basin as Malpensa, Linate, Bergamo, Turin and Cuneo.",
    ], size=10)
    d.panel(s, 5.06, 5.84, 4.66, 1.18, "Terzo Valico cuts both ways", [
        "The Genoa to Milan high-speed link, targeted at circa 1h07 from 2027, shortens the journey in both directions.",
    ], size=10, fill=RED)
    d.source(s, "Sources: Ministero delle Infrastrutture e dei Trasporti, Piano Nazionale degli Aeroporti 2026-2035, 15 July 2026; drive and rail "
                "times from routing and ticketing aggregators (Himmera, Omio, Trainline, calcolopercorso.it), retrieved 5 August 2026, all estimates; "
                "Terzo Valico timing per RFI reported by TrasportoEuropa and Regione Liguria, July 2026, a forecast.", size=7.0)

    s = d.content("Where Ligurians fly to the United States from today",
                  "The honest answer: nobody has measured it, and that is a scope of work")
    d.methodology(s, [
        ("What was checked",
         "ENAC and Assaeroporti data isolating Ligurian-resident origin passengers; a Regione Liguria or Camera di Commercio di Genova leakage "
         "study; academic and consultancy catchment work for Genoa. None of the three exists in the public record."),
        ("What the catchment case therefore rests on",
         "First, the national regulator's own classification of Genoa as sharing a traffic basin with Milan. Second, the drive and rail time "
         "comparison on the previous slide. Third, the structural fact that Genoa has no United States service at all, so every Ligurian "
         "transatlantic passenger today departs from somewhere else."),
        ("What it does not rest on",
         "A published leakage study, because none exists. A deck that implied otherwise would be misrepresenting the evidence."),
        ("What would close it",
         "A Sabre origin-destination pull keyed to Ligurian postcodes, or the DB1B-equivalent on the United States end, sized against the "
         "four-region population. That is a fortnight of work and it is the single highest-value thing anyone could commission before this "
         "route is taken to an airline."),
    ], y=1.44, w=5.90)
    d.panel(s, 6.24, 1.44, 3.48, 2.66, "The four-region population", [
        "Liguria 1.51m, Piemonte 4.25m, Lombardia 10.03m, Emilia-Romagna 4.49m: circa 20.3m.",
        "This is a four-region total, not a drive-time isochrone, and we label it as one.",
    ], size=10)
    d.panel(s, 6.24, 4.24, 3.48, 2.78, "Why we say so plainly", [
        "Presenting 20.3m as Genoa's catchment without that caveat is the claim that gets a deck dismissed by a network planner, because "
        "Malpensa sits inside the same four regions with a better road position.",
        "The forecast in section 5 uses the modelled catchment allocation, not the four-region total.",
    ], size=10, fill=TEAL)
    d.source(s, "Sources: ISTAT via tuttitalia.it and Regione Emilia-Romagna Servizio Statistica for population, reference dates between "
                "31 December 2023 and 1 January 2026; four-region total is an AviaSolutions calculation. Absence of a leakage study confirmed "
                "by search of ENAC, Assaeroporti, Regione Liguria and Camera di Commercio di Genova publications, 5 August 2026.", size=7.0)

    # ======================================================= SECTION 3
    d.divider("goa_div_port.jpg", "3", "The Genoa - New York corridor",
              "Freight, cruise, diaspora and corporate links that already run between the two cities")

    s = d.content("The corridor already exists, in freight",
                  "A quantified Genoa to New York commercial relationship with no air link")
    rows = [["Genoa and Savona port traffic that is United States linked", "circa 37%", "Autorita di Sistema Portuale via Genova24, October 2024"],
            ["Container units a year on that United States trade", "circa 336,000", "As above"],
            ["MSC Genoa to New York container service", "Direct, nine-day transit", "TrasportoEuropa (MEDUSEC service)"],
            ["Genoa cruise passengers, 2025", "1,630,593", "Stazioni Marittime SpA via ANSA, 13 January 2026"],
            ["Cruise sector value to the Genoa and Savona economy", "EUR 346m", "MedCruise and Ports of Genoa impact study"],
            ["ERG, Genoa headquartered, United States renewables acquisition", "USD 270m for 75% of a 317 MW platform", "ERG press release, 21 December 2023"]]
    d.table(s, M, 1.44, 9.40, ["Measure", "Figure", "Source"], rows,
            col_w=[3.9, 2.2, 3.3], size=9.5, hdr_size=9.0, row_h=0.54, hdr_h=0.44,
            aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
    d.callout(s, M, 4.86, 4.58, 0.90,
              ["The freight corridor is direct and scheduled.", "The passenger corridor is not."], size=12)
    d.panel(s, 5.12, 4.86, 4.60, 0.90, None, [
        "MSC is headquartered in Geneva, not Genoa. We do not describe it as a Genoese company.",
    ], size=9.5, fill=RED)
    d.bullets(s, M, 5.94, 9.40, [
        ("The Austrian Airlines proof point.", "Austrian flies a dedicated weekly Vienna to Genoa charter from 2 May to 31 October 2026 purely to feed the Costa cruise terminal. That is the demand mechanism this pitch depends on, already operating at short-haul scale."),
    ], size=11)
    d.source(s, "Sources as shown in the table; Austrian Airlines charter per GenovaToday, retrieved 5 August 2026 (secondary). "
                "Fincantieri's United States Navy work is deliberately not cited: four of the six remaining Constellation-class frigates were "
                "cancelled in December 2025.", size=7.0)

    s = d.content("The New York end of the demand",
                  "A large Italian-American base, and a declining one, stated as it is")
    rows = [["Italian ancestry, New York State and New Jersey combined", "circa 3.44m", "ACS 2020-2024 5-year table B04006 (secondary source, sum derived)"],
            ["Trend", "Declining circa 12% per decade on the New Jersey ACS evidence", "ACS, as above"],
            ["Ligurian or Genoese ancestry specifically", "No data exists at any level", "Checked at ACS and AIRE, 5 August 2026"],
            ["The \"two million emigrants sailed from Genoa\" claim", "Historically misleading: they emigrated through Genoa, mostly from the south", "AviaSolutions review of the historical record"]]
    d.table(s, M, 1.44, 6.10, ["Measure", "Figure", "Source and status"], rows,
            col_w=[2.4, 1.7, 2.0], size=9.2, hdr_size=8.8, row_h=0.66, hdr_h=0.44,
            aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
    d.panel(s, 6.56, 1.44, 3.16, 2.70, "How we use it", [
        "The visiting friends and relatives case is built on the pan-Italian-American population, which is large and Census-derived.",
        "It is not built on a Ligurian diaspora, because no such data exists.",
    ], size=10)
    d.panel(s, 6.56, 4.28, 3.16, 2.74, "Before publication", [
        "The 3.44m figure comes from a secondary aggregator of the ACS table and must be verified directly on data.census.gov before it goes to a client.",
    ], size=10, fill=RED)
    d.methodology(s, [
        ("Why this matters more than it looks",
         "Visiting friends and relatives traffic is the ballast under a thin long-haul route: it books early, it travels in the shoulder months "
         "and it is far less fare-elastic than leisure. It is also the flow most often overstated in route pitches, which is why the number here "
         "is the pan-Italian one and the Ligurian claim is left out."),
    ], y=3.78, w=6.10)
    d.source(s, "Source: US Census American Community Survey 2020-2024 5-year estimates, table B04006, accessed via a secondary aggregator; "
                "verification on data.census.gov outstanding. AviaSolutions analysis.", size=7.5)

    # ======================================================= SECTION 4
    d.divider("goa_div_coast.jpg", "4", "The Italian secondary-airport wave",
              "What has launched since 2022, what survived, and what it actually proves")

    s = d.content("Secondary Italian airports now carry United States service",
                  "Four carriers, seven or more nonstops, from cities smaller than Genoa's catchment")
    rows = [["Naples", "American Airlines, United, Delta", "Philadelphia, Newark, New York JFK", "Operating"],
            ["Catania", "Delta, Neos", "New York JFK", "Operating"],
            ["Palermo", "United", "Newark", "Operating, circa 85% load factor, extended through 16 December 2026"],
            ["Bari", "Neos, then United", "New York JFK, Newark", "Neos cut for 2026; United backfilled"],
            ["Olbia", "United", "Newark", "Operating"],
            ["Bologna", "American Airlines", "Philadelphia", "Died after one 2019 season"],
            ["Milan Malpensa", "La Compagnie", "Newark, all-business A321neo, 5x weekly since 15 April 2022", "Operating"]]
    d.table(s, M, 1.44, 9.40, ["Italian airport", "Carrier", "United States point", "Status"],
            rows, col_w=[1.5, 2.2, 3.0, 3.1], size=9.5, hdr_size=9.0, row_h=0.48,
            hdr_h=0.44, aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT,
                                PP_ALIGN.LEFT])
    d.callout(s, M, 5.10, 4.58, 0.88,
              ["La Compagnie has flown a narrowbody,", "premium-configured northern Italy to",
               "New York route for four years"], size=11)
    d.callout(s, 5.12, 5.10, 4.60, 0.88,
              ["United's Newark to Palermo runs at", "circa 85% load factor and has been",
               "extended into December 2026"], size=11, fill=NAVY, colour=WHITE)
    d.bullets(s, M, 6.16, 9.40, [
        ("Read the failures as well.", "Neos lost Bari and American lost Bologna. Both are in this table on purpose: the wave has casualties and the next slide sets out what that means."),
    ], size=11)
    d.source(s, "Sources: carrier route announcements and Italian and United States aviation trade press, retrieved 5 August 2026, secondary "
                "throughout; Palermo load factor per Travel And Tour World and Dream of Italy, January 2026; La Compagnie per carrier reporting "
                "via TourMag and Forbes. To be replaced with an OAG schedule extract before issue.", size=7.0)

    s = d.content("What the wave proves, and what it does not",
                  "The reframing that has to happen before this route is taken to an airline")
    d.methodology(s, [
        ("It proves that a secondary Italian city can hold a United States nonstop",
         "Naples, Catania, Palermo, Bari and Olbia now do. The old objection, that only Rome and Milan can carry transatlantic service from "
         "Italy, has been answered by the market rather than by argument."),
        ("It does not prove the market is growing",
         "ENAC's 2025 data shows the United States dropping out of Italy's top five international country markets, replaced by Poland. The "
         "United States National Travel and Tourism Office puts Italian visitation to the United States at 103% of 2019 in 2024, which is flat "
         "across five years. The wave is redistribution between Italian airports, not new demand."),
        ("Which is, on balance, the stronger argument for Genoa",
         "If the pool is fixed and it is being redistributed towards whichever secondary city puts an aircraft on the ground, then the case for "
         "Genoa is a share argument, not a growth argument, and share arguments are settled by who moves first. Genoa is the largest north-western "
         "Italian catchment with no United States service and no announced plan."),
        ("What that costs us",
         "It means the route cannot be underwritten on market growth. It has to be underwritten on capture, and capture is exactly the thing "
         "no one has yet measured for Liguria. The economics in section 5 are therefore run across a capture range rather than at a point."),
    ], y=1.46, w=9.40)
    d.source(s, "Sources: ENAC, Air Traffic Data 2025, published January 2026; United States National Travel and Tourism Office arrivals data, "
                "2024 against 2019. AviaSolutions analysis.", size=7.5)

    # ======================================================= SECTION 5
    d.divider("goa_div_aircraft.jpg", "5", "Route economics",
              "The route profit and loss, the fleet question, seasonality and the sensitivities")

    s = d.content("Route profit and loss per turnaround",
                  "Central planning case, US dollars per return rotation, daily A321XLR")
    rows = []
    for k, v in CASE["rev"].items():
        rows.append([k, "$%s" % fmt(v)])
    rows.append(["Net revenue", "$%s" % fmt(CASE["net_rev"])])
    rows.append(["Charges recovery", "$%s" % fmt(CASE["charges_recovery"])])
    rows.append(["Gross revenue", "$%s" % fmt(CASE["gross_rev"])])
    d.table(s, M, 1.42, 4.46, ["Revenue", "Per turnaround"], rows,
            col_w=[2.9, 1.56], size=9.5, hdr_size=9.0, row_h=0.36, hdr_h=0.40,
            total_row=True)
    crows = [[n, "$%s" % fmt(v)] for n, v in CASE["var"]]
    crows.append(["Variable cost", "$%s" % fmt(CASE["var_total"])])
    crows += [[n, "$%s" % fmt(v)] for n, v in CASE["fixed"]]
    crows.append(["Direct fixed", "$%s" % fmt(CASE["fixed_total"])])
    crows.append(["Indirect fixed, admin and sales", "$%s" % fmt(CASE["indirect_total"])])
    crows.append(["Total cost", "$%s" % fmt(CASE["total_cost"])])
    d.table(s, 5.26, 1.42, 4.46, ["Cost", "Per turnaround"], crows,
            col_w=[2.9, 1.56], size=9.0, hdr_size=9.0, row_h=0.315, hdr_h=0.40,
            total_row=True)
    d.callout(s, M, 5.36, 3.02, 0.80,
              ["Profit per turnaround", "$%s" % fmt(CASE["profit_turn"])], size=12)
    d.callout(s, 3.50, 5.36, 3.02, 0.80,
              ["Route margin", pct(CASE["margin"], 1)], size=12, fill=NAVY, colour=WHITE)
    d.callout(s, 6.70, 5.36, 3.02, 0.80,
              ["Breakeven load factor", pct(CASE["breakeven_lf"], 1)], size=12)
    d.panel(s, M, 6.30, 9.40, 0.74, None, [
        "Airport charges at both ends are indicative placeholders and are not yet verified against the Genoa and New York published tariffs. "
        "Crew at $1,200 per block hour is a low-cost-carrier judgement pending a citation. Both are in the assumptions register.",
    ], size=9.5, fill=RED)
    d.source(s, BASIS + " Maintenance from the validated Airbus 2024-2025 reserves; ownership from appraiser lease rates; "
                "fuel at $0.90 per kilogramme through-cycle planning assumption.", size=7.0)

    s = d.content("The result, and the fleet question behind it",
                  "The headline profit assumes the aircraft is busy when it is not flying to New York")
    d._pic(s, "goa_ch_cost.png", M, 1.40, w=5.86)
    rows = [["Network basis, fleet kept busy", "$%.2fm" % (CASE["annual_profit"] / 1e6)],
            ["Standalone, Genoa-only fleet", "$%.2fm" % (CASE["standalone_profit"] / 1e6)],
            ["Difference", "$%.2fm" % ((CASE["annual_profit"] - CASE["standalone_profit"]) / 1e6)],
            ["Ownership charged in the headline", "$%.2fm" % (CASE["own_charged"] / 1e6)],
            ["True ownership if dedicated", "$%.2fm" % (CASE["own_dedicated"] / 1e6)],
            ["Annual block hours flown", fmt(CASE["block_hours"])],
            ["Aircraft required", str(CASE["aircraft"])],
            ["Utilisation against a full year", pct(CASE["util_vs_full"], 1)],
            ["Spare block hours a year", fmt(CASE["spare_bh"])]]
    d.table(s, 6.32, 1.42, 3.40, ["Measure", "Value"], rows, col_w=[2.35, 1.05],
            size=8.8, hdr_size=8.4, row_h=0.42, hdr_h=0.40)
    d.bullets(s, M, 5.16, 5.86, [
        ("Read the difference, not the headline.", "The route carries $%.1fm of annual profit when the aircraft is kept busy elsewhere and $%.1fm when it is not. The gap, $%.1fm, is ownership recovery." % (
            CASE["annual_profit"] / 1e6, CASE["standalone_profit"] / 1e6,
            (CASE["annual_profit"] - CASE["standalone_profit"]) / 1e6)),
        ("The route is a fleet decision.", "At %s utilisation against a full year, an operator taking this route on a dedicated aircraft is buying %s spare block hours it has to place somewhere." % (
            pct(CASE["util_vs_full"], 1), fmt(CASE["spare_bh"]))),
    ], size=10.5)
    d.source(s, BASIS, size=7.5)

    s = d.content("Three United States points from one Genoa base",
                  "The network case: the same aircraft, three markets, one summer programme")
    rows = []
    for name, prog, pax, rev, cost, prof, marg, frames in CASE["network"]:
        rows.append([name, "A321XLR", prog, fmt(pax), "$%s" % fmt(rev),
                     "$%s" % fmt(cost), "$%s" % fmt(prof), pct(marg, 1),
                     "%.2f" % frames])
    t = CASE["network_total"]
    rows.append(["Network total", "", "", fmt(t[0]), "$%s" % fmt(t[1]),
                 "$%s" % fmt(t[2]), "$%s" % fmt(t[3]), pct(t[4], 1), "%.2f" % t[5]])
    d.table(s, M, 1.44, 9.40,
            ["Route", "Aircraft", "Programme", "Annual pax", "Annual revenue",
             "Annual cost", "Annual profit", "Margin", "Frames"],
            rows, col_w=[1.0, 0.8, 1.05, 0.95, 1.25, 1.25, 1.25, 0.65, 0.65],
            size=9.2, hdr_size=8.0, row_h=0.40, hdr_h=0.56, total_row=True,
            aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 6)
    d.callout(s, M, 3.94, 4.58, 0.86,
              ["One aircraft covers the whole", "summer programme, at 0.90 frames"], size=12)
    d.callout(s, 5.12, 3.94, 4.60, 0.86,
              ["Network margin %s on", "$%.1fm of annual revenue" % (t[1] / 1e6)],
              size=12, fill=NAVY, colour=WHITE)
    d.methodology(s, [
        ("Why this is the shape of the proposition",
         "Genoa to New York alone does not fill an aircraft. Genoa to New York, Newark and Boston does, and it does so within a single summer "
         "programme at 0.90 of a frame. That converts the route from a marginal standalone case into a base decision, which is a different and "
         "much easier conversation with an airline."),
        ("What it asks of the airport",
         "A base proposition needs overnight parking, crew facilities and a ground handling arrangement that will carry three United States "
         "rotations a week at the summer peak. Those are airport commitments, not airline ones, and they belong in the offer."),
    ], y=5.02, w=9.40)
    d.source(s, BASIS + " Network programme from the Genoa base summer schedule model.", size=7.5)

    s = d.content("Seasonality is the constraint, not the objection",
                  "Genoa's June to September concentration is identical to the Italian national figure")
    d._pic(s, "goa_ch_season.png", M, 1.40, w=5.90)
    rows = [["Nominal flat, planning load factor every month", "$%.2fm" % (CASE["season_a"] / 1e6)],
            ["Flat daily schedule, real monthly demand", "$%.2fm" % (CASE["season_b"] / 1e6)],
            ["Seasonal schedule, winter trimmed", "$%.2fm" % (CASE["season_c"] / 1e6)]]
    d.table(s, 6.32, 1.44, 3.40, ["Operating pattern", "Annual profit"], rows,
            col_w=[2.35, 1.05], size=9.0, hdr_size=8.6, row_h=0.62, hdr_h=0.44)
    d.panel(s, 6.32, 3.60, 3.40, 1.86, "The finding", [
        "Flying the same daily schedule through the winter costs $1.42m a year against a nominal flat case.",
        "Trimming January, February, November and December recovers $1.11m of it.",
    ], size=10)
    d.panel(s, 6.32, 5.58, 3.40, 1.44, "Genoa is normal for Italy", [
        "June to September is circa 41% of Genoa's annual traffic, the same as the 41% national figure.",
    ], size=10, fill=TEAL)
    d.bullets(s, M, 5.10, 5.90, [
        ("The honest reading.", "This is a seasonal proposition operated year round, or a year-round proposition operated seasonally. The airline will choose; the airport should price both."),
    ], size=10.5)
    d.source(s, "Sources: monthly demand index is an AviaSolutions working assumption pending a monthly Sabre pull, not measured Genoa to New York "
                "demand; Genoa's June to September concentration derived from operator monthly data with December verified against Assaeroporti; "
                "national figure from ENAC, January 2026.", size=7.0)

    s = d.content("What moves the answer, and what does not",
                  "Fuel moves it. Capture, across the tested range, does not.")
    d._pic(s, "goa_ch_fuel.png", M, 1.40, w=5.86)
    d.panel(s, 6.32, 1.44, 3.40, 2.50, "Fuel is the live variable", [
        "A move from $0.68 to $1.10 per kilogramme takes annual profit from $12.8m to $5.9m.",
        "The route stays profitable across the whole tested band, but the margin roughly halves.",
    ], size=10)
    d.panel(s, 6.32, 4.08, 3.40, 2.94, "Why capture does not move it", [
        "Across a capture range of 50% to 75% of leaked catchment demand, annual profit does not change, because the planning load-factor cap of 85% binds throughout.",
        "In plain terms: demand is not the binding constraint at this frequency and gauge. Seats are.",
        "That is a genuinely useful finding. It means the commercial risk sits in cost and fuel, not in whether the catchment shows up.",
    ], size=10, fill=TEAL)
    d.bullets(s, M, 5.16, 5.86, [
        ("The caveat that goes with it.", "The load-factor cap binding across the whole range is a property of the planning case, not a measurement. If real capture were far below 50%, the cap would stop binding and the result would move sharply."),
    ], size=10.5)
    d.source(s, BASIS + " Scenario grid from the business case; capture defined as the share of leaked catchment demand recovered by a "
                "Genoa service.", size=7.5)

    # ======================================================= SECTION 6
    d.divider("goa_div_counter.jpg", "6", "The counter-case",
              "The narrowbody thin-transatlantic record, including the routes that were cut")

    s = d.content("The narrowbody transatlantic record",
                  "Presented in full, because it is the first thing an airline will raise")
    rows = [["SAS", "Gothenburg - Newark", "A321LR", "Cut", "The nearest analogue to this proposition, and it did not hold"],
            ["SAS", "Aalborg - Newark", "A321LR", "Cut", "As above"],
            ["Aer Lingus", "Dublin - Minneapolis", "A321XLR", "Being cut", "Load factors bottoming at 30.0%, beaten head to head by Delta"],
            ["Neos", "Bari - New York JFK", "787", "Cut for 2026", "Backfilled by United, so the market held even though the operator did not"],
            ["American Airlines", "Philadelphia - Bologna", "Widebody", "Died after one 2019 season", "A secondary Italian city that did not hold"],
            ["La Compagnie", "Newark - Milan Malpensa", "A321neo, all business", "Operating since 15 April 2022", "The counter-example: four years on a narrowbody premium configuration"],
            ["United", "Newark - Palermo", "Widebody", "Operating, extended to December 2026", "Circa 85% load factor at a secondary Italian city"]]
    d.table(s, M, 1.42, 9.40,
            ["Carrier", "Route", "Aircraft", "Status", "What it tells us"], rows,
            col_w=[1.15, 1.85, 1.35, 1.45, 3.60], size=9.0, hdr_size=8.6,
            row_h=0.52, hdr_h=0.44,
            aligns=[PP_ALIGN.LEFT] * 5)
    d.methodology(s, [
        ("Our reading",
         "The failures share a shape: a thin secondary origin, a single United States point, a narrowbody, and no second market to move the "
         "aircraft to when the first one softens. That is precisely the shape the network case in section 5 avoids, and it is why we put the "
         "three-point Genoa base on the table rather than a single Genoa to New York rotation."),
    ], y=5.44, w=9.40)
    d.source(s, "Sources: carrier statements and aviation trade press, retrieved 5 August 2026, secondary throughout; Aer Lingus load factor "
                "reporting, 2026. AviaSolutions analysis.", size=7.5)

    s = d.content("What would have to be true",
                  "The four conditions this route depends on, and who can settle each one")
    conds = [("Capture", "That Genoa recovers a material share of Ligurian and western Lombard demand currently flying from Milan, Turin, Nice or Pisa.", "Avia, from a Sabre origin-destination pull. Circa a fortnight."),
             ("Aircraft performance", "That an A321XLR can depart Genoa at New York weight in summer temperatures on the preferential runway.", "Airbus, from a performance study. The airport should commission it."),
             ("Fuel", "That jet fuel stays inside the $0.68 to $1.10 per kilogramme band the case is run across.", "Nobody. It is a risk to be priced, not a question to be answered."),
             ("A funder", "That someone other than the airport company underwrites the launch risk.", "Regione Liguria, the chamber of commerce, the port authority or the tourism consortium.")]
    y = 1.46
    for i, (h, cond, who) in enumerate(conds):
        d._rect(s, M, y, 9.40, 1.22, fill=LIGHT if i % 2 == 0 else WHITE)
        d._rect(s, M, y, 0.055, 1.22, fill=ORANGE)
        d._text(s, M + 0.20, y + 0.12, 1.95, 0.42, [(h, 14, True, NAVY)])
        d._text(s, M + 0.20, y + 0.58, 1.95, 0.56,
                [("Settled by:", 9, True, GREY)])
        d._text(s, M + 0.20, y + 0.78, 1.95, 0.40, [(who.split(",")[0], 9, False, BODY)])
        d._text(s, M + 2.35, y + 0.14, 4.30, 0.98, [(cond, 11, False, BODY)])
        d._text(s, 6.85, y + 0.14, 2.80, 0.98, [(who, 10, False, TEAL)])
        y += 1.28
    d.callout(s, M, 6.60, 9.40, 0.56,
              "Two of the four can be closed before this route is taken to an airline, and both should be",
              size=12, fill=NAVY, colour=WHITE)
    d.source(s, "Source: AviaSolutions analysis. Fuel band per the scenario grid in section 5.", size=7.5)

    # ======================================================= SECTION 7
    d.divider("goa_div_method.jpg", "7", "Methodology and track record",
              "How the case is built, and how the engine behind it has been tested")

    s = d.content("Basis of the case", "What is modelled, what is assumed and what is not yet known")
    d.methodology(s, [
        ("Catchment and demand",
         "Catchment allocation from GeoNames population and OSRM road times, validated. The New York origin and destination market split and "
         "the fares come from Sabre ODPOO 2024, validated. Fares are one-way: economy $345 and business $1,400, which is an entrant-yield "
         "judgement rather than an observed Genoa fare, because no Genoa to New York fare exists to observe."),
        ("Economics",
         "The cost stack is anchored to the Avia route economics module. Maintenance comes from the validated Airbus 2024-2025 reserves. "
         "Ownership comes from appraiser lease rates, and is directional. Fuel at $0.90 per kilogramme is a through-cycle planning assumption. "
         "Crew at $1,200 per block hour is a low-cost-carrier judgement pending a citation."),
        ("What is explicitly not verified",
         "Airport charges at Genoa and at the New York end are indicative placeholders. The monthly seasonality profile is an assumption "
         "pending a monthly Sabre pull. Both are flagged on the slides where they appear, and both are in the assumptions register in the "
         "appendix."),
        ("The rule we work to",
         "Flag rather than fill. Where no clean source exists the cell is left open and listed, rather than approximated to make a slide "
         "look finished."),
    ], y=1.46, w=9.40)
    d.source(s, "Source: AviaSolutions assumptions register, Genoa - New York business case, version of 5 August 2026.", size=7.5)

    s = d.content("Tested against 2,915 real route launches",
                  "The engine behind Avia's route forecasts, graded against what actually flew")
    d._pic(s, "QSI_accuracy_distribution_fitted.png", M, 1.44, w=5.70)
    d.panel(s, 6.20, 1.44, 3.52, 2.70, "How the test was built", [
        "Every genuinely new route launched worldwide in 2016-2019 and 2025, 2,915 of them, taken from the complete OAG schedule archive.",
        "For each one the engine saw only the world as it stood the month before launch.",
        "Its first-year forecast was then compared with the passengers who actually flew.",
    ], size=10.5)
    d.callout(s, 6.20, 4.28, 3.52, 1.02,
              ["89% within 20% of outturn", "82% within 10%"], size=15)
    d.panel(s, 6.20, 5.42, 3.52, 1.60, "And on routes never seen", [
        "Forecasting portfolios of twenty unseen routes, the portfolio total came within 20% of the actual total 94% of the time.",
    ], size=10.5)
    d.source(s, "Source: Avia Cortex QSI backtest programme, runs of 5 August 2026, n=2,915 launches, 2016-2019 and 2025; the pandemic years "
                "2020-2023 are deliberately excluded. Outturn is US DOT DB1B for United States domestic routes and Sabre MIDT elsewhere. "
                "AviaSolutions analysis.", size=7.5)

    # ======================================================= SECTION 8
    d.divider("goa_div_offer.jpg", "8", "The airport, the aircraft and the ask",
              "What Genoa can do today, what is still open, and what we are asking for")

    s = d.content("Genoa airport: what is settled and what is open",
                  "The runway question is the one item that has to be closed before an airline conversation")
    rows = [["Long-haul or transatlantic service today", "None, from any airline", "Settled"],
            ["Summer 2026 seat capacity, all routes", "1,375,000 seats, up 10% on summer 2025", "Settled"],
            ["Rank among Italian commercial airports", "19th of 44", "Settled"],
            ["Widebody landing capability", "An ITA A330-900neo diverted in on 6 April 2026", "A landing, not a transatlantic departure"],
            ["Runway length", "Sources give 2,916 m, 2,925 m and 3,066 m", "OPEN. Not printed here until the ENAV AIP chart is in hand"],
            ["A321XLR performance out of Genoa at New York weight", "No Airbus study exists", "OPEN. The airport should commission it"],
            ["Ownership", "Port authority 60%, Camera di Commercio di Genova 40%", "Settled"],
            ["Airport company financial capacity", "EUR 1.5m EBITDA and EUR 105,450 net profit in 2025", "Settled, and it matters: see the ask"]]
    d.table(s, M, 1.42, 9.40, ["Item", "Position", "Status"], rows,
            col_w=[3.1, 3.6, 2.7], size=9.2, hdr_size=8.8, row_h=0.50, hdr_h=0.44,
            aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
    d.panel(s, M, 5.60, 9.40, 1.42, "Why the runway figure is not on this slide", [
        "The three published figures differ by 150 metres and the largest comes from social media. The whole aircraft case rests on it. "
        "Avia does not print a number it cannot source, and a range comparison against the A321XLR's published 4,700 nm is safe where a "
        "feasibility assertion is not.",
    ], size=10.5, fill=RED)
    d.source(s, "Sources: ENAC, Air Traffic Data 2025, January 2026; GuidaViaggi reporting the airport operator, 19 March 2026; Genova24, "
                "6 April 2026, for the diversion; Shipmag.it and Primocanale, 8 June 2026, for the shareholding; airport company results as "
                "reported at the May 2026 shareholders' meeting. Runway figures deliberately not quoted.", size=7.0)

    s = d.content("The New York end", "Neither airport is an open door, and the deck should not pretend otherwise")
    rows = [["New York JFK", "Level 3 slot controlled", "The strongest catchment and the Italian-American base; hardest to enter"],
            ["Newark", "Under an FAA capacity order through 24 October 2026", "United's hub, and the airport that has taken every recent secondary Italian route"],
            ["LaGuardia", "Cannot take the aircraft or the sector", "Excluded"]]
    d.table(s, M, 1.44, 9.40, ["Airport", "Constraint", "Assessment"], rows,
            col_w=[1.6, 3.0, 4.8], size=9.8, hdr_size=9.2, row_h=0.62, hdr_h=0.44,
            aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
    d.methodology(s, [
        ("Our view on the operator",
         "The carriers that have actually launched secondary Italian routes since 2022 are United, Delta, American, Neos and La Compagnie. "
         "United is the closest fit on evidence: it holds Newark, it has taken Palermo, Bari and Olbia, and it is the carrier that backfilled a "
         "route another operator abandoned. La Compagnie is the interesting second, because its all-business narrowbody model on Newark to "
         "Malpensa is the closest commercial template to a premium-configured Genoa proposition."),
        ("What we would not claim",
         "That any carrier has expressed interest. None has. There is no announced plan, no memorandum of understanding and no feasibility "
         "study for a Genoa transatlantic route, and the airport's public statements on the subject are aspirations rather than plans."),
    ], y=3.70, w=9.40)
    d.source(s, "Sources: FAA Newark capacity order, effective through 24 October 2026; IATA Level 3 designation for New York JFK; carrier "
                "route announcements 2022-2026, secondary. AviaSolutions analysis.", size=7.5)

    s = d.content("The ask", "Three things, and a question about who pays for them")
    asks = [("Commission the Airbus performance study", "An A321XLR study for departure at New York weight from Genoa's preferential runway in summer conditions. Without it the aircraft case is a range comparison, not a feasibility statement."),
            ("Commission the catchment and leakage measurement", "A Sabre origin-destination pull keyed to Ligurian and western Lombard postcodes. It is the number every airline will ask for first and it is the one number nobody has."),
            ("Assemble the funding package before the airline conversation", "The airport company made EUR 105,450 of net profit in 2025. It cannot underwrite a transatlantic launch. Regione Liguria, the Camera di Commercio, the port authority and the cruise and tourism sector are the candidates, and the package has to exist before a carrier is approached, not after.")]
    y = 1.48
    for i, (h, b) in enumerate(asks):
        d._rect(s, M, y, 9.40, 1.44, fill=LIGHT if i % 2 == 0 else WHITE)
        d._rect(s, M, y, 0.055, 1.44, fill=ORANGE)
        d._text(s, M + 0.20, y + 0.14, 3.30, 0.90, [(h, 13, True, NAVY)])
        d._text(s, M + 3.70, y + 0.16, 5.50, 1.12, [(b, 11, False, BODY)])
        y += 1.52
    d.callout(s, M, 6.20, 9.40, 0.72,
              "Proposed next step: a working session with Regione Liguria and the port authority, autumn 2026, to fund the two studies",
              size=13, fill=NAVY, colour=WHITE)
    d.source(s, "Source: AviaSolutions analysis. Airport company results as reported at the May 2026 shareholders' meeting.", size=7.5)

    s = d.content("Choose Genoa", "The case in six lines")
    picks = [("An existing corridor", "Circa 37% of Genoa and Savona port traffic is United States linked, and MSC already sails Genoa to New York direct in nine days."),
             ("Growing fastest", "Up 18.1% to 1,587,761 passengers in 2025, against Malpensa's 8.6% and Fiumicino's 4.5%."),
             ("A second gateway already at scale", "1,630,593 cruise passengers in 2025, more than the airport handled, worth EUR 346m to the local economy."),
             ("Within the aircraft", "3,509.6 nautical miles, circa 75% of the A321XLR's published range and shorter than Madrid to Boston."),
             ("A base, not a route", "Three United States points from one Genoa base fill 0.90 of a frame at a %s network margin." % pct(CASE["network_total"][4], 1)),
             ("Nobody else has moved", "The largest north-western Italian catchment with no United States service and no announced plan.")]
    y = 1.48
    for i, (h, b) in enumerate(picks):
        d._rect(s, M, y, 9.40, 0.86, fill=LIGHT if i % 2 == 0 else WHITE)
        d._rect(s, M, y, 0.055, 0.86, fill=ORANGE)
        d._text(s, M + 0.20, y + 0.12, 2.60, 0.62, [(h, 13, True, NAVY)],
                anchor=MSO_ANCHOR.MIDDLE)
        d._text(s, M + 2.90, y + 0.11, 6.40, 0.64, [(b, 11, False, BODY)],
                anchor=MSO_ANCHOR.MIDDLE)
        y += 0.90
    d.source(s, "Sources as cited on the preceding slides. Full citations in the Project Liguria evidence pack, 5 August 2026.", size=7.5)

    # thank you
    s = d._slide()
    d._hero(s, "goa_thanks.jpg", 0, 0, SW, SH, "closing photograph")
    d._rect(s, 0, 2.70, SW, 2.10, fill=NAVY)
    d._text(s, 0, 3.00, SW, 0.60, [("Thank you", 34, True, WHITE)],
            align=PP_ALIGN.CENTER)
    d._text(s, 0, 3.72, SW, 0.36,
            [("A direct link between Genoa and New York", 15, False,
              RGBColor(0xD5, 0xE2, 0xF2))], align=PP_ALIGN.CENTER)
    d._text(s, 0, 4.14, SW, 0.30,
            [("Avia Solutions Limited  |  %s  |  Commercial in Confidence" % CODENAME,
              11, True, RGBColor(0x7F, 0xC6, 0xF0))], align=PP_ALIGN.CENTER)

    # appendix
    d.divider("goa_div_appendix.jpg", "9", "Appendix",
              "The assumptions register and the data gaps, stated openly")

    s = d.content("Assumptions register", "Every input, its source, and its status")
    rows = [["Catchment population and road times", "GeoNames population, OSRM road times", "Validated"],
            ["New York market split and fares", "Sabre ODPOO 2024", "Validated"],
            ["Fares: economy $345, business $1,400 one-way", "Entrant-yield judgement", "Judgement, no Genoa fare exists to observe"],
            ["Cost stack", "Avia route economics module", "Directional"],
            ["Maintenance", "Airbus 2024-2025 maintenance reserves", "Validated"],
            ["Ownership", "Appraiser lease rates", "Directional"],
            ["Jet fuel $0.90 per kilogramme", "Through-cycle planning assumption", "Assumption, tested across $0.68 to $1.10"],
            ["Crew $1,200 per block hour", "Low-cost-carrier judgement", "Pending a citation"],
            ["Genoa and New York airport charges", "Indicative placeholders", "NOT VERIFIED"],
            ["Monthly seasonality profile", "Leisure-weighted assumption", "Pending a monthly Sabre pull"],
            ["2026 transatlantic booking environment", "Demand watchpoint", "Monitored, not modelled"]]
    d.table(s, M, 1.42, 9.40, ["Input", "Source", "Status"], rows,
            col_w=[3.5, 3.4, 2.5], size=9.2, hdr_size=8.8, row_h=0.42, hdr_h=0.42,
            aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
    d.panel(s, M, 6.20, 9.40, 0.80, None, [
        "Two inputs are marked not verified and one is pending a citation. None of the three is load-bearing on the direction of the result, "
        "but all three move the margin, and they are named here rather than buried.",
    ], size=10, fill=RED)
    d.source(s, "Source: AviaSolutions assumptions register, Genoa - New York business case, 5 August 2026.", size=7.5)

    s = d.content("Data gaps to close before issue", "Stated openly rather than filled with estimates")
    gaps = [["Genoa runway length and declared distances", "The A321XLR feasibility case rests on it; sources differ by 150 m", "ENAV AIP chart for LIMJ"],
            ["A321XLR performance out of Genoa", "No study exists at New York weight in summer conditions", "Commission from Airbus"],
            ["Ligurian leakage to Milan, Turin, Nice and Pisa", "No study exists anywhere; it is the heart of the pitch", "Sabre origin-destination pull by postcode"],
            ["Genoa annual passengers 2015 to 2023", "Could not be tied to a named publisher", "ENAC Dati di traffico PDFs and Assaeroporti monthly files; circa one analyst-hour"],
            ["Genoa summer 2026 route map", "Sourced from regional press, not schedule data", "OAG extract, which Avia holds in-house"],
            ["Italian ancestry in New York and New Jersey", "Taken from a secondary aggregator of the Census table", "Verify on data.census.gov"],
            ["Genoa and New York airport charges", "Indicative placeholders in the economics", "Published tariff tables at both ends"]]
    d.table(s, M, 1.44, 9.40, ["Gap", "Why it matters", "How to close it"], gaps,
            col_w=[3.0, 3.7, 2.7], size=9.2, hdr_size=8.8, row_h=0.60, hdr_h=0.44,
            aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
    d.panel(s, M, 5.94, 9.40, 1.08, None, [
        "Four claims in the evidence pack are single-sourced or contradicted and are excluded from this deck entirely: the reported Emirates "
        "Malpensa to JFK suspension, a probable-false American Philadelphia to Bologna relaunch, United's first A321XLR delivery date, and "
        "American's Philadelphia to Naples 2026 frequency.",
    ], size=9.5, fill=RED)
    d.source(s, "Source: AviaSolutions evidence register for %s, 5 August 2026." % CODENAME, size=7.5)

    d.save(OUT, title="%s - Genoa to New York route business case" % CODENAME,
           subject="Route business case, prepared for Aeroporto di Genova and Regione Liguria")
    print("slides:", len(d.prs.slides._sldIdLst))
    print("written:", OUT)
    return [p for _, _, p in d.contents]


if __name__ == "__main__":
    found = main()
    main(pages=[3] + found)
