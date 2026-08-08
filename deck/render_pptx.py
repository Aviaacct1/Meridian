#!/usr/bin/env python3
"""Observatory PowerPoint renderer.

Reads the same deck_spec.py the HTML renderer reads and emits a .pptx. Nothing
in the spec knows this file exists: the spec carries content, this carries the
Observatory house style, and a third renderer will carry a client's template.
That separation is the whole point of the three-style product.

House style, from Brand Guidelines v1.3:
  - 16:9, the PowerPoint equivalent of the 1920x1080 HTML canvas
  - ink 0F1B28, paper E7E4DD, gold B8862F on paper and D4A249 on ink
  - hairline rules, never filled header bands
  - type never sits over a photograph. Covers and dividers put the type on its
    own ink panel and the photograph in a panel beside it.
  - one evidence plate to a page, four to a report

Fonts. Newsreader and IBM Plex Mono are the brand faces and the HTML renderer
can fetch them. PowerPoint cannot: it substitutes whatever the reader has, so a
client opening the .pptx without Newsreader installed sees something else.
Run with --safe-fonts for Cambria and Arial, which ship with Office and hold
the serif and sans contrast. Deliver as PDF and the brand faces are embedded,
so the default stays on brand.

Avia Solutions Limited. All rights reserved.
"""

import argparse
import math
import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

import deck_spec as S

# --- canvas ----------------------------------------------------------------
# The HTML renderer works in a 1920x1080 pixel canvas. Keeping the same numbers
# here means a measurement can be read off one renderer and used in the other.
# Who publishes the file. Product decks go out under The Aviation Observatory, the
# separate company; an Avia client deliverable built through the same renderer sets
# meta["author"] to Avia Solutions instead. Anything else fails verify().
AUTHOR = os.environ.get("AVIA_DECK_AUTHOR", "The Aviation Observatory")
PUBLISHERS = {"The Aviation Observatory", "Avia Solutions"}

PX = 1920.0
SLIDE_W_IN = 13.3333
SLIDE_H_IN = 7.5
EMU_PER_PX = int(round(Emu(int(SLIDE_W_IN * 914400)) / PX))


def X(px):
    return Emu(int(round(px * EMU_PER_PX)))


# --- text metrics ----------------------------------------------------------
# Layout here is computed, not placed by hand, so every box has to know how
# tall its text will actually be. The canvas is 1920 px across a 13.333 inch
# slide, so one point is exactly two canvas pixels. Average character width for
# mixed-case text runs near half the point size. Tracking widens every
# character, which matters: the house eyebrow tracks at 1.2pt, nearly 30 per
# cent extra at 8.5pt, and ignoring it ran the source lines into the footer.
PX_PER_PT = 2.0
CHAR_W = 0.50
SLACK = 0.35        # of a line, so a renderer's leading cannot clip a descender


def cpl(width_px, size, spacing=0.0):
    """Characters that fit on one line."""
    return max(6.0, width_px / ((CHAR_W * size + spacing) * PX_PER_PT))


def nlines(s, width_px, size, spacing=0.0):
    """Wrapped line count, respecting any hard breaks."""
    total = 0
    for part in str(s or "").split("\n"):
        total += max(1, int(math.ceil(len(part) / cpl(width_px, size, spacing))))
    return total


def cpl_width(chars, size, spacing=0.0):
    """The box width that holds this many characters on a line."""
    return chars * (CHAR_W * size + spacing) * PX_PER_PT


def block_h(s, width_px, size, line=1.4, spacing=0.0):
    """Height in canvas pixels a run of text will occupy, with slack."""
    return (nlines(s, width_px, size, spacing) + SLACK) * line * size * PX_PER_PT


# --- palette ---------------------------------------------------------------
INK_BG = "0F1B28"
INK_FG = "F4F1EA"
INK_MUTE = "9AA7B3"
INK_LABEL = "7D8B98"
INK_RULE = "2A3A49"
GOLD = "D4A249"

PAPER_BG = "E7E4DD"
PAPER_FG = "141C25"
PAPER_MUTE = "5A6470"
PAPER_LABEL = "6B7480"
PAPER_RULE = "D8D2C4"
GOLD_D = "B8862F"
CAUTION = "8A3A2A"
POSITIVE = "3F6B4A"

SERIF, MONO = "Newsreader", "IBM Plex Mono"
SAFE_SERIF, SAFE_MONO = "Cambria", "Arial"

# margins, in canvas pixels
ML, MR = 132, 132
MT, MB = 108, 96
COL = PX - ML - MR


def rgb(hexstr):
    return RGBColor.from_string(hexstr)


class Theme(object):
    """Two surfaces, ink and paper. Everything else derives from which is up."""

    def __init__(self, dark=False, safe_fonts=False):
        self.dark = dark
        self.bg = INK_BG if dark else PAPER_BG
        self.fg = INK_FG if dark else PAPER_FG
        self.mute = INK_MUTE if dark else PAPER_MUTE
        self.label = INK_LABEL if dark else PAPER_LABEL
        self.rule = INK_RULE if dark else PAPER_RULE
        self.gold = GOLD if dark else GOLD_D
        self.serif = SAFE_SERIF if safe_fonts else SERIF
        self.mono = SAFE_MONO if safe_fonts else MONO

    def tone(self, name):
        return {"accent": self.gold, "caution": CAUTION,
                "positive": POSITIVE}.get(name, self.fg)


# ---------------------------------------------------------------------------
# Drawing primitives
# ---------------------------------------------------------------------------
def fill(slide, colour, x=0, y=0, w=PX, h=1080):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, X(x), X(y), X(w), X(h))
    box.fill.solid()
    box.fill.fore_color.rgb = rgb(colour)
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def rule(slide, x, y, w, colour, weight=1.0):
    """A hairline. The house style uses these instead of filled header bands."""
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, X(x), X(y), X(w),
                                X(max(weight, 0.75)))
    ln.fill.solid()
    ln.fill.fore_color.rgb = rgb(colour)
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def text(slide, s, x, y, w, h, *, size=16, colour="000000", font=SERIF,
         bold=False, italic=False, align=PP_ALIGN.LEFT, spacing=None,
         caps=False, anchor=MSO_ANCHOR.TOP, line=1.25, space_after=0):
    """One text box, one paragraph per line of s (str or list of str)."""
    box = slide.shapes.add_textbox(X(x), X(y), X(w), X(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = s if isinstance(s, (list, tuple)) else [s]
    for i, raw in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        if space_after:
            p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = (raw or "").upper() if caps else (raw or "")
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = rgb(colour)
        if spacing:
            _char_spacing(run, spacing)
    return box


def _char_spacing(run, points):
    """python-pptx exposes no letter spacing, so set a:rPr/@spc directly."""
    run.font._rPr.set("spc", str(int(round(points * 100))))


def mono_label(slide, s, x, y, w, theme, *, size=9, colour=None, gold=False):
    """The house eyebrow: mono, caps, wide tracking, muted."""
    return text(slide, s, x, y, w, 20, size=size, font=theme.mono, caps=True,
                colour=colour or (theme.gold if gold else theme.label),
                spacing=1.6)


def picture_fit(slide, path, x, y, w, h, align="center"):
    """Place the whole image inside the box, to scale, and centre it.

    picture_cover is right for photography, where the frame matters more than
    the edges of the shot. It is wrong for a chart or a map: cropping a map
    removes geography and cropping a chart removes an axis, and handing
    python-pptx both dimensions distorts whatever is left. A generated figure
    keeps its aspect ratio and gives back the space it does not need.

    Returns (picture, drawn_height) so the caller can set what follows against
    what was actually placed rather than against the box it offered.
    """
    from PIL import Image
    try:
        iw, ih = Image.open(path).size
        scale = min(float(w) / iw, float(h) / ih)
        dw, dh = iw * scale, ih * scale
    except Exception:
        dw, dh = w, h
    dx = (w - dw) / 2.0 if align == "center" else 0.0
    return slide.shapes.add_picture(path, X(x + dx), X(y), X(dw), X(dh)), dh


def picture_cover(slide, path, x, y, w, h):
    """Fill the box with the image, cropping the overflow rather than squashing.

    python-pptx scales to the box it is given, so an image with a different
    aspect ratio would distort. Cropping keeps the geometry honest.
    """
    from PIL import Image
    pic = slide.shapes.add_picture(path, X(x), X(y), X(w), X(h))
    try:
        iw, ih = Image.open(path).size
    except Exception:
        return pic
    want, have = float(w) / float(h), float(iw) / float(ih)
    if abs(want - have) < 0.001:
        return pic
    if have > want:                      # image is wider: trim left and right
        trim = (1.0 - want / have) / 2.0
        pic.crop_left = pic.crop_right = trim
    else:                                # image is taller: trim top and bottom
        trim = (1.0 - have / want) / 2.0
        pic.crop_top = pic.crop_bottom = trim
    return pic


def _fit(title, base, at, smaller):
    """Long titles step down rather than overflow their band."""
    return smaller if len(title or "") > at else base


def fit_to_width(lines, width_px, cap, floor=20):
    """Largest size at which every line still fits on one line."""
    longest = max((len(l or "") for l in lines), default=1)
    return max(floor, min(cap, 0.92 * width_px / max(longest, 1)))


# ---------------------------------------------------------------------------
# Chrome
# ---------------------------------------------------------------------------
def chrome(slide, theme, meta, page=None, section=None):
    """Running head and foot. Hairlines, no bands."""
    rule(slide, ML, MT - 34, COL, theme.rule)
    mono_label(slide, meta["codename"], ML, MT - 58, 500, theme, gold=True)
    if section:
        text(slide, section, ML + 500, MT - 58, COL - 500, 20, size=9,
             font=theme.mono, caps=True, colour=theme.label, spacing=1.6,
             align=PP_ALIGN.RIGHT)
    rule(slide, ML, 1080 - MB + 26, COL, theme.rule)
    foot = "%s  -  %s" % (meta["confidentiality"], meta["date"])
    text(slide, foot, ML, 1080 - MB + 42, COL - 120, 20, size=8.5,
         font=theme.mono, caps=True, colour=theme.label, spacing=1.4)
    if page is not None:
        text(slide, "P: %d" % page, ML + COL - 120, 1080 - MB + 42, 120, 20,
             size=8.5, font=theme.mono, caps=True, colour=theme.label,
             spacing=1.4, align=PP_ALIGN.RIGHT)


PROSE_PT, PROSE_PT_MAX, PROSE_FILL = 12.5, 22.0, 0.86


def prose_size(paras, width_px, room_px, base=PROSE_PT, cap=PROSE_PT_MAX,
               fill=PROSE_FILL):
    """The point size that lets a prose page fill its own height.

    A section paragraph is written to a 430-character budget, which is about 65
    words, and that is deliberate: it is one argument, and lengthening it to
    fill the slide would be padding. Set at 12.5pt beside a callout column, one
    such paragraph occupies under a third of the page, and the reader reads the
    emptiness as an unfinished document rather than as a deliberate opening
    statement.

    So the page is filled by typography, not by words. The size steps up until
    the block occupies `fill` of the available height, capped so the measure
    stays readable. A page that is already full, the methodology page with its
    four paragraphs, finds no headroom and stays at the base size, so this
    changes nothing where nothing needed changing.
    """
    def total(size):
        h = 0.0
        for para in paras:
            body = para[1] if isinstance(para, (tuple, list)) else para
            if isinstance(para, (tuple, list)) and para[0]:
                h += 26
            h += block_h(body, min(width_px, cpl_width(78, size)), size, 1.5) + 16
        return h

    if room_px <= 0 or not paras:
        return base
    best = base
    size = base
    while size + 0.5 <= cap:
        size += 0.5
        if total(size) > room_px * fill:
            break
        best = size
    return best


def heading(slide, theme, sl, y=MT):
    """Slide title and optional subtitle. Returns the y to carry on from."""
    title = sl.get("title")
    if not title:
        return y
    size = _fit(title, 32, 52, 26)
    th = block_h(title, COL, size, 1.12)
    text(slide, title, ML, y, COL, th, size=size, font=theme.serif,
         colour=theme.fg, line=1.12)
    y += th + 12
    if sl.get("subtitle"):
        sh = block_h(sl["subtitle"], COL * 0.78, 13.5, 1.35)
        text(slide, sl["subtitle"], ML, y, COL * 0.78, sh, size=13.5,
             font=theme.serif, colour=theme.mute, italic=True, line=1.35)
        y += sh + 8
    return y + 22


def source_line(slide, theme, sl):
    """Sources sit clear of the footer rule, however many lines they run to."""
    if not sl.get("source"):
        return 1080 - MB + 6
    h = block_h(sl["source"], COL, 8.5, 1.45, spacing=1.2)
    y = 1080 - MB + 10 - h
    text(slide, sl["source"], ML, y, COL, h, size=8.5, font=theme.mono,
         colour=theme.label, spacing=1.2, line=1.45)
    return y


# ---------------------------------------------------------------------------
# Components
# ---------------------------------------------------------------------------
def draw_panel(slide, theme, p, x, y, w):
    """A titled list. Tone colours the title and the leading hairline only."""
    tone = theme.tone(p.get("tone", "neutral"))
    rule(slide, x, y, 46, tone, 1.6)
    mono_label(slide, p["title"], x, y + 14, w, theme, colour=tone)
    yy = y + 40
    for item in p["items"]:
        h = block_h(item, w, 11.5, 1.42)
        text(slide, item, x, yy, w, h, size=11.5, font=theme.serif,
             colour=theme.fg, line=1.42)
        yy += h + 10
    return yy


def draw_callout(slide, theme, c, x, y, w):
    tone = theme.tone(c.get("tone", "accent"))
    sizes = [12.5 if i == 0 else 11.5 for i in range(len(c["lines"]))]
    heights = [block_h(ln, w - 26, sz, 1.35)
               for ln, sz in zip(c["lines"], sizes)]
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, X(x), X(y), X(2.5),
                                 X(sum(heights) + 8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(tone)
    bar.line.fill.background()
    bar.shadow.inherit = False
    yy = y + 2
    for i, ln in enumerate(c["lines"]):
        text(slide, ln, x + 18, yy, w - 26, heights[i], size=sizes[i],
             font=theme.serif, colour=theme.fg if i == 0 else theme.mute,
             bold=(i == 0), line=1.35)
        yy += heights[i]
    return yy + 12


def draw_table(slide, theme, spec, x, y, w, max_h=None):
    """Hairline table. No filled header, no banding: the rules do the work."""
    head, rows = spec.get("head") or [], spec.get("rows") or []
    align = spec.get("align") or []
    ncol = len(head) or (len(rows[0]) if rows else 0)
    if not ncol:
        return y
    # the spec carries relative column weights, not fractions: normalise.
    # Reading a weight of 1.0 as the whole width collapsed a nine-column
    # table to two visible columns.
    widths = list(spec.get("widths") or [1.0] * ncol)[:ncol]
    widths += [1.0] * (ncol - len(widths))
    tot = sum(widths) or 1.0
    widths = [x / tot for x in widths]
    xs, acc = [], x
    for frac in widths:
        xs.append(acc)
        acc += w * frac
    rowh = spec.get("row_height", 30)
    if max_h and rowh * (len(rows) + 1) > max_h:
        rowh = max(20, max_h / float(len(rows) + 1))
    size = 12 if rowh >= 28 else 11
    # a nine-column table cannot carry twelve-point type
    narrowest = min(widths) * w
    longest = max([len(str(c)) for r in rows for c in r[:ncol]] or [1])
    size = max(7.5, min(size, 0.92 * narrowest / max(longest, 1)))

    if head:
        for i, cell in enumerate(head):
            a = PP_ALIGN.RIGHT if (align[i:i + 1] or ["l"])[0] == "r" \
                else PP_ALIGN.LEFT
            text(slide, cell, xs[i], y, w * widths[i] - 12, 22, size=8.5,
                 font=theme.mono, caps=True, colour=theme.label, spacing=1.4,
                 align=a)
        y += 26
        rule(slide, x, y, w, theme.gold, 1.2)
        y += 8

    for r, row in enumerate(rows):
        for i, cell in enumerate(row[:ncol]):
            a = PP_ALIGN.RIGHT if (align[i:i + 1] or ["l"])[0] == "r" \
                else PP_ALIGN.LEFT
            numeric = (align[i:i + 1] or ["l"])[0] == "r"
            text(slide, str(cell), xs[i], y + 4, w * widths[i] - 12, rowh,
                 size=size, font=theme.mono if numeric else theme.serif,
                 colour=theme.fg, align=a,
                 bold=bool(spec.get("bold_rows")) and r in spec["bold_rows"])
        y += rowh
        if r < len(rows) - 1:
            rule(slide, x, y - 6, w, theme.rule, 0.75)
    return y


def draw_stats(slide, theme, stats, x, y, w):
    """Big numbers in a row, each over its own hairline."""
    n = max(len(stats), 1)
    gap = 28
    cw = (w - gap * (n - 1)) / float(n)
    for i, st in enumerate(stats):
        lab, val, accent = (list(st) + [False])[:3]
        cx = x + i * (cw + gap)
        rule(slide, cx, y, cw, theme.gold if accent else theme.rule,
             1.4 if accent else 1.0)
        mono_label(slide, lab, cx, y + 14, cw, theme,
                   colour=theme.gold if accent else theme.label)
        val = str(val)
        text(slide, val, cx, y + 54, cw, 74, size=_fit(val, 40, 9, 30),
             font=theme.serif, colour=theme.gold if accent else theme.fg,
             line=1.05)
    return y + 132


# ---------------------------------------------------------------------------
# Slide kinds
# ---------------------------------------------------------------------------
def slide_cover(deck, spec, sl, meta, safe):
    theme = Theme(dark=True, safe_fonts=safe)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    fill(slide, theme.bg)
    # Type on its own panel, photograph in its own panel. Chapter 14: type is
    # never set over a photograph.
    panel_w = PX * 0.52
    if sl.get("image") and os.path.exists(sl["image"]):
        picture_cover(slide, sl["image"], panel_w, 0, PX - panel_w, 1080)
    else:
        fill(slide, INK_RULE, panel_w, 0, PX - panel_w, 1080)
        text(slide, "No cleared image for this slot", panel_w, 520,
             PX - panel_w, 40, size=10, font=theme.mono, caps=True,
             colour=INK_LABEL, spacing=1.6, align=PP_ALIGN.CENTER)
    fill(slide, theme.bg, 0, 0, panel_w, 1080)

    x, w = ML, panel_w - ML - 60
    mono_label(slide, meta["codename"], x, 210, w, theme, size=10, gold=True)
    rule(slide, x, 248, 90, theme.gold, 1.6)
    size = fit_to_width(sl["title_lines"], w, 42)
    lh = size * 1.14 * PX_PER_PT
    y = 320
    for ln in sl["title_lines"]:
        text(slide, ln, x, y, w, lh + 6, size=size, font=theme.serif,
             colour=theme.fg, line=1.14)
        y += lh
    if meta.get("strap"):
        sh = block_h(meta["strap"], w, 14, 1.4)
        text(slide, meta["strap"], x, y + 28, w, sh, size=14,
             font=theme.serif, colour=INK_MUTE, italic=True, line=1.4)
    rule(slide, x, 830, w, INK_RULE)
    text(slide, "Prepared for %s" % meta["prepared_for"], x, 852, w, 30,
         size=12.5, font=theme.serif, colour=theme.fg)
    text(slide, [meta["date"], meta["confidentiality"]], x, 884, w, 60,
         size=9, font=theme.mono, caps=True, colour=INK_LABEL, spacing=1.6,
         line=1.6)
    return slide


def slide_divider(deck, spec, sl, meta, page, safe):
    theme = Theme(dark=True, safe_fonts=safe)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    fill(slide, theme.bg)
    panel_w = PX * 0.57
    if sl.get("image") and os.path.exists(sl["image"]):
        picture_cover(slide, sl["image"], panel_w, 0, PX - panel_w, 1080)
    else:
        fill(slide, INK_RULE, panel_w, 0, PX - panel_w, 1080)
    fill(slide, theme.bg, 0, 0, panel_w, 1080)

    x, w = ML, panel_w - ML - 60
    num = sl["number"]
    text(slide, "%02d" % num if isinstance(num, int) else str(num),
         x, 380, 200, 130, size=76,
         font=theme.serif, colour=theme.gold, line=1.0)
    rule(slide, x, 512, 90, theme.gold, 1.6)
    tsize = fit_to_width([sl.get("title") or ""], w, 34, 22)
    th = block_h(sl.get("title") or "", w, tsize, 1.14)
    text(slide, sl.get("title") or "", x, 546, w, th, size=tsize,
         font=theme.serif, colour=theme.fg, line=1.14)
    if sl.get("strap"):
        sh = block_h(sl["strap"], w, 13.5, 1.45)
        text(slide, sl["strap"], x, 556 + th, w, sh, size=13.5,
             font=theme.serif, colour=INK_MUTE, italic=True, line=1.45)
    text(slide, "P: %d" % page, x, 1080 - MB + 42, 200, 20, size=8.5,
         font=theme.mono, caps=True, colour=INK_LABEL, spacing=1.4)
    return slide


def slide_contents(deck, spec, sl, meta, page, safe):
    theme = Theme(safe_fonts=safe)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    fill(slide, theme.bg)
    chrome(slide, theme, meta, page)
    y = heading(slide, theme, dict(sl, title=sl.get("title") or "Contents"))
    items = sl["items"]
    half = (len(items) + 1) // 2
    colw = (COL - 80) / 2.0
    for c, chunk in enumerate((items[:half], items[half:])):
        yy = y
        for num, title, pg in chunk:
            cx = ML + c * (colw + 80)
            label = "%02d" % num if isinstance(num, int) else str(num)
            text(slide, label, cx, yy + 2, 60, 30, size=11,
                 font=theme.mono, colour=theme.gold, spacing=1.2)
            text(slide, title, cx + 62, yy, colw - 130, 34, size=15,
                 font=theme.serif, colour=theme.fg, line=1.3)
            if pg:
                text(slide, str(pg), cx + colw - 60, yy + 2, 56, 30, size=11,
                     font=theme.mono, colour=theme.label,
                     align=PP_ALIGN.RIGHT)
            yy += 40
            rule(slide, cx, yy - 8, colw - 4, theme.rule, 0.75)
    return slide


def slide_body(deck, spec, sl, meta, page, safe):
    """Everything with a heading and a content area under it."""
    theme = Theme(safe_fonts=safe)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    fill(slide, theme.bg)
    chrome(slide, theme, meta, page, sl.get("section"))
    y = heading(slide, theme, sl)
    src_h = block_h(sl["source"], COL, 8.5, 1.45, 1.2) if sl.get("source") else 0
    bottom = 1080 - MB - 6 - src_h
    kind = sl["type"]

    panels = sl.get("panels") or []
    callouts = sl.get("callouts") or []
    # a stat row and a wide table need the whole page; their callouts go
    # underneath in a row rather than into a side column, which was printing
    # them straight over the third statistic
    full_width = kind == "stat_row" or (
        kind == "table" and len((sl.get("table") or {}).get("head") or []) > 5)
    has_side = bool(panels or callouts) and not full_width
    main_w = COL * (0.62 if has_side else 1.0)
    side_x = ML + COL * 0.66
    side_w = COL * 0.34
    side_top = y

    if kind == "stat_row":
        y = draw_stats(slide, theme, sl["stats"], ML, y, COL)
        below = callouts if full_width else []
        cal_h = 0
        if below:
            cal_h = max(block_h(c["lines"][0],
                                (COL - 30 * (len(below) - 1)) / len(below) - 26,
                                12, 1.35) for c in below) + 24
        if sl.get("table"):
            y = draw_table(slide, theme, sl["table"], ML, y + 12, COL,
                           bottom - y - cal_h - 24)
        fig = sl.get("figure")
        if fig and os.path.exists(fig):
            fh = max(bottom - y - cal_h - 28, 120)
            _pic, fh = picture_fit(slide, fig, ML, y + 14, COL, fh)
            y += fh + 14
        if below:
            cw2 = (COL - 30 * (len(below) - 1)) / float(len(below))
            cy = bottom - cal_h + 12
            for i, c in enumerate(below):
                draw_callout(slide, theme, c, ML + i * (cw2 + 30), cy, cw2)

    elif kind == "keynumbers":
        # Row height follows the TALLER of the two columns, which is usually the
        # display number and not the statement. Sizing the row on the statement
        # alone is what put the hairlines through the figures: a 36pt number
        # occupies about 86px once its top offset is counted, against a 73px row
        # for a one-line statement, so every row overlapped the next by 13px all
        # the way down the page.
        items = sl["items"]
        sw = COL - 320
        vw = 300
        sizes = [_fit(str(v), 36, 8, 28) for v, _s in items]
        heights = [max(block_h(str(v), vw, vs, 1.05) + 20,
                       block_h(st, sw, 12.5, 1.4) + 26)
                   for (v, st), vs in zip(items, sizes)]
        room = bottom - y
        scale = min(1.0, room / max(sum(heights), 1))
        for (val, stmt), vs, h in zip(items, sizes, heights):
            h *= scale
            rule(slide, ML, y, COL, theme.rule, 0.75)
            val = str(val)
            # Both columns centred in the row rather than hung from its top. The
            # number sets at 36pt and the statement at 12.5, so a shared top edge
            # never reads as aligned; a shared centre does, whatever either one
            # does with its own height.
            text(slide, val, ML, y + 10, vw, h - 20, size=vs,
                 font=theme.serif, colour=theme.gold, line=1.05,
                 anchor=MSO_ANCHOR.MIDDLE)
            text(slide, stmt, ML + 320, y + 10, sw, h - 20, size=12.5,
                 font=theme.serif, colour=theme.fg, line=1.4,
                 anchor=MSO_ANCHOR.MIDDLE)
            y += h

    elif kind == "table":
        y2 = draw_table(slide, theme, sl["table"], ML, y, main_w, bottom - y)
        if full_width and callouts:
            cy = y2 + 20
            for c in callouts:
                cy = draw_callout(slide, theme, c, ML, cy, COL) + 10
        if sl.get("bullets"):
            yy = y2 + 18
            for b in sl["bullets"]:
                h = block_h(b, main_w, 11.5, 1.4)
                text(slide, b, ML, yy, main_w, h, size=11.5,
                     font=theme.serif, colour=theme.mute, line=1.4)
                yy += h + 8

    elif kind == "figure":
        img = sl.get("image")
        fw = main_w
        # the bullets under a figure take their own height rather than a fixed
        # 60px, which clipped the third line of the route facts
        bull_h = sum(block_h(b, fw, 11.5, 1.35) + 6
                     for b in (sl.get("bullets") or [])) + 14
        fh = max(bottom - y - (bull_h if sl.get("bullets") else 0), 160)
        if img and os.path.exists(img):
            _pic, fh = picture_fit(slide, img, ML, y, fw, fh)
        else:
            fill(slide, PAPER_RULE, ML, y, fw, fh)
            text(slide, "Figure not generated", ML, y + fh / 2 - 10, fw, 30,
                 size=10, font=theme.mono, caps=True, colour=theme.label,
                 spacing=1.6, align=PP_ALIGN.CENTER)
        if sl.get("bullets"):
            yy = y + fh + 14
            for b in sl["bullets"]:
                h = block_h(b, fw, 11.5, 1.35)
                text(slide, b, ML, yy, fw, h, size=11.5, font=theme.serif,
                     colour=theme.mute, line=1.35)
                yy += h + 6

    elif kind == "prose":
        size = prose_size(sl["paras"], main_w, bottom - y)
        yy = y
        for para in sl["paras"]:
            if isinstance(para, (tuple, list)):
                head, body = para
                if head:
                    mono_label(slide, head, ML, yy, main_w, theme, gold=True)
                    yy += 26
            else:
                body = para
            # cap the measure: full-width 12.5pt runs to circa 110 characters
            # a line, which is too wide to read. 78 is the comfortable ceiling.
            pw = min(main_w, cpl_width(78, size))
            h = block_h(body, pw, size, 1.5)
            text(slide, body, ML, yy, pw, h, size=size, font=theme.serif,
                 colour=theme.fg, line=1.5)
            yy += h + 16

    elif kind == "grid":
        rows = sl["rows"]
        accent = set(sl.get("accent_rows") or [])
        cw = (COL - 44) / 2.0
        colx = [ML, ML + cw + 44]
        room = bottom - y - (70 if sl.get("callout") else 0)

        blocks = []
        for head, body in rows:
            hh = block_h(head, cw, 14.5, 1.25)
            bh = block_h(body, cw, 11.5, 1.42)
            blocks.append((head, body, hh, bh, hh + bh + 40))
        # split where the running height first passes half the total, so two
        # long entries do not stack in one column while the other sits empty
        total = sum(b[4] for b in blocks)
        run, cut = 0, len(blocks)
        for i, blk in enumerate(blocks):
            if run + blk[4] / 2.0 > total / 2.0:
                cut = i
                break
            run += blk[4]
        chunks = (blocks[:cut], blocks[cut:])
        scale = min(1.0, room / max(sum(b[4] for b in ch) or 1
                                    for ch in chunks))
        top, deepest = y, y
        for c, chunk in enumerate(chunks):
            yy = top                     # both columns start level
            for j, (head, body, hh, bh, _t) in enumerate(chunk):
                idx = j + (cut if c else 0)
                tone = theme.gold if idx in accent else theme.fg
                rule(slide, colx[c], yy, 46, tone, 1.6)
                text(slide, head, colx[c], yy + 14, cw, hh, size=14.5,
                     font=theme.serif, colour=tone, line=1.25)
                text(slide, body, colx[c], yy + 16 + hh, cw, bh, size=11.5,
                     font=theme.serif, colour=theme.mute, line=1.42)
                yy += (hh + bh + 40) * scale
            deepest = max(deepest, yy)
        y = deepest
        if sl.get("callout"):
            ch = block_h(sl["callout"]["lines"][0], COL - 26, 12.5, 1.35)
            draw_callout(slide, theme, sl["callout"], ML,
                         min(y + 22, bottom - ch - 12), COL)

    # the side column
    if has_side:
        yy = side_top
        for p in panels:
            yy = draw_panel(slide, theme, p, side_x, yy, side_w) + 18
        for c in callouts:
            yy = draw_callout(slide, theme, c, side_x, yy, side_w) + 12

    source_line(slide, theme, sl)
    return slide


def slide_plate(deck, spec, sl, meta, page, safe):
    """An evidence plate. Body grid only, one to a page, captioned."""
    theme = Theme(safe_fonts=safe)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    fill(slide, theme.bg)
    chrome(slide, theme, meta, page, sl.get("section"))
    y = heading(slide, theme, sl)
    src_h = block_h(sl["source"], COL, 8.5, 1.45, 1.2) if sl.get("source") else 0
    bottom = 1080 - MB - 6 - src_h

    img_w = COL * 0.58
    img_h = min(img_w * 2.0 / 3.0, bottom - y - 40)
    if sl.get("image") and os.path.exists(sl["image"]):
        picture_cover(slide, sl["image"], ML, y, img_w, img_h)
        # the house hairline frame, square corners
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, X(ML), X(y),
                                       X(img_w), X(img_h))
        frame.fill.background()
        frame.line.color.rgb = rgb(theme.gold)
        frame.line.width = Pt(0.75)
        frame.shadow.inherit = False
    else:
        # rule 04: a mono-ruled diagram, never a stock substitute
        fill(slide, theme.bg, ML, y, img_w, img_h)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, X(ML), X(y),
                                     X(img_w), X(img_h))
        box.fill.background()
        box.line.color.rgb = rgb(theme.rule)
        box.line.width = Pt(0.75)
        box.shadow.inherit = False
        for i in range(1, 6):
            rule(slide, ML, y + img_h * i / 6.0, img_w, theme.rule, 0.75)
        text(slide, "No cleared photograph for this slot", ML,
             y + img_h / 2 - 12, img_w, 30, size=10, font=theme.mono,
             caps=True, colour=theme.label, spacing=1.6,
             align=PP_ALIGN.CENTER)

    cap = "  /  ".join([p for p in (sl.get("subject"), sl.get("date"),
                                    sl.get("credit")) if p])
    text(slide, cap, ML, y + img_h + 12, img_w, 24, size=9, font=theme.mono,
         caps=True, colour=theme.label, spacing=1.6)

    sx, sw = ML + img_w + 56, COL - img_w - 56
    yy = y
    if sl.get("supports"):
        mono_label(slide, "What this shows", sx, yy, sw, theme, gold=True)
        yy += 26
        h = block_h(sl["supports"], sw, 12.5, 1.5)
        text(slide, sl["supports"], sx, yy, sw, h, size=12.5,
             font=theme.serif, colour=theme.fg, line=1.5)
        yy += h + 16
    for para in sl.get("body") or []:
        h = block_h(para, sw, 11.5, 1.45)
        text(slide, para, sx, yy, sw, h, size=11.5, font=theme.serif,
             colour=theme.mute, line=1.45)
        yy += h + 12
    for p in sl.get("panels") or []:
        yy = draw_panel(slide, theme, p, sx, yy + 8, sw) + 14
    source_line(slide, theme, sl)
    return slide


def slide_thanks(deck, spec, sl, meta, safe):
    theme = Theme(dark=True, safe_fonts=safe)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    fill(slide, theme.bg)
    if sl.get("image") and os.path.exists(sl["image"]):
        picture_cover(slide, sl["image"], 0, 0, PX, 1080)
        # the closing page is the one full-bleed frame, so the type needs a
        # panel of its own rather than sitting on the photograph
        fill(slide, theme.bg, 0, 640, PX, 440)
    tt = sl.get("title") or "Thank you"
    tsize = fit_to_width([tt], COL * 0.7, 38, 24)
    th = block_h(tt, COL * 0.7, tsize, 1.12)
    text(slide, tt, ML, 700, COL * 0.7, th, size=tsize, font=theme.serif,
         colour=theme.fg, line=1.12)
    if sl.get("strap"):
        sh = block_h(sl["strap"], COL * 0.62, 14, 1.4)
        text(slide, sl["strap"], ML, 710 + th, COL * 0.62, sh, size=14,
             font=theme.serif, colour=INK_MUTE, italic=True, line=1.4)
    rule(slide, ML, 900, COL, INK_RULE)
    text(slide, [meta.get("author") or AUTHOR, meta["confidentiality"]],
         ML, 920, COL, 60, size=9, font=theme.mono, caps=True,
         colour=INK_LABEL, spacing=1.6, line=1.7)
    return slide


# ---------------------------------------------------------------------------
class Assets(object):
    """Turn a spec image field into a path, as the HTML renderer does.

    A generated figure is not photography: charts and maps in the assets
    folder resolve straight to a file. Everything else is a slot, and the
    resolver decides, applying the rights rules in avia_images.auto_ok.
    """

    def __init__(self, assets_dir="assets", resolver=None, compress=True,
                 quality=82, max_px=2200):
        self.dir = assets_dir
        self.resolver = resolver
        self.compress = compress
        self.quality = quality
        self.max_px = max_px
        self._cache = {}
        self._tmp = None

    def __call__(self, name, kind="mood", subjects=None, family=None,
                 prefer=None):
        if not name:
            return None
        local = os.path.join(self.dir, name)
        if os.path.exists(local) and name.lower().endswith((".png", ".jpg")):
            return local
        if self.resolver is not None:
            p, _src = self.resolver.resolve(name, family=family,
                                            subjects=subjects, kind=kind,
                                            prefer=prefer)
            if not (p and os.path.exists(p)):
                return None
            return self._photo(p)
        return local if os.path.exists(local) else None

    def _photo(self, path):
        """Re-encode photography to JPEG. Charts and maps are left alone.

        The library frames are full-size PNG, which put the Liguria deck at
        20 MB, too big to send. A photograph loses nothing at JPEG 82 and the
        deck lands near a tenth of that. Line art would show the artefacts, so
        only the resolved photography goes through here.
        """
        if not self.compress:
            return path
        if path in self._cache:
            return self._cache[path]
        try:
            import tempfile
            from PIL import Image
            if self._tmp is None:
                self._tmp = tempfile.mkdtemp(prefix="avia_deck_")
            im = Image.open(path).convert("RGB")
            if max(im.size) > self.max_px:
                r = self.max_px / float(max(im.size))
                im = im.resize((max(1, int(im.width * r)),
                                max(1, int(im.height * r))), Image.LANCZOS)
            dest = os.path.join(self._tmp, "%s.jpg"
                                % os.path.splitext(os.path.basename(path))[0])
            im.save(dest, "JPEG", quality=self.quality, optimize=True,
                    progressive=True)
            self._cache[path] = dest
            return dest
        except Exception:
            return path


def render(spec, path, safe_fonts=False, credits=None, assets_dir="assets",
           resolver=None, compress=True):
    assets = Assets(assets_dir, resolver, compress=compress)
    deck = Presentation()
    deck.slide_width = Emu(int(SLIDE_W_IN * 914400))
    deck.slide_height = Emu(int(SLIDE_H_IN * 914400))
    meta = spec["meta"]

    page = 0
    for sl in spec["slides"]:
        kind = sl["type"]
        if sl.get("image"):
            sl = dict(sl, image=assets(
                sl["image"],
                kind="evidence" if kind == "plate" else "mood",
                subjects=sl.get("subjects"), family=sl.get("family"),
                prefer=sl.get("prefer")))
        if kind == "plate" and not sl.get("image"):
            sl = dict(sl, image=assets(sl.get("slot"), kind="evidence",
                                       subjects=sl.get("subjects")))
        if kind == "cover":
            slide_cover(deck, spec, sl, meta, safe_fonts)
            continue
        page += 1
        if kind == "divider":
            slide_divider(deck, spec, sl, meta, page, safe_fonts)
        elif kind == "contents":
            slide_contents(deck, spec, sl, meta, page, safe_fonts)
        elif kind == "plate":
            slide_plate(deck, spec, sl, meta, page, safe_fonts)
        elif kind == "thanks":
            slide_thanks(deck, spec, sl, meta, safe_fonts)
        else:
            slide_body(deck, spec, sl, meta, page, safe_fonts)
        if sl.get("notes"):
            deck.slides[-1].notes_slide.notes_text_frame.text = sl["notes"]

    # after rendering, so the list is what was actually placed
    credits = credits or (resolver.credits() if resolver else None)
    if credits:
        _credits_slide(deck, meta, credits, safe_fonts)

    _metadata(deck, meta)
    deck.save(path)
    _language(path)
    return path


def _credits_slide(deck, meta, credits, safe):
    theme = Theme(safe_fonts=safe)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    fill(slide, theme.bg)
    chrome(slide, theme, meta)
    y = heading(slide, theme, {"title": "Image credits",
                               "subtitle": "Required attribution for the "
                                           "photography carried in this "
                                           "document."})
    for line in credits:
        text(slide, line, ML, y, COL, 26, size=10.5, font=theme.mono,
             colour=theme.mute, line=1.4)
        y += 24


def _metadata(deck, meta):
    """Never the generating library.

    These decks are published by The Aviation Observatory, a separate company from
    Avia Solutions, so they carry the Observatory as author. That is a deliberate
    departure from the Avia house rule and applies to product output only: an Avia
    client deliverable still carries Avia Solutions. The name comes from the spec's
    meta, so a deck built for either can set it.
    """
    cp = deck.core_properties
    cp.author = meta.get("author") or AUTHOR
    cp.last_modified_by = cp.author
    cp.company = cp.author
    cp.title = "%s - %s" % (meta["codename"], meta["title"])
    cp.category = meta["confidentiality"]
    cp.comments = ""


_LANG_PARTS = ("ppt/slides/slide", "ppt/slideLayouts/slideLayout",
               "ppt/slideMasters/slideMaster", "ppt/notesSlides/notesSlide",
               "ppt/notesMasters/notesMaster", "ppt/presentation.xml")


def _language(path):
    """Force en-GB on every run.

    python-pptx writes no language on the text it creates, so PowerPoint
    auto-detects each paragraph and regularly tags UK English as French. A
    run-level language beats the document default, so setting it in the
    interface reverts at once. This rewrites the packed XML, which is the only
    place it sticks.

    The masters, the layouts and presentation.xml are rewritten too, not the
    slides alone. python-pptx ships a default template that carries en-US in all
    three, and while that does not change how the built text is tagged, it is
    what any paragraph a user types into the deck afterwards inherits. Fixing
    the slides and leaving the master sets the file up to drift back to US
    English the first time somebody edits it.
    """
    import shutil
    import zipfile
    tmp = path + ".tmp"
    with zipfile.ZipFile(path) as zin, \
            zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith(_LANG_PARTS) \
                    and item.filename.endswith(".xml"):
                xml = data.decode("utf-8")
                xml = re.sub(r'\slang="[^"]*"', '', xml)
                xml = re.sub(r'\saltLang="[^"]*"', '', xml)
                for tag in ("a:rPr", "a:endParaRPr", "a:defRPr"):
                    xml = re.sub(r'<%s(?![^>]*\blang=)' % tag,
                                 '<%s lang="en-GB"' % tag, xml)
                data = xml.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(tmp, path)
    return path


def verify(path):
    """Post-build checks. Fails loud rather than shipping something wrong."""
    import zipfile
    from pptx import Presentation as P
    problems = []
    d = P(path)
    cp = d.core_properties
    if cp.author not in PUBLISHERS or cp.last_modified_by != cp.author:
        problems.append("author metadata is %r / %r, not one of %s"
                        % (cp.author, cp.last_modified_by, ", ".join(sorted(PUBLISHERS))))
    with zipfile.ZipFile(path) as z:
        bad = set()
        for n in z.namelist():
            if n.startswith("ppt/slides/slide") and n.endswith(".xml"):
                for m in re.findall(r'lang="([^"]+)"', z.read(n).decode("utf-8")):
                    if m != "en-GB":
                        bad.add(m)
        if bad:
            problems.append("proofing language present other than en-GB: %s"
                            % ", ".join(sorted(bad)))
    text_all = []
    for s in d.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                text_all.append(sh.text_frame.text)
    blob = "\n".join(text_all)
    for ch in ("—", "–"):
        if ch in blob:
            problems.append("%s dash found in the deck text"
                            % ("em" if ch == "—" else "en"))
    return problems


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("spec", help="python module exposing build() -> spec dict")
    ap.add_argument("-o", "--out", required=True)
    ap.add_argument("--assets", default="assets")
    ap.add_argument("--library", default="observatory_library",
                    help="the Observatory brand library")
    ap.add_argument("--store", help="the licensed subject store")
    ap.add_argument("--uploads", help="the client's uploaded photography")
    ap.add_argument("--region", default=None,
                    help="globe region for the cover, e.g. europe")
    ap.add_argument("--use", default="confidential",
                    choices=("confidential", "published"),
                    help="panorama risk blocks a published use")
    ap.add_argument("--embed-fonts", metavar="DIR", nargs="?", const="",
                    help="embed the staged brand faces. DIR is optional; the "
                         "default is font_store from avia_config.json, since "
                         "font binaries are data and never live in the repo")
    ap.add_argument("--full-size-images", action="store_true",
                    help="skip JPEG re-encoding of photography")
    ap.add_argument("--safe-fonts", action="store_true",
                    help="Cambria and Arial instead of Newsreader and IBM "
                         "Plex Mono, for a .pptx a client opens themselves")
    a = ap.parse_args()

    import importlib
    mod = importlib.import_module(a.spec.replace(".py", ""))
    spec = mod.build()
    S.paginate(spec)
    resolver = None
    if os.path.isdir(a.library) or a.store or a.uploads:
        import avia_slots
        resolver = avia_slots.SlotResolver(
            uploads_dir=a.uploads, subject_store=a.store,
            brand_library=a.library if os.path.isdir(a.library) else None,
            project=spec["meta"]["codename"], region=a.region, use=a.use)
    render(spec, a.out, a.safe_fonts, assets_dir=a.assets, resolver=resolver,
           compress=not a.full_size_images)
    if a.embed_fonts is not None:
        import avia_fonts
        store = avia_fonts.font_store(a.embed_fonts or None)
        faces = avia_fonts.staged(store)
        if not faces:
            print("   NOTE: no staged fonts in %s, so none embedded. Run "
                  "avia_fonts.py prepare first." % store)
        else:
            avia_fonts.embed(a.out, faces)
            for k, v in avia_fonts.check(a.out):
                print("   font: %-22s %s" % (k, v))
    problems = verify(a.out)
    if resolver:
        print(resolver.report())
    print("wrote %s (%d slides)" % (a.out, len(spec["slides"])))
    for p in problems:
        print("   PROBLEM: %s" % p)
    if not problems:
        print("   checks passed: metadata, en-GB, no em or en dashes")


if __name__ == "__main__":
    main()
