#!/usr/bin/env python3
"""Genoa-New York forecast + P&L deck (Avia house style). Numbers from the validated
genoa_nyc.py run (realistic premium case)."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2E, 0x5E, 0x8C)
GREY = RGBColor(0x59, 0x59, 0x59)
LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
GREEN = RGBColor(0x00, 0x61, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = "Arial"
out = sys.argv[1] if len(sys.argv) > 1 else "Genoa_NYC_forecast_and_PnL.pptx"

DISCLAIMER = ("Indicative, for directional guidance only. Built on generic published assumptions and "
              "Sabre/road-time data, not any airline's actual LOPA, contract terms or internal P&L. Actual results will differ.")

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def rect(slide, l, t, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def txt(slide, l, t, w, h, text, size, *, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
        italic=False, anchor=MSO_ANCHOR.TOP, font=BODY):
    b = slide.shapes.add_textbox(l, t, w, h); tf = b.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font; f.color.rgb = color
    return b


def stat(slide, l, t, w, value, label, vcolor=NAVY, vsize=34):
    b = slide.shapes.add_textbox(l, t, w, Inches(1.2)); tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = value
    r.font.size = Pt(vsize); r.font.bold = True; r.font.name = BODY; r.font.color.rgb = vcolor
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(11.5); r2.font.name = BODY; r2.font.color.rgb = GREY
    return b


# ---------------------------------------------------------------- 1) title
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SW, SH, NAVY)
txt(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.6), "Genoa - New York", 54, bold=True, color=WHITE)
txt(s, Inches(0.95), Inches(3.7), Inches(11.5), Inches(0.8), "A321XLR route assessment: catchment, demand and economics", 22, color=LIGHT)
txt(s, Inches(0.95), Inches(6.5), Inches(6), Inches(0.5), "Avia Solutions", 14, bold=True, color=WHITE)
txt(s, Inches(8.0), Inches(6.5), Inches(4.4), Inches(0.5), "Indicative, directional guidance", 11, color=LIGHT, align=PP_ALIGN.RIGHT, italic=True)

# ---------------------------------------------------------------- 2) forecast
s = prs.slides.add_slide(blank)
txt(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.7), "Catchment and demand: New York from Genoa", 32, bold=True, color=NAVY)
txt(s, Inches(0.62), Inches(1.15), Inches(12.1), Inches(0.5),
    "Genoa's New York travellers leak to Milan today; a GOA nonstop repatriates its own catchment", 15, color=GREY)
# stat row
stat(s, Inches(0.6), Inches(1.95), Inches(3.0), "19.5m", "catchment population (220 km)")
stat(s, Inches(3.7), Inches(1.95), Inches(3.0), "553,000", "New York O&D today (per year)")
stat(s, Inches(6.8), Inches(1.95), Inches(3.0), "85,500", "GOA catchment leaking to Milan")
stat(s, Inches(9.9), Inches(1.95), Inches(3.0), "55,600", "repatriable by a nonstop", vcolor=GREEN)
# leakage bars - where Genoa-region NYC demand departs today
txt(s, Inches(0.6), Inches(3.35), Inches(6.2), Inches(0.4), "Where the region's New York demand departs today", 15, bold=True, color=NAVY)
split = [("Milan (MXP)", 0.784), ("Milan Linate", 0.100), ("Bologna", 0.080),
         ("Turin", 0.023), ("Genoa (GOA)", 0.013), ("Bergamo", 0.002)]
y = Inches(3.85); barL = Inches(2.4); maxW = 4.0
for name, sh in split:
    txt(s, Inches(0.6), y, Inches(1.8), Inches(0.32), name, 12, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    w = Inches(max(maxW * sh, 0.04))
    col = GREEN if name.startswith("Genoa") else BLUE
    rect(s, barL, y + Inches(0.03), w, Inches(0.26), col)
    txt(s, barL + w + Inches(0.05), y, Inches(1.0), Inches(0.32), f"{sh:.0%}", 12, bold=True, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.46)
# right: the repatriation story
rx = Inches(7.4)
rect(s, rx, Inches(3.7), Inches(5.3), Inches(2.95), LIGHT)
txt(s, rx + Inches(0.25), Inches(3.85), Inches(4.9), Inches(0.4), "A Genoa nonstop", 16, bold=True, color=NAVY)
lines = ("GOA's New York catchment: 92,500 / year\n"
         "Carried by GOA today: 7,000  (1.3%)\n"
         "Leaking to Milan: 85,500\n"
         "Repatriated at 65% capture: 55,600 each way\n"
         "Directional demand with a nonstop: 62,600\n"
         "Fills a daily A321XLR at 85% load (≈6,300 spilled)")
txt(s, rx + Inches(0.25), Inches(4.35), Inches(4.9), Inches(2.2), lines, 13.5, color=NAVY)
txt(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4),
    "Source: GeoNames population, OSRM road times, Sabre point-of-origin O&D; catchment calibrated to the observed airport split.",
    9.5, color=GREY, italic=True)

# ---------------------------------------------------------------- 3) P&L
s = prs.slides.add_slide(blank)
txt(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.7), "Route economics: A321XLR, Genoa - New York", 32, bold=True, color=NAVY)
txt(s, Inches(0.62), Inches(1.15), Inches(12.1), Inches(0.5),
    "Daily service, one-way economy fare $345, planned load factor 85%; realistic premium case", 15, color=GREY)
stat(s, Inches(0.6), Inches(1.95), Inches(3.0), "$15,200", "operating profit per rotation")
stat(s, Inches(3.7), Inches(1.95), Inches(3.0), "10.7%", "operating margin", vcolor=GREEN)
stat(s, Inches(6.8), Inches(1.95), Inches(3.0), "74%", "breakeven load factor")
stat(s, Inches(9.9), Inches(1.95), Inches(3.0), "$5.5m", "annual profit (daily)", vcolor=GREEN)
# per-rotation summary (left)
txt(s, Inches(0.6), Inches(3.35), Inches(5.7), Inches(0.4), "Per rotation (return)", 15, bold=True, color=NAVY)
pnl = [("Revenue", "$141,800"), ("Operating cost", "$126,600"), ("Operating profit", "$15,200"), ("Margin", "10.7%")]
y = Inches(3.85)
for k, v in pnl:
    rect(s, Inches(0.6), y, Inches(5.4), Inches(0.45), LIGHT if k != "Operating profit" else NAVY)
    c = WHITE if k == "Operating profit" else NAVY
    txt(s, Inches(0.75), y, Inches(3.6), Inches(0.45), k, 13, bold=(k in ("Operating profit", "Margin")), color=c, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.2), y, Inches(1.65), Inches(0.45), v, 13, bold=True, color=c, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.52)
# assumptions + sensitivity (right)
rx = Inches(6.6)
txt(s, rx, Inches(3.35), Inches(6.1), Inches(0.4), "Assumptions and sensitivity", 15, bold=True, color=NAVY)
asum = ("Aircraft: A321XLR, GOA-JFK 3,500 nm, daily\n"
        "Fare: $345 one-way economy (Sabre, round-trip halved)\n"
        "Load factor: 85% planned (demand fills the aircraft)\n"
        "Cabin mix: economy-led with modest premium")
txt(s, rx, Inches(3.85), Inches(6.1), Inches(1.5), asum, 13, color=NAVY)
# sensitivity boxes
rect(s, rx, Inches(5.45), Inches(2.95), Inches(1.25), LIGHT)
txt(s, rx + Inches(0.2), Inches(5.55), Inches(2.6), Inches(0.4), "Realistic premium", 13, bold=True, color=NAVY)
txt(s, rx + Inches(0.2), Inches(5.95), Inches(2.6), Inches(0.7), "10.7% margin\n$5.5m / year", 14, bold=True, color=GREEN)
rect(s, rx + Inches(3.15), Inches(5.45), Inches(2.95), Inches(1.25), LIGHT)
txt(s, rx + Inches(3.35), Inches(5.55), Inches(2.6), Inches(0.4), "Full business cabin", 13, bold=True, color=NAVY)
txt(s, rx + Inches(3.35), Inches(5.95), Inches(2.6), Inches(0.7), "19.9% margin\n$11.6m / year", 14, bold=True, color=NAVY)
txt(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4), DISCLAIMER, 9.5, color=GREY, italic=True)

cp = prs.core_properties
cp.author = "Avia Solutions"; cp.last_modified_by = "Avia Solutions"; cp.title = "Genoa-New York route assessment"
prs.save(out)
print("saved", out)
