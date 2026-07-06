#!/usr/bin/env python3
"""
Avia Solutions - Route economics P&L slide.
Renders an aircraft_economics RoutePnL as a single, on-brand PowerPoint slide
(turnaround P&L) for dropping into a client deck. Carries the standard disclaimer.

  from economics_slide import slide_from_route_pnl
  slide_from_route_pnl(route_pnl, "Genoa_NYC_economics.pptx", title="Genoa-New York: A321XLR economics")
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1F, 0x38, 0x64)
GREY = RGBColor(0x59, 0x59, 0x59)
LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
GREEN = RGBColor(0x00, 0x61, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = "Arial"

try:
    from aircraft_economics import DISCLAIMER_SHORT
except Exception:
    DISCLAIMER_SHORT = ("Indicative, for directional guidance only. Built on generic published "
                        "assumptions, not any airline's actual LOPA, MTOW, contract terms or internal P&L. "
                        "Actual results will differ.")


def _txt(box, text, size, *, bold=False, color=NAVY, align=PP_ALIGN.LEFT, italic=False, font=BODY):
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font; f.color.rgb = color
    return tf


def _stat(slide, left, top, w, h, value, label, value_color=NAVY):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = value
    r.font.size = Pt(40); r.font.bold = True; r.font.name = BODY; r.font.color.rgb = value_color
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(12); r2.font.name = BODY; r2.font.color.rgb = GREY
    return box


def slide_from_compute(x, meta, out_path):
    """Build the slide from a compute() dict (x) and a meta dict
    (airline, origin, dest, aircraft, airline_type, sector_h, title)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    title = meta.get('title') or f"{meta.get('origin','')}-{meta.get('dest','')}: {meta.get('aircraft','')} route economics"
    _txt(slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7)),
         title, 30, bold=True, color=NAVY)
    sub = (f"Turnaround economics, indicative. {meta.get('airline','')}  "
           f"{meta.get('aircraft','')}  {meta.get('airline_type','')} operation  "
           f"sector {meta.get('sector_h',0):.1f}h.")
    _txt(slide.shapes.add_textbox(Inches(0.55), Inches(1.02), Inches(12.2), Inches(0.4)),
         sub, 13, color=GREY)

    # ---- Left column: stat callouts ----
    lx = Inches(0.55)
    _stat(slide, lx, Inches(1.75), Inches(4.4), Inches(1.1),
          f"${x['profit']:,.0f}", "PROFIT PER TURNAROUND (standalone)")
    _stat(slide, lx, Inches(2.95), Inches(2.1), Inches(1.0), f"{x['margin']:.1%}", "MARGIN")
    _stat(slide, Inches(2.85), Inches(2.95), Inches(2.1), Inches(1.0), f"{x['breakeven_lf']:.0%}", "BREAKEVEN LF")

    if x.get('incentive_value'):
        card = slide.shapes.add_shape(1, lx, Inches(4.2), Inches(4.4), Inches(1.25))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.fill.background()
        tf = card.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(10); tf.margin_top = Pt(8)
        p = tf.paragraphs[0]; r = p.add_run(); r.text = "With airport incentive"
        r.font.size = Pt(13); r.font.bold = True; r.font.name = BODY; r.font.color.rgb = NAVY
        p2 = tf.add_paragraph(); r2 = p2.add_run()
        r2.text = f"Profit ${x['profit_with_incentive']:,.0f}   Margin {x['margin_with_incentive']:.1%}"
        r2.font.size = Pt(15); r2.font.bold = True; r2.font.name = BODY; r2.font.color.rgb = GREEN
        p3 = tf.add_paragraph(); r3 = p3.add_run()
        r3.text = f"Airport contributes ${x['incentive_value']:,.0f} per turn"
        r3.font.size = Pt(11); r3.font.name = BODY; r3.font.color.rgb = GREY

    # ---- Right column: P&L table ----
    other_var = x['variable'] - x['fuel'] - x['maintenance']
    crew_ins = x['crew'] + x['insurance']
    rows = [
        ("Gross revenue", x['gross_rev'], True, False),
        ("Fuel", x['fuel'], False, False),
        ("Maintenance", x['maintenance'], False, False),
        ("Other variable (catering, charges, nav, handling)", other_var, False, False),
        ("Ownership (cost of capital)", x['ownership'], False, False),
        ("Crew & insurance", crew_ins, False, False),
        ("Overheads (admin & sales)", x['indirect_fixed'], False, False),
        ("Total cost", x['total_cost'], True, False),
        ("Profit per turnaround", x['profit'], True, True),
    ]
    tx, ty, tw = Inches(5.35), Inches(1.75), Inches(7.45)
    rh = Inches(0.46)
    tbl = slide.shapes.add_table(len(rows) + 1, 2, tx, ty, tw, rh).table
    tbl.columns[0].width = Inches(5.35); tbl.columns[1].width = Inches(2.10)
    # header
    for j, h in enumerate(["Turnaround P&L (USD)", "Per turn"]):
        c = tbl.cell(0, j); c.text = ""
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        run = pr.add_run(); run.text = h
        run.font.size = Pt(12); run.font.bold = True; run.font.name = BODY; run.font.color.rgb = WHITE
    for i, (label, val, bold, profit) in enumerate(rows, 1):
        c0 = tbl.cell(i, 0); c1 = tbl.cell(i, 1)
        for c in (c0, c1):
            c.fill.solid(); c.fill.fore_color.rgb = WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = c0.text_frame.paragraphs[0]; p0.alignment = PP_ALIGN.LEFT
        r0 = p0.add_run(); r0.text = label
        r0.font.size = Pt(12); r0.font.bold = bold; r0.font.name = BODY
        r0.font.color.rgb = GREEN if profit else NAVY
        p1 = c1.text_frame.paragraphs[0]; p1.alignment = PP_ALIGN.RIGHT
        r1 = p1.add_run(); r1.text = f"{val:,.0f}"
        r1.font.size = Pt(12); r1.font.bold = bold; r1.font.name = BODY
        r1.font.color.rgb = GREEN if profit else NAVY

    # ---- Provenance + disclaimer ----
    prov = (f"Maintenance: {x.get('maint_basis','')}.  Ownership: {x.get('own_basis','')}.")
    _txt(slide.shapes.add_textbox(Inches(0.55), Inches(6.35), Inches(12.25), Inches(0.5)),
         prov, 9, italic=True, color=GREY)
    _txt(slide.shapes.add_textbox(Inches(0.55), Inches(6.95), Inches(12.25), Inches(0.45)),
         DISCLAIMER_SHORT, 9, italic=True, color=GREY)

    prs.core_properties.author = "Avia Solutions"
    prs.core_properties.title = title
    prs.save(out_path)
    return out_path


def slide_from_route_pnl(route_pnl, out_path, title=None):
    x = route_pnl.compute()
    meta = dict(airline=route_pnl.airline, origin=route_pnl.origin, dest=route_pnl.dest,
                aircraft=route_pnl.aircraft, airline_type=route_pnl.airline_type,
                sector_h=x.get('sector_fh', 0), title=title)
    return slide_from_compute(x, meta, out_path)


def slide_from_network(net, out_path, title=None, subtitle=None):
    """Build a Network P&L slide from an aircraft_economics.network_pnl() result."""
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _txt(slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.7)),
         title or "Network economics", 30, bold=True, color=NAVY)
    n_routes = len(net.get('routes', []))
    sub = subtitle or f"Annual network P&L, indicative. {n_routes} routes sharing one fleet."
    _txt(slide.shapes.add_textbox(Inches(0.55), Inches(1.02), Inches(12.2), Inches(0.4)),
         sub, 13, color=GREY)

    # Left: network stat callouts
    lx = Inches(0.55)
    _stat(slide, lx, Inches(1.8), Inches(4.4), Inches(1.1),
          f"${net['annual_profit']/1e6:,.2f}m", "ANNUAL NETWORK PROFIT")
    _stat(slide, lx, Inches(3.05), Inches(2.1), Inches(1.0), f"{net['margin']:.1%}", "NETWORK MARGIN")
    _stat(slide, Inches(2.85), Inches(3.05), Inches(2.1), Inches(1.0),
          f"{net['aircraft_required']}", "AIRCRAFT (SHARED FLEET)")
    _txt(slide.shapes.add_textbox(lx, Inches(4.25), Inches(4.4), Inches(0.9)),
         f"{net['aircraft_required_fractional']:.2f} frames of flying; one aircraft covers the "
         f"base, with spare capacity across the network.", 12, color=GREY, italic=True)

    # Right: per-route table
    rows = net.get('routes', [])
    tbl = slide.shapes.add_table(len(rows) + 2, 4, Inches(5.35), Inches(1.8),
                                 Inches(7.45), Inches(0.46)).table
    for w, col in zip((3.55, 1.7, 1.1, 1.1), range(4)):
        tbl.columns[col].width = Inches(w)
    for j, h in enumerate(["Route", "Annual profit", "Margin", "Frames"]):
        c = tbl.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = NAVY
        pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        run = pr.add_run(); run.text = h
        run.font.size = Pt(12); run.font.bold = True; run.font.name = BODY; run.font.color.rgb = WHITE
    for i, (o, d, ac, c) in enumerate(rows, 1):
        vals = [f"{o}-{d}  {ac}", f"{c['annual_profit']:,.0f}", f"{c['margin']:.1%}",
                f"{c['aircraft_required_fractional']:.2f}"]
        for j, v in enumerate(vals):
            cell = tbl.cell(i, j); cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
            r = p.add_run(); r.text = v
            r.font.size = Pt(12); r.font.name = BODY; r.font.color.rgb = NAVY
    # total row
    ti = len(rows) + 1
    tvals = ["Network total", f"{net['annual_profit']:,.0f}", f"{net['margin']:.1%}",
             f"{net['aircraft_required_fractional']:.2f}"]
    for j, v in enumerate(tvals):
        cell = tbl.cell(ti, j); cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        r = p.add_run(); r.text = v
        r.font.size = Pt(12); r.font.bold = True; r.font.name = BODY
        r.font.color.rgb = GREEN if j == 1 else NAVY

    _txt(slide.shapes.add_textbox(Inches(0.55), Inches(6.4), Inches(12.25), Inches(0.4)),
         "Costs on the current-data stack: maintenance (Airbus reserves, sector), ownership "
         "(appraiser lease/value, sector and age), utilisation (Airbus, sector) and crew (by airline model).",
         9, italic=True, color=GREY)
    _txt(slide.shapes.add_textbox(Inches(0.55), Inches(6.95), Inches(12.25), Inches(0.45)),
         DISCLAIMER_SHORT, 9, italic=True, color=GREY)

    prs.core_properties.author = "Avia Solutions"
    prs.core_properties.title = title or "Network economics"
    prs.save(out_path)
    return out_path
