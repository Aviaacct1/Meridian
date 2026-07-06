#!/usr/bin/env python3
"""Plain-English explainer of the Genoa-New York terms. One-off."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1F, 0x38, 0x64); BLUE = RGBColor(0x2E, 0x5E, 0x8C)
GREY = RGBColor(0x4A, 0x4A, 0x4A); LIGHT = RGBColor(0xEC, 0xF1, 0xF8)
GREEN = RGBColor(0x00, 0x61, 0x00); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xB8, 0x86, 0x0B); BODY = "Arial"
out = sys.argv[1] if len(sys.argv) > 1 else "Genoa_NYC_explainer.pptx"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height; blank = prs.slide_layouts[6]


def rect(s, l, t, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def txt(s, l, t, w, h, text, size, *, bold=False, color=NAVY, align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(l, t, w, h); tf = b.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(3); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = BODY; f.color.rgb = color
    return b


def card(s, l, t, w, h, term, value, definition, vcolor=GREEN):
    rect(s, l, t, w, h, LIGHT)
    txt(s, l + Inches(0.18), t + Inches(0.1), w - Inches(2.0), Inches(0.4), term, 14.5, bold=True, color=NAVY)
    txt(s, l + w - Inches(2.0), t + Inches(0.1), Inches(1.85), Inches(0.4), value, 15, bold=True, color=vcolor, align=PP_ALIGN.RIGHT)
    txt(s, l + Inches(0.18), t + Inches(0.52), w - Inches(0.36), h - Inches(0.6), definition, 11.5, color=GREY)


def header(s, title, sub):
    txt(s, Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.65), title, 30, bold=True, color=NAVY)
    txt(s, Inches(0.62), Inches(1.05), Inches(12.1), Inches(0.45), sub, 14, color=GREY)


# ---- 1 title
s = prs.slides.add_slide(blank); rect(s, 0, 0, SW, SH, NAVY)
txt(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.7), "Reading the numbers", 50, bold=True, color=WHITE)
txt(s, Inches(0.95), Inches(3.8), Inches(11.5), Inches(0.8), "Genoa - New York: what each term means, in plain English", 22, color=RGBColor(0xCA, 0xDC, 0xFC))
txt(s, Inches(0.95), Inches(6.5), Inches(6), Inches(0.5), "Avia Solutions", 14, bold=True, color=WHITE)

# ---- 2 demand terms
s = prs.slides.add_slide(blank)
header(s, "The demand, in plain English", "From the people near Genoa to the passengers a New York nonstop would carry")
L, R = Inches(0.6), Inches(6.85); W = Inches(5.9); H = Inches(1.5)
rows = [Inches(1.65), Inches(3.3), Inches(4.95)]
card(s, L, rows[0], W, H, "Catchment population", "3.3m",
     "The people for whom Genoa is the nearest, most convenient airport, out to roughly a 2-hour drive.", vcolor=NAVY)
card(s, R, rows[0], W, H, "Total air demand at Genoa", "~1.0m / yr",
     "Air trips those people actually make from Genoa today, all destinations. Most longer-haul trips leak to Milan.", vcolor=NAVY)
card(s, L, rows[1], W, H, "New York demand (latent)", "92,500 / yr",
     "All New York trips the catchment generates per year, wherever they depart from today.", vcolor=NAVY)
card(s, R, rows[1], W, H, "Leakage", "85,500 / yr",
     "Catchment New York demand departing other airports today (mainly Milan) because Genoa has no nonstop.", vcolor=GOLD)
card(s, L, rows[2], W, H, "Capture & repatriation", "65% -> 55,600",
     "Share of that leaked demand a Genoa nonstop wins back (assumed 65%) = passengers repatriated, each way.", vcolor=GREEN)
card(s, R, rows[2], W, H, "Directional demand", "62,600 each way",
     "Passengers flying one way (Genoa to New York). The return leg is counted separately; together they fill the aircraft.", vcolor=NAVY)
txt(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.4),
    "Propensity: New York trips per person per year, about 28 per 1,000 people here, turns catchment population into demand.", 11, color=GREY, italic=True)

# ---- 3 economics terms + the two questions
s = prs.slides.add_slide(blank)
header(s, "The economics, in plain English", "What load factor, margin and breakeven actually mean for this route")
card(s, L, rows[0], W, H, "Load factor", "85% planned",
     "How full the aircraft is: seats sold divided by seats flown. We plan at 85%; no route runs 95% year-round.", vcolor=NAVY)
card(s, R, rows[0], W, H, "Operating margin", "10.7%",
     "Operating profit per $100 of revenue, AT the planned 85% load. Fuller lifts it, emptier lowers it.", vcolor=GREEN)
card(s, L, rows[1], W, H, "Breakeven load factor", "74%",
     "The fullness at which profit is exactly zero. Below 74% the route loses money; above it, it profits.", vcolor=GOLD)
card(s, R, rows[1], W, H, "Per rotation / annual", "$15,200 -> $5.5m",
     "A rotation is one return trip (GOA-JFK-GOA): $15,200 profit. Flown daily (364/yr) that is $5.5m a year.", vcolor=GREEN)
# the two direct questions
rect(s, L, rows[2], Inches(12.15), H, NAVY)
txt(s, L + Inches(0.22), rows[2] + Inches(0.12), Inches(11.7), Inches(0.4),
    "Your two questions, directly:", 13.5, bold=True, color=WHITE)
txt(s, L + Inches(0.22), rows[2] + Inches(0.52), Inches(11.7), Inches(0.9),
    "Is 10.7% margin if the aircraft is 85% full?   Yes, the margin is calculated at the 85% planned load.\n"
    "Breakeven 74% means zero profit there?   Yes, at 74% full the route just covers its costs; at 85% you sit 11 points above it, and that cushion is the 10.7% margin.",
    12.5, color=RGBColor(0xE6, 0xEE, 0xFA))

# ---- 4 the Genoa end funnel
s = prs.slides.add_slide(blank)
header(s, "The Genoa end: from catchment to carried passengers", "How the 3.3m people narrow down to the New York passengers a nonstop carries")
funnel = [("Catchment population (nearest to Genoa)", "3,300,000", NAVY),
          ("Air trips from Genoa today, all destinations", "~1,000,000 / yr", BLUE),
          ("New York demand the catchment generates", "92,500 / yr", BLUE),
          ("Carried by Genoa today (rest leak to Milan)", "7,000 / yr", GOLD),
          ("Repatriated by a Genoa nonstop (65% capture)", "55,600 each way", GREEN)]
y = Inches(1.7); fullW = 11.4
for i, (lab, val, col) in enumerate(funnel):
    w = Inches(fullW * (1 - i * 0.13))
    rect(s, Inches(0.6), y, w, Inches(0.78), col)
    txt(s, Inches(0.8), y, w - Inches(2.4), Inches(0.78), lab, 13.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(0.6) + w - Inches(2.3), y, Inches(2.1), Inches(0.78), val, 15, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.92)
txt(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.5),
    "New York is one destination within Genoa's catchment. The same method sizes any route; the catchment (3.3m) is shared, the destination demand differs.",
    11, color=GREY, italic=True)

cp = prs.core_properties; cp.author = "Avia Solutions"; cp.last_modified_by = "Avia Solutions"
cp.title = "Genoa-New York: plain-English explainer"
prs.save(out); print("saved", out)
