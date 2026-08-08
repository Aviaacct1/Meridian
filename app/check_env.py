#!/usr/bin/env python3
"""Is this Python actually able to run the tool. Answers in one screen.

    py -3.12 check_env.py

Written on 7 August 2026 after installing matplotlib and basemap into the global
3.12 downgraded numpy from 2.5.1 to 2.3.5 and could not delete the old copy,
leaving four ~-prefixed folders in site-packages. pip reported that as a warning
and exited zero. An environment can be broken and still look installed, and the
tool is the thing that finds out.

Three things are checked:

  1. WHERE this Python is, and whether it is a virtual environment. Installing
     tool dependencies into a shared global Python is how one tool's install
     breaks another's runtime.
  2. WHAT is importable, with versions, split into what the engine needs, what
     the deck generator needs, and what is optional.
  3. WHETHER a previous install left wreckage behind: ~-prefixed directories in
     site-packages are the Windows signature of an uninstall that could not
     delete locked files, and they shadow or corrupt the real package.

Exit code is 0 when everything required is present and working, 1 when it is
not, so this can gate a deploy.

Avia Solutions Limited. All rights reserved.
"""

import importlib
import os
import sys

# name to import, pip name, why. The pip name only differs where it has to.
ENGINE = [
    ("fastapi", "fastapi", "the portal server"),
    ("uvicorn", "uvicorn", "serves it"),
    ("duckdb", "duckdb", "every store read; the engine cannot run without it"),
    ("airportsdata", "airportsdata", "airport reference, used in 15 modules"),
    ("openpyxl", "openpyxl", "workbooks and the business case"),
    ("pptx", "python-pptx", "deck output"),
    ("numpy", "numpy", "catchment and drive-time maths"),
]
DECK = [
    ("matplotlib", "matplotlib", "every chart"),
    ("PIL", "pillow", "picture fitting and the JPEG re-encode"),
    ("fontTools", "fonttools", "embedded typefaces"),
]
OPTIONAL = [
    ("mpl_toolkits.basemap", "basemap", "the route map; without it the deck is one page shorter"),
    ("anthropic", "anthropic", "the live research pass; replay runs need no key"),
    ("global_land_mask", "global-land-mask", "water-boundary catchment check"),
    ("xlrd", "xlrd", "legacy .xls reads"),
    ("piexif", "piexif", "image rights metadata"),
    ("pandas", "pandas", "the legacy Streamlit tool only"),
    ("streamlit", "streamlit", "the legacy Streamlit tool only"),
    ("docx", "python-docx", "docx packs"),
]

OK, BAD = [], []


def _version(mod):
    for attr in ("__version__", "version", "VERSION"):
        v = getattr(mod, attr, None)
        if isinstance(v, str):
            return v
    return "?"


def group(title, items, required):
    print("\n%s" % title)
    worst = True
    for name, pip_name, why in items:
        try:
            m = importlib.import_module(name)
            print("   %-24s %-12s %s" % (name, _version(m), why))
            OK.append(name)
        except Exception as e:
            worst = False
            mark = "MISSING " if isinstance(e, ImportError) else "BROKEN  "
            print("   %-24s %-12s %s" % (name, mark, why))
            print("   %-24s %s" % ("", "%s: %s" % (type(e).__name__, e)))
            (BAD if required else OK).append(name)
            if required:
                MISSING_PIP.append(pip_name)
    return worst


MISSING_PIP = []


def leftovers():
    """~-prefixed folders in site-packages, from an uninstall that could not delete.

    pip prints "Failed to remove contents in a temporary directory", carries on,
    and exits zero. Two quite different things look identical in that message and
    this separates them, because a check that cries wolf gets ignored:

      ORPHANED   the real package is present beside it and imports. The leftover
                 is dead weight, a few megabytes, and cannot shadow anything
                 because "~umpy.libs" is not an importable name. Housekeeping.
      SHADOWING  the real package is absent, so the leftover is all that is left
                 of it and the import will fail or half work. A real fault.

    The smoke tests are the ground truth either way, which is why only they and a
    missing requirement set the exit code.
    """
    found = []
    seen = set()
    for p in sys.path:
        if not p or not p.endswith(("site-packages", "dist-packages")):
            continue
        real = os.path.realpath(p)
        if real in seen or not os.path.isdir(real):
            continue
        seen.add(real)
        for entry in os.listdir(real):
            if not entry.startswith("~"):
                continue
            # "~umpy.libs" lost its first character to the tilde, so the real
            # name is any sibling that matches on everything after it
            tail = entry[1:]
            twin = any(o != entry and o.endswith(tail) and len(o) == len(entry)
                       for o in os.listdir(real))
            found.append((os.path.join(real, entry),
                          "orphaned" if twin else "SHADOWING"))
    return found


def smoke():
    """Do the things the tool actually does, not just import the packages."""
    tests = []

    def run(label, fn):
        try:
            fn()
            tests.append((label, None))
        except Exception as e:
            tests.append((label, "%s: %s" % (type(e).__name__, e)))

    def _numpy():
        import numpy as np
        assert float(np.array([1.0, 2.0]).sum()) == 3.0

    def _mpl():
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        fig, ax = plt.subplots(figsize=(1, 1))
        ax.plot([0, 1], [0, 1])
        plt.close(fig)

    def _basemap():
        from mpl_toolkits.basemap import Basemap
        Basemap(projection="merc", llcrnrlon=-10, urcrnrlon=10,
                llcrnrlat=40, urcrnrlat=60, resolution=None)

    def _pptx():
        from pptx import Presentation
        Presentation().slides.add_slide(Presentation().slide_layouts[6])

    def _duckdb():
        import duckdb
        assert duckdb.connect().execute("select 1").fetchone()[0] == 1

    run("numpy arithmetic", _numpy)
    run("matplotlib draws a figure", _mpl)
    run("basemap builds a projection", _basemap)
    run("python-pptx makes a slide", _pptx)
    run("duckdb answers a query", _duckdb)
    return tests


def main():
    print("=" * 74)
    print("AVIA CORTEX QSI: environment check")
    print("=" * 74)
    print("interpreter   %s" % sys.executable)
    print("version       %s" % ".".join(str(v) for v in sys.version_info[:3]))
    venv = sys.prefix != getattr(sys, "base_prefix", sys.prefix)
    print("virtualenv    %s" % ("yes, %s" % sys.prefix if venv
                                else "NO, this is a shared Python"))
    if not venv:
        print("              Installing one tool's dependencies here changes every other")
        print("              tool that uses this interpreter. See requirements.txt.")

    group("ENGINE, required", ENGINE, True)
    group("DECK GENERATOR, required for the Observatory deck path", DECK, True)
    group("OPTIONAL", OPTIONAL, False)

    junk = leftovers()
    print("\nFAILED UNINSTALLS")
    if not junk:
        print("   none")
    else:
        for path, kind in junk:
            print("   %-10s %s" % (kind, path))
            if kind == "SHADOWING":
                BAD.append("shadowing leftover: %s" % os.path.basename(path))
        if any(k == "SHADOWING" for _p, k in junk):
            print("   A SHADOWING folder is all that is left of a package. Fix it.")
        if any(k == "orphaned" for _p, k in junk):
            print("   Orphaned folders are dead weight, not a fault: the real package is")
            print("   beside them and a ~-prefixed name cannot be imported. Delete them")
            print("   when convenient. If the delete is refused, a process still has the")
            print("   DLLs mapped; find it with")
            print("      Get-Process | ForEach-Object { $p=$_; try { $p.Modules |")
            print("        Where-Object { $_.FileName -like '*~*' } | ForEach-Object {")
            print("          '{0} (PID {1})' -f $p.ProcessName, $p.Id } } catch {} }")
            print("   Stop it, or clear them on the next reboot.")

    print("\nSMOKE TESTS")
    for label, err in smoke():
        if err:
            print("   FAIL  %-34s %s" % (label, err))
            if "basemap" not in label:
                BAD.append(label)
            else:
                print("         Optional. The deck builds without it, one page shorter.")
        else:
            print("   ok    %s" % label)

    print("\n" + "=" * 74)
    if BAD:
        print("NOT READY. %d problem(s): %s" % (len(BAD), ", ".join(sorted(set(BAD)))))
        if MISSING_PIP:
            print("\n   %s -m pip install %s"
                  % (sys.executable, " ".join(sorted(set(MISSING_PIP)))))
        print("=" * 74)
        return 1
    print("READY. Everything required is present and working.")
    print("=" * 74)
    return 0


if __name__ == "__main__":
    sys.exit(main())
