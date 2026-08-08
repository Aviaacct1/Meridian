"""
city_pair_pptx_generator.py
===========================
Python-based PPTX generator for Avia Solutions city pair presentations.
Uses python-pptx (no Node.js dependency).

Generates a branded presentation with:
  1. Cover slide
  2. Contents / agenda
  3. Executive summary
  4. Economic context
  5. Demand drivers (corporate, tourism, trade, etc.)
  6. Airport overview
  7. Route forecast summary
  8. Why this route (4 key points)
  9. Closing slide

Auto-populates from:
  - Route config (sidebar)
  - Research findings (executor)
  - Forecast results (pipeline)
"""

import os
import json
import re
import tempfile
from datetime import datetime
from typing import Dict, List, Optional, Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ============================================================================
# COLOUR PALETTE (Avia Solutions branding)
# ============================================================================
NAVY = RGBColor(0x1E, 0x27, 0x61) if HAS_PPTX else None
DARK_NAVY = RGBColor(0x14, 0x1B, 0x4D) if HAS_PPTX else None
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC) if HAS_PPTX else None
WHITE = RGBColor(0xFF, 0xFF, 0xFF) if HAS_PPTX else None
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2) if HAS_PPTX else None
ACCENT_GOLD = RGBColor(0xD4, 0xA5, 0x37) if HAS_PPTX else None
ACCENT_TEAL = RGBColor(0x02, 0x80, 0x90) if HAS_PPTX else None
MED_GREY = RGBColor(0x5A, 0x7A, 0x9A) if HAS_PPTX else None
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C) if HAS_PPTX else None

SLIDE_W = Inches(13.333) if HAS_PPTX else 0
SLIDE_H = Inches(7.5) if HAS_PPTX else 0


def _set_bg(slide, colour):
    """Set slide background to solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_text_box(slide, left, top, width, height, text,
                  font_size=14, font_colour=None, bold=False,
                  italic=False, alignment=PP_ALIGN.LEFT,
                  font_name="Calibri", valign=MSO_ANCHOR.TOP,
                  line_spacing=1.15):
    """Add a text box with consistent formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_colour or DARK_TEXT
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    tf.auto_size = None
    txBox.text_frame.paragraphs[0].space_after = Pt(0)
    txBox.text_frame.paragraphs[0].space_before = Pt(0)
    return txBox


def _add_shape_rect(slide, left, top, width, height, fill_colour,
                    border_colour=None, border_width=0):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if border_colour:
        shape.line.color.rgb = border_colour
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def _add_multiline(slide, left, top, width, height, lines,
                   font_size=12, font_colour=None, font_name="Calibri",
                   alignment=PP_ALIGN.LEFT, bold_first=False,
                   bullet=False, line_spacing=1.3):
    """Add text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_colour or DARK_TEXT
        p.font.name = font_name
        p.alignment = alignment
        p.line_spacing = Pt(font_size * line_spacing)
        if bold_first and i == 0:
            p.font.bold = True
        if bullet and i > 0:
            p.level = 0
    return txBox


# ============================================================================
# SLIDE BUILDERS
# ============================================================================

def _slide_cover(prs, config):
    """Slide 1: Cover with airline/route/date."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, NAVY)

    # Gold accent line at top
    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GOLD)

    # Title
    headline = config.get("headline", "")
    if not headline:
        headline = (f"{config.get('origin_city', '')} - {config.get('dest_city', '')}\n"
                    f"Route Assessment for {config.get('airline_name', '')}")

    _add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(2),
                  headline, font_size=40, font_colour=WHITE, bold=True,
                  alignment=PP_ALIGN.LEFT, font_name="Georgia")

    # Subtitle line
    subtitle = (f"{config.get('origin', '')} - {config.get('destination', '')} | "
                f"{config.get('frequency', 7)}x weekly | {config.get('aircraft_type', '')}")
    _add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.6),
                  subtitle, font_size=18, font_colour=ICE_BLUE,
                  alignment=PP_ALIGN.LEFT)

    # Date and branding
    _add_text_box(slide, Inches(0.8), Inches(5.2), Inches(5), Inches(0.5),
                  config.get("date", datetime.now().strftime("%B %Y")),
                  font_size=14, font_colour=MED_GREY)

    _add_text_box(slide, Inches(0.8), Inches(5.8), Inches(5), Inches(0.5),
                  "Prepared by Avia Solutions", font_size=12,
                  font_colour=MED_GREY, italic=True)

    if config.get("client_name"):
        _add_text_box(slide, Inches(0.8), Inches(6.2), Inches(5), Inches(0.5),
                      f"For {config['client_name']}", font_size=12,
                      font_colour=MED_GREY)

    # Bottom gold line
    _add_shape_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT_GOLD)


def _slide_contents(prs, config, sections):
    """Slide 2: Contents / agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                  "Contents", font_size=36, font_colour=NAVY, bold=True,
                  font_name="Georgia")

    # Gold accent line
    _add_shape_rect(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.04), ACCENT_GOLD)

    for i, (num, title) in enumerate(sections):
        y = 1.8 + i * 0.55
        # Number
        _add_text_box(slide, Inches(0.8), Inches(y), Inches(0.6), Inches(0.45),
                      num, font_size=16, font_colour=ACCENT_GOLD, bold=True)
        # Title
        _add_text_box(slide, Inches(1.5), Inches(y), Inches(8), Inches(0.45),
                      title, font_size=16, font_colour=DARK_TEXT)

    # Footer
    _add_text_box(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
                  f"Avia Solutions - {config.get('origin', '')}-{config.get('destination', '')} Assessment",
                  font_size=9, font_colour=MED_GREY)


def _slide_section_header(prs, section_num, section_title):
    """Section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_NAVY)

    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GOLD)

    _add_text_box(slide, Inches(0.8), Inches(2.5), Inches(2), Inches(1),
                  section_num, font_size=60, font_colour=ACCENT_GOLD,
                  bold=True, font_name="Georgia")

    _add_text_box(slide, Inches(0.8), Inches(3.8), Inches(10), Inches(1),
                  section_title, font_size=36, font_colour=WHITE,
                  bold=True, font_name="Georgia")

    _add_shape_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT_GOLD)


def _slide_executive_summary(prs, config, research):
    """Executive summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
                  "Executive Summary", font_size=32, font_colour=NAVY,
                  bold=True, font_name="Georgia")
    _add_shape_rect(slide, Inches(0.8), Inches(1.25), Inches(2), Inches(0.04), ACCENT_GOLD)

    # Summary text
    summary = research.get("executive_summary", "")
    if not summary:
        summary = (f"This assessment evaluates the commercial viability of "
                   f"{config.get('airline_name', '')} operating "
                   f"{config.get('origin_city', '')} ({config.get('origin', '')}) - "
                   f"{config.get('dest_city', '')} ({config.get('destination', '')}) "
                   f"at {config.get('frequency', 7)}x weekly frequency.")

    _add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(2),
                  summary, font_size=14, font_colour=DARK_TEXT,
                  line_spacing=1.4)

    # Key metrics boxes (if forecast available)
    forecast = config.get("forecast", {})
    if forecast:
        metrics = [
            ("Total Passengers", f"{forecast.get('grand_total', 0):,.0f}"),
            ("Load Factor", f"{forecast.get('load_factor', 0):.0%}" if isinstance(forecast.get('load_factor'), float)
             else str(forecast.get('load_factor', 'N/A'))),
            ("P2P / Connecting", f"{forecast.get('p2p_total', 0):,.0f} / {(forecast.get('cnx_home_total', 0) + forecast.get('cnx_dest_total', 0)):,.0f}"),
        ]

        for i, (label, value) in enumerate(metrics):
            x = 0.8 + i * 3.8
            _add_shape_rect(slide, Inches(x), Inches(4.2), Inches(3.2), Inches(1.8),
                            LIGHT_GREY, border_colour=ICE_BLUE, border_width=1)
            _add_text_box(slide, Inches(x + 0.2), Inches(4.4), Inches(2.8), Inches(0.4),
                          label, font_size=11, font_colour=MED_GREY, bold=True)
            _add_text_box(slide, Inches(x + 0.2), Inches(4.9), Inches(2.8), Inches(0.8),
                          str(value), font_size=28, font_colour=NAVY, bold=True)

    _add_text_box(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
                  "Avia Solutions - Confidential", font_size=9, font_colour=MED_GREY)


_VERIFICATION_COUNT = re.compile(
    r"^\s*(no independently sourced figures|\d+\s+sourced finding)", re.I)


def _is_verification_count(text):
    """True for pitch_verify.block_summary output, which is audit, not content."""
    return bool(text) and bool(_VERIFICATION_COUNT.match(str(text)))


def _slide_research_block(prs, config, block_id, block_name, findings, summary_text=""):
    """Content slide for a research block with findings."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                  block_name, font_size=28, font_colour=NAVY,
                  bold=True, font_name="Georgia")
    _add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.04), ACCENT_GOLD)

    # Summary paragraph. A caller may pass real prose, and that belongs here. What
    # does NOT belong here is the verification count, "3 sourced findings (3
    # verified against the cited page)", which pitch_report supplies for the audit
    # trail. It was printing as the opening line of every section, so the deck told
    # the reader how many facts had been checked instead of what they meant.
    if _is_verification_count(summary_text):
        summary_text = ""
    if summary_text:
        _add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1),
                      summary_text, font_size=13, font_colour=DARK_TEXT,
                      italic=True, line_spacing=1.4)
        findings_y = 2.7
    else:
        findings_y = 1.5

    # Findings as bullet points (max 5 per slide)
    display_findings = findings[:5] if findings else []
    lines = []
    sources = []
    for f in display_findings:
        claim = f.get("claim", "") if isinstance(f, dict) else getattr(f, "claim", "")
        value = f.get("value", "") if isinstance(f, dict) else getattr(f, "value", "")
        if value:
            lines.append(f"• {claim} ({value})")
        else:
            lines.append(f"• {claim}")

        # Collect source info
        if isinstance(f, dict):
            src = f.get("source_name", "")
            yr = f.get("year", "")
        else:
            src = ""
            if hasattr(f, "citations") and f.citations:
                src = f.citations[0].source_name
            yr = getattr(f, "year", "")
        if src:
            sources.append(f"{src} ({yr})" if yr else src)

    if lines:
        _add_multiline(slide, Inches(0.8), Inches(findings_y), Inches(11), Inches(3.5),
                       lines, font_size=13, font_colour=DARK_TEXT, line_spacing=1.5)

    # Source citations at bottom
    if sources:
        unique_sources = list(dict.fromkeys(sources))[:4]
        src_text = "Sources: " + "; ".join(unique_sources)
        _add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
                      src_text, font_size=9, font_colour=MED_GREY, italic=True)

    _add_text_box(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
                  f"Avia Solutions - {config.get('origin', '')}-{config.get('destination', '')}",
                  font_size=9, font_colour=MED_GREY)


def _slide_forecast(prs, config):
    """Forecast summary slide with key numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                  "Route Forecast Summary", font_size=32, font_colour=NAVY,
                  bold=True, font_name="Georgia")
    _add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.04), ACCENT_GOLD)

    forecast = config.get("forecast", {})
    if not forecast or not forecast.get('grand_total'):
        _add_text_box(slide, Inches(0.8), Inches(2), Inches(10), Inches(1),
                      "Forecast data not yet available - run the QSI pipeline first.",
                      font_size=16, font_colour=MED_GREY, italic=True)
        return

    # Large metric cards
    gt = forecast.get('grand_total', 0)
    lf = forecast.get('load_factor', 0)
    lf_str = f"{lf:.0%}" if isinstance(lf, float) and lf < 10 else f"{lf:.1%}" if isinstance(lf, float) else str(lf)
    p2p = forecast.get('p2p_total', 0)
    cnx_dest = forecast.get('cnx_dest_total', 0)
    cnx_home = forecast.get('cnx_home_total', 0)
    cnx_total = cnx_dest + cnx_home

    cards = [
        ("Annual Passengers", f"{gt:,.0f}", "Total forecast demand"),
        ("Load Factor", lf_str, "Annual average"),
        ("P2P Passengers", f"{p2p:,.0f}", f"{p2p/gt:.0%} of total" if gt > 0 else ""),
        ("Connecting Pax", f"{cnx_total:,.0f}", f"{cnx_total/gt:.0%} of total" if gt > 0 else ""),
    ]

    for i, (label, value, sub) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 5.8
        y = 1.8 + row * 2.4

        _add_shape_rect(slide, Inches(x), Inches(y), Inches(5.2), Inches(2),
                        LIGHT_GREY, border_colour=ICE_BLUE, border_width=1)
        _add_text_box(slide, Inches(x + 0.3), Inches(y + 0.2), Inches(4.6), Inches(0.4),
                      label, font_size=12, font_colour=MED_GREY, bold=True)
        _add_text_box(slide, Inches(x + 0.3), Inches(y + 0.7), Inches(4.6), Inches(0.8),
                      str(value), font_size=36, font_colour=NAVY, bold=True)
        _add_text_box(slide, Inches(x + 0.3), Inches(y + 1.5), Inches(4.6), Inches(0.3),
                      sub, font_size=10, font_colour=MED_GREY)

    # Service parameters
    svc = (f"{config.get('airline_name', '')} | "
           f"{config.get('aircraft_type', '')} ({config.get('seats', '')} seats) | "
           f"{config.get('frequency', '')}x weekly")
    _add_text_box(slide, Inches(0.8), Inches(6.8), Inches(10), Inches(0.3),
                  svc, font_size=11, font_colour=DARK_TEXT)

    _add_text_box(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
                  "Avia Solutions - Confidential", font_size=9, font_colour=MED_GREY)


def _slide_forecast_breakdown(prs, config):
    """Detailed forecast breakdown slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                  "Forecast Breakdown", font_size=28, font_colour=NAVY,
                  bold=True, font_name="Georgia")
    _add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.04), ACCENT_GOLD)

    forecast = config.get("forecast", {})
    if not forecast:
        return

    gt = forecast.get('grand_total', 0)
    p2p = forecast.get('p2p_total', 0)
    cnx_home = forecast.get('cnx_home_total', 0)
    cnx_dest = forecast.get('cnx_dest_total', 0)
    lf = forecast.get('load_factor', 0)
    freq = config.get('frequency', 7)
    seats = config.get('seats', 214)
    capacity = freq * 52 * seats * 2

    # Table-style layout
    rows = [
        ("Point-to-Point (P2P)", f"{p2p:,.0f}", f"{p2p/gt:.0%}" if gt > 0 else ""),
        (f"Connecting @ {config.get('origin', 'Home')}", f"{cnx_home:,.0f}", f"{cnx_home/gt:.0%}" if gt > 0 else ""),
        (f"Connecting @ {config.get('destination', 'Dest')}", f"{cnx_dest:,.0f}", f"{cnx_dest/gt:.0%}" if gt > 0 else ""),
        ("", "", ""),
        ("TOTAL PASSENGERS", f"{gt:,.0f}", "100%"),
        ("", "", ""),
        ("Annual Capacity (two-way)", f"{capacity:,.0f}", ""),
        ("Load Factor", f"{lf:.1%}" if isinstance(lf, float) else str(lf), ""),
        ("Frequency", f"{freq}x weekly", f"{'Daily' if freq == 7 else f'{freq}x/week'}"),
        ("Aircraft", f"{config.get('aircraft_type', '')} ({seats} seats)", ""),
    ]

    for i, (label, value, pct) in enumerate(rows):
        y = 1.6 + i * 0.48
        is_total = label.startswith("TOTAL")
        is_blank = label == ""

        if is_blank:
            continue

        if is_total:
            _add_shape_rect(slide, Inches(0.8), Inches(y - 0.05), Inches(10.4), Inches(0.45),
                            LIGHT_GREY)

        _add_text_box(slide, Inches(0.8), Inches(y), Inches(5), Inches(0.4),
                      label, font_size=13, font_colour=DARK_TEXT,
                      bold=is_total)
        _add_text_box(slide, Inches(6.5), Inches(y), Inches(2.5), Inches(0.4),
                      value, font_size=13, font_colour=NAVY,
                      bold=is_total, alignment=PP_ALIGN.RIGHT)
        if pct:
            _add_text_box(slide, Inches(9.5), Inches(y), Inches(1.5), Inches(0.4),
                          pct, font_size=11, font_colour=MED_GREY,
                          alignment=PP_ALIGN.RIGHT)

    _add_text_box(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
                  "Avia Solutions - Confidential", font_size=9, font_colour=MED_GREY)


def _slide_assumptions(prs, config):
    """Key assumptions and parameters slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                  "Key Assumptions", font_size=28, font_colour=NAVY,
                  bold=True, font_name="Georgia")
    _add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.04), ACCENT_GOLD)

    assumptions = config.get("assumptions", {})
    forecast = config.get("forecast", {})

    lines = [
        f"Airline: {config.get('airline_name', 'N/A')}",
        f"Route: {config.get('origin', '')} - {config.get('destination', '')} ({config.get('origin_city', '')} - {config.get('dest_city', '')})",
        f"Aircraft: {config.get('aircraft_type', 'N/A')} with {config.get('seats', 'N/A')} seats",
        f"Frequency: {config.get('frequency', 'N/A')}x weekly",
        f"Demand source: Sabre MIDT (indirect O&D passengers)",
        f"Schedule source: OAG published schedules",
    ]

    if assumptions:
        if 'qsi_adjustment' in assumptions:
            lines.append(f"QSI adjustment factor: {assumptions['qsi_adjustment']:.3f}")
        if 'qsi_ceiling' in assumptions:
            lines.append(f"QSI ceiling: {assumptions['qsi_ceiling']:.2f}")
        if 'stimulation' in assumptions:
            lines.append(f"P2P stimulation factor: {assumptions['stimulation']:.2f}")
        if 'capture_rate' in assumptions:
            lines.append(f"P2P capture rate: {assumptions['capture_rate']:.0%}")

    lines.extend([
        "",
        "Methodology: Avia Solutions QSI (Quality of Service Index)",
        "All connecting traffic modelled through hub connection analysis",
        "P2P demand derived from indirect O&D flows with stimulation for new direct service",
    ])

    _add_multiline(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5),
                   [l for l in lines if l], font_size=13, font_colour=DARK_TEXT,
                   line_spacing=1.5)

    _add_text_box(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
                  "Avia Solutions - Confidential", font_size=9, font_colour=MED_GREY)


def _slide_connecting_cities(prs, config):
    """Top connecting cities slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    dest = config.get('destination', 'Hub')
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                  f"Top Connecting Markets via {dest}", font_size=28, font_colour=NAVY,
                  bold=True, font_name="Georgia")
    _add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.04), ACCENT_GOLD)

    cnx_cities = config.get("connecting_cities", [])
    if not cnx_cities:
        _add_text_box(slide, Inches(0.8), Inches(2), Inches(10), Inches(1),
                      "Connecting city detail available in the QSI output workbook.",
                      font_size=14, font_colour=MED_GREY, italic=True)
        return

    # Header row
    _add_shape_rect(slide, Inches(0.8), Inches(1.5), Inches(10.4), Inches(0.45), NAVY)
    _add_text_box(slide, Inches(1.0), Inches(1.52), Inches(4), Inches(0.4),
                  "City", font_size=12, font_colour=WHITE, bold=True)
    _add_text_box(slide, Inches(5.5), Inches(1.52), Inches(2.5), Inches(0.4),
                  "Passengers", font_size=12, font_colour=WHITE, bold=True,
                  alignment=PP_ALIGN.RIGHT)
    _add_text_box(slide, Inches(8.5), Inches(1.52), Inches(2.5), Inches(0.4),
                  "Share", font_size=12, font_colour=WHITE, bold=True,
                  alignment=PP_ALIGN.RIGHT)

    # Top 15 cities
    total_cnx = sum(c.get('pax', 0) for c in cnx_cities)
    for i, city in enumerate(cnx_cities[:15]):
        y = 2.05 + i * 0.33
        bg = LIGHT_GREY if i % 2 == 0 else WHITE
        _add_shape_rect(slide, Inches(0.8), Inches(y - 0.02), Inches(10.4), Inches(0.33), bg)
        _add_text_box(slide, Inches(1.0), Inches(y), Inches(4), Inches(0.3),
                      city.get('city', ''), font_size=11, font_colour=DARK_TEXT)
        _add_text_box(slide, Inches(5.5), Inches(y), Inches(2.5), Inches(0.3),
                      f"{city.get('pax', 0):,.0f}", font_size=11, font_colour=DARK_TEXT,
                      alignment=PP_ALIGN.RIGHT)
        share = city.get('pax', 0) / total_cnx if total_cnx > 0 else 0
        _add_text_box(slide, Inches(8.5), Inches(y), Inches(2.5), Inches(0.3),
                      f"{share:.1%}", font_size=11, font_colour=MED_GREY,
                      alignment=PP_ALIGN.RIGHT)

    _add_text_box(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
                  "Avia Solutions - Confidential", font_size=9, font_colour=MED_GREY)


def _slide_methodology(prs, config):
    """QSI methodology overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
                  "Methodology", font_size=28, font_colour=NAVY,
                  bold=True, font_name="Georgia")
    _add_shape_rect(slide, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.04), ACCENT_GOLD)

    lines = [
        "Avia Solutions QSI (Quality of Service Index) Forecast Model",
        "",
        "The forecast uses Avia Solutions' proprietary QSI methodology, developed over "
        "30+ years of aviation consulting experience. The model evaluates the quality of "
        "service offered by every competing itinerary between each city pair and allocates "
        "demand proportionally based on relative service quality.",
        "",
        "Key Model Components:",
        "",
        "1. Schedule Analysis: OAG-published flight schedules are processed to identify "
        "all feasible itineraries (direct and connecting) between the route endpoints and "
        "every potential connecting market.",
        "",
        "2. Connection Quality: Each connecting itinerary is scored on elapsed time, "
        "connection time, circuity, number of stops, and carrier/alliance factors.",
        "",
        "3. QSI Scoring: Quality of Service scores are calculated for each itinerary, "
        "with the proposed new service competing against all existing options.",
        "",
        "4. Demand Allocation: Total market demand (from Sabre MIDT booking data) is "
        "allocated across competing itineraries proportional to their QSI scores.",
        "",
        "5. Calibration: Results are calibrated against commercial benchmarks to ensure "
        "forecasts reflect real-world load factors and competitive dynamics.",
    ]

    _add_multiline(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(5),
                   [l for l in lines if l], font_size=12, font_colour=DARK_TEXT,
                   line_spacing=1.35)

    _add_text_box(slide, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
                  "Avia Solutions - Confidential", font_size=9, font_colour=MED_GREY)


def _slide_why_route(prs, config):
    """'Why this route' slide with 4 cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, NAVY)

    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                  "Why This Route?", font_size=32, font_colour=WHITE,
                  bold=True, font_name="Georgia")
    _add_shape_rect(slide, Inches(0.8), Inches(1.25), Inches(2), Inches(0.04), ACCENT_GOLD)

    why_points = config.get("why_points", [
        {"title": "Market Size", "text": "DATA REQUIRED"},
        {"title": "Strategic Fit", "text": "DATA REQUIRED"},
        {"title": "Competitive Gap", "text": "DATA REQUIRED"},
        {"title": "Economic Drivers", "text": "DATA REQUIRED"},
    ])

    for i, point in enumerate(why_points[:4]):
        col = i % 2
        row = i // 2
        x = 0.8 + col * 5.8
        y = 1.8 + row * 2.6

        # Card background
        _add_shape_rect(slide, Inches(x), Inches(y), Inches(5.2), Inches(2.2),
                        DARK_NAVY, border_colour=ACCENT_GOLD, border_width=1)
        # Title
        _add_text_box(slide, Inches(x + 0.3), Inches(y + 0.2), Inches(4.6), Inches(0.5),
                      point.get("title", ""), font_size=18, font_colour=ACCENT_GOLD, bold=True)
        # Text
        _add_text_box(slide, Inches(x + 0.3), Inches(y + 0.8), Inches(4.6), Inches(1.2),
                      point.get("text", "DATA REQUIRED"), font_size=13,
                      font_colour=ICE_BLUE, line_spacing=1.4)


def _slide_closing(prs, config):
    """Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, NAVY)

    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_GOLD)

    _add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.2),
                  "Thank You", font_size=48, font_colour=WHITE, bold=True,
                  font_name="Georgia", alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.8),
                  "Avia Solutions Limited", font_size=20,
                  font_colour=ICE_BLUE, alignment=PP_ALIGN.CENTER)

    contact = config.get("client_name", "")
    if contact:
        _add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5),
                      f"Prepared for {contact}", font_size=14,
                      font_colour=MED_GREY, alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
                  "john.carter@aviasolutions.com | +44 (0)7XXX XXX XXX",
                  font_size=12, font_colour=MED_GREY, alignment=PP_ALIGN.CENTER)

    _add_shape_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT_GOLD)


# ============================================================================
# MAIN GENERATOR
# ============================================================================

def generate_presentation(config: Dict[str, Any],
                          research_blocks: Optional[Dict] = None,
                          output_path: Optional[str] = None) -> str:
    """
    Generate a complete city pair PPTX presentation.

    Args:
        config: Route/airline/forecast configuration dict
        research_blocks: Dict of block_id -> {findings: [...], summary: str}
        output_path: Where to save the .pptx file

    Returns:
        Path to generated .pptx file
    """
    if not HAS_PPTX:
        raise ImportError("python-pptx is required. Install with: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Prepare research data
    if research_blocks is None:
        research_blocks = {}

    research = {
        "executive_summary": config.get("executive_summary", ""),
    }

    # Build section list
    sections = [
        ("01", "Executive Summary"),
        ("02", "Economic Context & Demand Drivers"),
        ("03", "Airport Overview"),
        ("04", "Route Forecast"),
        ("05", "Key Assumptions & Methodology"),
        ("06", "Why This Route"),
    ]

    # ---- Build slides ----

    # 1. Cover
    _slide_cover(prs, config)

    # 2. Contents
    _slide_contents(prs, config, sections)

    # 3. Executive summary
    _slide_executive_summary(prs, config, research)

    # 4. Research block slides
    block_display_order = [
        ("economic_context", "Economic Context"),
        ("corporate_links", "Corporate & Technology Links"),
        ("tourism", "Tourism & Visitor Economy"),
        ("trade", "Trade & Investment"),
        ("airport_overview", "Airport Overview"),
        ("diaspora", "Diaspora & VFR Population"),
        ("passenger_profile", "Passenger Profile"),
        ("non_cannibalization", "Market Stimulation Evidence"),
        ("case_study", "Comparable Route Case Study"),
        ("education", "Education & Student Links"),
    ]

    for block_id, display_name in block_display_order:
        if block_id in research_blocks:
            block_data = research_blocks[block_id]
            findings = block_data.get("findings", [])
            summary = block_data.get("summary", "")
            if findings or summary:
                _slide_research_block(prs, config, block_id, display_name,
                                      findings, summary)

    # 5. Forecast slides
    _slide_forecast(prs, config)
    if config.get("forecast", {}).get("grand_total"):
        _slide_forecast_breakdown(prs, config)

    # 6. Connecting cities (if available)
    if config.get("connecting_cities"):
        _slide_connecting_cities(prs, config)

    # 7. Assumptions
    _slide_assumptions(prs, config)

    # 8. Methodology
    _slide_methodology(prs, config)

    # 9. Why this route
    _slide_why_route(prs, config)

    # 10. Closing
    _slide_closing(prs, config)

    # Save
    if output_path is None:
        output_path = os.path.join(
            tempfile.mkdtemp(),
            f"{config.get('origin', 'XXX')}_{config.get('destination', 'XXX')}_"
            f"{config.get('airline_name', 'Airline').replace(' ', '_')}.pptx"
        )

    # House rule: every generated file carries Avia Solutions as author and as
    # last modified by, never the generating library. The decks the site has been
    # producing carried an empty author, which is what a reader's file properties
    # showed when they opened one.
    # These decks are published by The Aviation Observatory, the separate company,
    # not by Avia Solutions. Deliberate departure from the Avia house rule, product
    # output only.
    publisher = os.environ.get("AVIA_DECK_AUTHOR", "The Aviation Observatory")
    cp = prs.core_properties
    cp.author = publisher
    cp.last_modified_by = publisher
    cp.company = publisher
    cp.title = "%s to %s route pitch" % (config.get("origin_city", ""),
                                         config.get("dest_city", ""))
    prs.save(output_path)
    return output_path


# ============================================================================
# HELPER: Extract research from executor for presentation
# ============================================================================

def extract_research_for_pptx(executor) -> Dict[str, Dict]:
    """
    Convert a ResearchExecutor's blocks into the format
    expected by generate_presentation().
    """
    result = {}
    for block_id, bf in executor.blocks.items():
        findings_dicts = []
        for f in bf.findings:
            fd = {
                "claim": f.claim,
                "value": f.value or "",
                "year": f.year or "",
                "source_name": f.citations[0].source_name if f.citations else "",
            }
            findings_dicts.append(fd)

        result[block_id] = {
            "findings": findings_dicts,
            "summary": bf.summary or bf.presentation_text or "",
        }
    return result
