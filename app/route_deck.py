#!/usr/bin/env python3
"""
Avia Solutions - route deck: forecast + P&L on brand, driven by LIVE tool output.
=================================================================================
build_deck(out, forecast, pnl, meta) renders a 3-slide Avia deck (title, demand forecast,
route P&L) from plain dicts, so genoa_nyc.py (or any route case) emits its own client deck:
the P&L slide shows the REAL cost stack from aircraft_economics RoutePnL.compute(), not a
hand-typed summary. Pure python-pptx; no economics import here (takes the compute() dict).

  forecast = dict(pop, nyc_od, leaked, repatriated, directional, split=[(label,share)...], fit_lines=[...])
  pnl      = an aircraft_economics RoutePnL.compute() dict
  meta     = dict(title, subtitle, origin, dest, aircraft, sector_nm, fare_ow, plan_lf,
                  frequency, annual_profit, disclaimer, [sensitivity=[(label,margin,annual)...]])
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1F, 0x38, 0x64); BLUE = RGBColor(0x2E, 0x5E, 0x8C)
GREY = RGBColor(0x59, 0x59, 0x59); LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
GREEN = RGBColor(0x00, 0x61, 0x00); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = "Arial"


def _rect(s, l, t, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def _txt(s, l, t, w, h, text, size, *, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
         italic=False, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(l, t, w, h); tf = b.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = BODY; f.color.rgb = color
    return b


def _stat(s, l, t, w, value, label, vcolor=NAVY, vsize=33):
    b = s.shapes.add_textbox(l, t, w, Inches(1.2)); tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = value
    r.font.size = Pt(vsize); r.font.bold = True; r.font.name = BODY; r.font.color.rgb = vcolor
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(11.5); r2.font.name = BODY; r2.font.color.rgb = GREY


def _money(v):
    return f"${v:,.0f}"


def build_deck(out, forecast, pnl, meta):
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height; blank = prs.slide_layouts[6]
    # SEASON labelling: volume figures are for the operating season; the addressable market stays annual.
    _season = meta.get("season") or {}
    _smode = _season.get("mode", "annual")
    _seasonal = _smode in ("summer", "winter")
    _pnoun = _smode if _seasonal else "annual"        # "summer" / "winter" / "annual"
    _padj = _smode.capitalize() if _seasonal else "Annual"
    _weeks = float(_season.get("weeks") or 52)
    _days = (_weeks * 7.0) if _seasonal else 365.0     # season operating days for pdew -> volume

    # 1) title
    s = prs.slides.add_slide(blank); _rect(s, 0, 0, SW, SH, NAVY)
    _txt(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.6), meta['title'], 50, bold=True, color=WHITE)
    _txt(s, Inches(0.95), Inches(3.8), Inches(11.5), Inches(0.8), meta.get('subtitle', ''), 22, color=LIGHT)
    _txt(s, Inches(0.95), Inches(6.5), Inches(6), Inches(0.5), "Avia Solutions", 14, bold=True, color=WHITE)
    _txt(s, Inches(8.0), Inches(6.5), Inches(4.4), Inches(0.5), "Indicative, directional guidance", 11,
         color=LIGHT, align=PP_ALIGN.RIGHT, italic=True)

    # 2) forecast
    s = prs.slides.add_slide(blank)
    oname = meta.get('origin_name', meta['origin'])
    _txt(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.7),
         f"Catchment and demand: {meta['dest']} from {oname}", 32, bold=True)
    _txt(s, Inches(0.62), Inches(1.15), Inches(12.1), Inches(0.5), forecast.get('subtitle', ''), 15, color=GREY)
    _stat(s, Inches(0.6), Inches(1.95), Inches(3.0), forecast['market'],
          f"addressable market, each way ({forecast.get('market_2w','')} both ways/yr)")
    _stat(s, Inches(3.7), Inches(1.95), Inches(3.0), forecast['captured'], "captured point-to-point")
    _stat(s, Inches(6.8), Inches(1.95), Inches(3.0), forecast['feed'], "connecting feed")
    _stat(s, Inches(9.9), Inches(1.95), Inches(3.0), forecast['total'],
          f"total forecast, each way ({forecast.get('total_2w','')} both ways/yr)", vcolor=GREEN)
    _txt(s, Inches(0.6), Inches(3.35), Inches(6.2), Inches(0.4),
         f"Where the region's {meta['dest']} demand departs today", 15, bold=True)
    y = Inches(3.85); barL = Inches(2.4); maxW = 4.0
    for name, sh in forecast['split']:
        _txt(s, Inches(0.6), y, Inches(1.8), Inches(0.32), name, 12, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        w = Inches(max(maxW * sh, 0.04)); col = GREEN if name == forecast.get('home_label') else BLUE
        _rect(s, barL, y + Inches(0.03), w, Inches(0.26), col)
        _txt(s, barL + w + Inches(0.05), y, Inches(1.0), Inches(0.32), f"{sh:.0%}", 12, bold=True, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
        y = y + Inches(0.46)
    rx = Inches(7.4); _rect(s, rx, Inches(3.7), Inches(5.3), Inches(2.95), LIGHT)
    _txt(s, rx + Inches(0.25), Inches(3.85), Inches(4.9), Inches(0.4), f"A {oname} nonstop", 16, bold=True)
    _txt(s, rx + Inches(0.25), Inches(4.35), Inches(4.9), Inches(2.2), "\n".join(forecast['fit_lines']), 13.5)
    _txt(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4),
         "Source: GeoNames population, OSRM road times, Sabre point-of-origin O&D; catchment calibrated to the observed airport split.",
         9.5, color=GREY, italic=True)

    # 2b) five-year demand build (grown at the measured market trend)
    pj = forecast.get('projection')
    if pj and pj.get('build'):
        RED = RGBColor(0xB0, 0x00, 0x00)
        s = prs.slides.add_slide(blank)
        _txt(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.7),
             f"Five-year demand build: {meta['dest']} from {oname}", 32, bold=True)
        _cg = pj.get('cagr', 0.0)
        _txt(s, Inches(0.62), Inches(1.15), Inches(12.1), Inches(0.5),
             f"Demand grown at the measured market trend of {_cg*100:+.1f}% a year; capacity held on the chosen schedule",
             15, color=GREY)
        y0 = Inches(2.15); TW = Inches(9.2)
        _rect(s, Inches(0.6), y0, TW, Inches(0.42), NAVY)
        _txt(s, Inches(0.75), y0, Inches(2.0), Inches(0.42), "Year", 12.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        _txt(s, Inches(2.7), y0, Inches(2.6), Inches(0.42), "Demand (each way)", 12.5, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _txt(s, Inches(5.5), y0, Inches(2.3), Inches(0.42), "Carried", 12.5, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _txt(s, Inches(8.0), y0, Inches(1.6), Inches(0.42), "Load", 12.5, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        yy = Inches(2.62)
        for i, x in enumerate(pj['build']):
            if i % 2 == 0:
                _rect(s, Inches(0.6), yy, TW, Inches(0.38), LIGHT)
            _lbl = f"{x['year']} (base)" if x.get('offset') == 0 else str(x['year'])
            _spill = (x.get('spill', 0) or 0) > 0.02 * max(x.get('demand', 1), 1)
            _ld = f"{round((x.get('load') or 0)*100)}%" if x.get('load') is not None else "-"
            _txt(s, Inches(0.75), yy, Inches(2.2), Inches(0.38), _lbl, 12.5, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
            _txt(s, Inches(2.7), yy, Inches(2.6), Inches(0.38), f"{x['demand']:,}", 12.5, color=NAVY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            _txt(s, Inches(5.5), yy, Inches(2.3), Inches(0.38), f"{x['carried']:,}", 12.5, color=NAVY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            _txt(s, Inches(8.0), yy, Inches(1.6), Inches(0.38), _ld, 12.5, bold=_spill, color=(RED if _spill else NAVY), align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            yy = yy + Inches(0.42)
        _b = pj['build']; _full = next((r for r in _b if (r.get('spill', 0) or 0) > 0.02 * max(r.get('demand', 1), 1)), None)
        _last = _b[-1]; _lastld = round((_last.get('load') or 0) * 100) if _last.get('load') is not None else None
        _note = (f"Demand fills the aircraft by {_full['year']}; beyond that it spills, the case for added frequency or a larger gauge."
                 if _full else f"Within the chosen schedule throughout, reaching {_lastld}% load by {_last['year']}." if _lastld is not None
                 else f"Demand grows steadily to {_last['year']}.")
        _txt(s, Inches(0.6), yy + Inches(0.25), Inches(9.2), Inches(0.7), _note, 13.5, color=NAVY)
        _txt(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4),
             "Growth is the measured Sabre O&D trend for this market, clamped to a sustainable range; a launch-year figure, not the current year.",
             9.5, color=GREY, italic=True)

    # 3) route P&L with the real cost stack
    s = prs.slides.add_slide(blank)
    _txt(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.7),
         f"Route economics: {meta['aircraft']}, {meta['origin']} - {meta['dest']}", 32, bold=True)
    _txt(s, Inches(0.62), Inches(1.15), Inches(12.1), Inches(0.5), meta.get('pnl_subtitle', ''), 15, color=GREY)
    # THE HEADLINE IS CONTRIBUTION, NOT PROFIT. Avia cannot source a lease rate and does not publish
    # one, and the airport charges here are a generic placeholder rather than this airport pair. So
    # the slide leads with what the route contributes towards ownership after all cash costs, and
    # states the ownership cost at which it stops working. Profit stays on the cost stack, below the
    # contribution line and with both plugs labelled.
    _own = pnl['ownership'] + pnl['insurance']
    _contrib = pnl['profit'] + _own
    pos = GREEN if _contrib >= 0 else RGBColor(0xB0, 0x00, 0x00)
    _stat(s, Inches(0.6), Inches(1.95), Inches(3.0), _money(_contrib),
          "contribution towards ownership, per rotation", vcolor=pos)
    _stat(s, Inches(3.7), Inches(1.95), Inches(3.0), _money(_own),
          "ownership plug in this case", vcolor=GREY)
    _stat(s, Inches(6.8), Inches(1.95), Inches(3.0), f"{pnl['breakeven_lf']:.0%}", "breakeven load factor")
    _stat(s, Inches(9.9), Inches(1.95), Inches(3.0),
          _money(meta['annual_profit'] + _own * float(meta.get('frequency') or 0) * 52.0),
          f"{_pnoun} contribution ({meta['frequency']}x/week)", vcolor=pos)
    # cost stack table (left)
    _txt(s, Inches(0.6), Inches(3.3), Inches(5.7), Inches(0.4), "Per rotation (return)", 15, bold=True)
    rows = [("Revenue", pnl['gross_rev'], False),
            ("Fuel", -pnl['fuel'], False),
            ("Maintenance", -pnl['maintenance'], False),
            ("Crew", -pnl['crew'], False),
            ("Airport charges (plug)", -(pnl['landing'] + pnl['per_pax'] + pnl['handling']), False),
            ("En-route (nav, plug)", -pnl['nav'], False),
            ("Catering", -pnl['catering'], False),
            ("Overhead & sales", -(pnl['admin'] + pnl['sales']), False),
            ("Contribution to ownership", _contrib, True),
            ("Ownership & insurance (plug)", -_own, False),
            ("Operating profit, plugs as set", pnl['profit'], True)]
    y = Inches(3.78)
    for k, v, hl in rows:
        _rect(s, Inches(0.6), y, Inches(5.6), Inches(0.4), NAVY if hl else LIGHT)
        c = WHITE if hl else NAVY
        _txt(s, Inches(0.75), y, Inches(3.7), Inches(0.4), k, 12.5, bold=hl, color=c, anchor=MSO_ANCHOR.MIDDLE)
        _txt(s, Inches(4.35), y, Inches(1.7), Inches(0.4), _money(v), 12.5, bold=hl, color=c,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        y = y + Inches(0.305)
    # assumptions + sensitivity (right)
    rx = Inches(6.7)
    _txt(s, rx, Inches(3.3), Inches(6.0), Inches(0.4), "Assumptions and basis", 15, bold=True)
    asum = (f"Aircraft: {meta['aircraft']}, {meta['origin']}-{meta['dest']} {meta['sector_nm']:,} nm, {meta['frequency']}x/week\n"
            f"Fare: ${meta['fare_ow']:,} one-way economy (Sabre)\n"
            f"Load factor: {meta['plan_lf']:.0%} planned; CASK ${pnl['cask']*100:.1f}c/seat-km\n"
            f"Maintenance: {meta.get('maint_basis','')}\n"
            f"Ownership: {meta.get('own_basis','')}")
    _txt(s, rx, Inches(3.78), Inches(6.0), Inches(2.0), asum, 12, color=NAVY)
    sens = meta.get('sensitivity')
    if sens:
        bx = rx; bw = Inches(2.9)
        for i, (lab, mg, ann) in enumerate(sens[:2]):
            x = bx + (Inches(3.1) if i else Inches(0))
            _rect(s, x, Inches(5.55), bw, Inches(1.15), LIGHT)
            _txt(s, x + Inches(0.2), Inches(5.63), bw - Inches(0.4), Inches(0.4), lab, 12.5, bold=True)
            _txt(s, x + Inches(0.2), Inches(6.0), bw - Inches(0.4), Inches(0.6),
                 f"{mg:.1%} margin\n{_money(ann)} / year", 13.5, bold=True, color=(GREEN if i == 0 else NAVY))
    _txt(s, Inches(0.6), Inches(6.98), Inches(12.1), Inches(0.4), meta.get('disclaimer', ''), 9.5, color=GREY, italic=True)

    if meta.get('full_report'):
        # 4) connecting feed detail (PDEW each way)
        s = prs.slides.add_slide(blank)
        _txt(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.7), "Connecting feed detail", 32, bold=True)
        _txt(s, Inches(0.62), Inches(1.15), Inches(12.1), Inches(0.5),
             "Top connecting markets each way, passengers per day each way (PDEW)", 15, color=GREY)

        def _feedtbl(x, title, rows):
            _txt(s, x, Inches(1.95), Inches(5.8), Inches(0.4), title, 14, bold=True)
            _rect(s, x, Inches(2.45), Inches(5.8), Inches(0.36), NAVY)
            _txt(s, x + Inches(0.1), Inches(2.45), Inches(2.6), Inches(0.36), "Market", 11.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
            _txt(s, x + Inches(3.5), Inches(2.45), Inches(1.05), Inches(0.36), "PDEW", 11.5, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            _txt(s, x + Inches(4.6), Inches(2.45), Inches(1.1), Inches(0.36), _padj, 11.5, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            yy = Inches(2.85)
            for i, row in enumerate(rows[:10]):
                if i % 2 == 0:
                    _rect(s, x, yy, Inches(5.8), Inches(0.32), LIGHT)
                _txt(s, x + Inches(0.1), yy, Inches(3.4), Inches(0.32), f"{row.get('name')} ({row.get('code')})", 11, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
                _txt(s, x + Inches(3.4), yy, Inches(1.15), Inches(0.32), f"{row.get('pdew', 0):,.1f}", 11, color=NAVY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
                _txt(s, x + Inches(4.55), yy, Inches(1.15), Inches(0.32), f"{row.get('pdew', 0) * _days:,.0f}", 11, color=NAVY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
                yy = yy + Inches(0.34)
        _feedtbl(Inches(0.6), f"Behind {meta['origin']}", forecast.get('behind_pdew') or [])
        _feedtbl(Inches(6.9), f"Beyond {meta['dest']}", forecast.get('beyond_pdew') or [])
        _txt(s, Inches(0.6), Inches(6.98), Inches(12.1), Inches(0.4),
             "PDEW = passengers per day each way. Onward O&D on the selected airline, alliance-weighted and circuity-screened.",
             9.5, color=GREY, italic=True)

        # 5) catchment and capture assumptions
        s = prs.slides.add_slide(blank)
        _txt(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.7), "Catchment and capture assumptions", 32, bold=True)
        _txt(s, Inches(0.62), Inches(1.15), Inches(12.1), Inches(0.5),
             f"How {meta.get('origin_name', meta['origin'])}'s demand is measured and captured", 15, color=GREY)
        _txt(s, Inches(0.6), Inches(1.95), Inches(6.0), Inches(0.4), "Catchment split across airports today", 14, bold=True)
        crows = forecast.get('catchment_rows') or forecast.get('split') or []
        mx = max([sh for _, sh in crows] or [1]) or 1
        yy = Inches(2.5)
        for name, shr in crows[:9]:
            _txt(s, Inches(0.6), yy, Inches(2.0), Inches(0.3), name, 11, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
            w = Inches(max(3.6 * shr / mx, 0.04)); col = GREEN if name == forecast.get('home_label') else BLUE
            _rect(s, Inches(2.7), yy + Inches(0.03), w, Inches(0.24), col)
            _txt(s, Inches(2.7) + w + Inches(0.05), yy, Inches(0.9), Inches(0.3), f"{shr:.0%}", 11, bold=True, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
            yy = yy + Inches(0.42)
        rx = Inches(7.4); _rect(s, rx, Inches(2.4), Inches(5.3), Inches(3.95), LIGHT)
        _txt(s, rx + Inches(0.25), Inches(2.55), Inches(4.9), Inches(3.7), meta.get('catchment_text', ''), 13)
        _txt(s, Inches(0.6), Inches(6.98), Inches(12.1), Inches(0.4),
             "Where measured survey or mobility data exists, it overrides the modelled capture.", 9.5, color=GREY, italic=True)

        # 6) methodology
        s = prs.slides.add_slide(blank)
        _txt(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.7), "Methodology", 32, bold=True)
        _txt(s, Inches(0.62), Inches(1.15), Inches(12.1), Inches(0.5), "How Cortex builds the forecast, step by step", 15, color=GREY)
        steps = ["Catchment", "Market (Sabre O&D)", "Capture (QSI + access)", "Stimulation", "Connecting feed", "Capacity cap", "Forecast + economics"]
        bw = Inches(1.68); gap = Inches(0.07)
        for i, st in enumerate(steps):
            x = Inches(0.6) + i * (bw + gap)
            _rect(s, x, Inches(2.1), bw, Inches(1.0), NAVY if i == len(steps) - 1 else BLUE)
            _txt(s, x + Inches(0.06), Inches(2.1), bw - Inches(0.12), Inches(1.0), st, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        expl = ("Catchment: resident population within drive time of the origin, from GeoNames and least-cost road times.\n"
                "Market: measured Sabre point-of-origin O&D each way in that catchment.\n"
                "Capture: the QSI and access share the new nonstop wins from competing airports and airlines; measured survey or mobility data overrides the model where held.\n"
                "Stimulation: the demand a new nonstop creates over the indirect base.\n"
                "Connecting feed: onward O&D behind the origin and beyond the destination on the chosen airline, alliance-weighted and circuity-screened.\n"
                "Capacity cap: demand bounded by the aircraft, frequency and planned load factor.\n"
                f"Forecast and economics: the bounded total each way, then a turnaround and {_pnoun} P&L on validated type costs.")
        _txt(s, Inches(0.6), Inches(3.4), Inches(12.1), Inches(3.3), expl, 13, color=NAVY)
        _txt(s, Inches(0.6), Inches(6.98), Inches(12.1), Inches(0.4), meta.get('disclaimer', ''), 9.5, color=GREY, italic=True)

    cp = prs.core_properties
    cp.author = "Avia Solutions"; cp.last_modified_by = "Avia Solutions"; cp.title = meta['title']
    prs.save(out)
    return out
